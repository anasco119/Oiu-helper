import os
import sqlite3
import time # <--- Ø£Ø¶Ù Ù‡Ø°Ø§ Ø§Ù„Ø³Ø·Ø±
from datetime import date
from datetime import datetime, timedelta
import telebot
from telebot.types import ChatPermissions
import threading
import logging
from dotenv import load_dotenv
from telebot.types import InlineKeyboardMarkup, InlineKeyboardButton
import docx
import fitz                     # PyMuPDF
import google.generativeai as genai
import requests
import cohere
from groq import Groq
import json
import re
from pptx import Presentation
import traceback
import threading
import queue
import time
import random
from flask import Flask, render_template


# Ù…ØªØºÙŠØ±Ø§Øª Ø§Ù„Ø¨ÙŠØ¦Ø©
load_dotenv()
BOT_TOKEN = os.getenv("BOT_TOKEN")
GROUP_ID = int(os.getenv("GROUP_ID"))
ADMIN_ID = int(os.getenv("ADMIN_ID"))
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
WEBHOOK_URL = os.getenv('WEBHOOK_URL')
allowed_channels = set()
env_channels = os.getenv("ALLOWED_CHANNELS", "")
if env_channels.strip():
    allowed_channels = set(map(int, env_channels.split(",")))
BOT_TOKEN_2 = os.getenv("BOT_TOKEN_2")
BOT_TOKEN_3 = os.getenv("BOT_TOKEN_3")

OCR_API_KEY = os.getenv("OCR_SPACE_API_KEY", "helloworld")  # Ù…ÙØªØ§Ø­ Ø§ÙØªØ±Ø§Ø¶ÙŠ Ù„Ù„Ø§Ø®ØªØ¨Ø§Ø±
bot = telebot.TeleBot(BOT_TOKEN)
bot2 = telebot.TeleBot(BOT_TOKEN_2)
bot3 = telebot.TeleBot(BOT_TOKEN_3)




# -------------------------------------------------------------------
# --------- Notofication code section ---
# -------------------------------------------------------------------

def notify_admin(action: str, username: str, user_id: int):
    """
    ÙŠØ±Ø³Ù„ Ø¥Ø´Ø¹Ø§Ø± Ù„Ù„Ø¥Ø¯Ù…Ù† Ø¨Ø®ØµÙˆØµ Ø¹Ù…Ù„ÙŠØ© Ù…Ø¹ÙŠÙ†Ø©.
    action: Ù†ÙˆØ¹ Ø§Ù„Ø¹Ù…Ù„ÙŠØ© (Ù…Ø«Ù„Ø§Ù‹: "ØªÙˆÙ„ÙŠØ¯ Ø£Ù†ÙƒÙŠ Ø¢Ù„ÙŠ")
    """
    bot3.send_message(
        ADMIN_ID,
        f"Ù†Ø¬Ø§Ø­âœ”ï¸ØŒ Ø¹Ù…Ù„ÙŠØ© {action} Ù„: {username} | UID: {user_id}"
    )


def notify_process_info(uid, file_id, username):
    try:

        # Ø¥Ø±Ø³Ø§Ù„ Ø¥Ø´Ø¹Ø§Ø± Ø¥Ù„Ù‰ Ø§Ù„Ø£Ø¯Ù…Ù† ÙÙ‚Ø·
        # Ø¨Ø¯Ù„ user_id Ø§Ø³ØªØ®Ø¯Ù… uid
        bot3.send_message(
            ADMIN_ID,
            f"ğŸ“ ØªÙ… Ø§Ù„Ø§Ù†ØªÙ‡Ø§Ø¡ Ù…Ù† Ù…Ø¹Ø§Ù„Ø¬Ø© Ù…Ù„Ù Ù…Ù† Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…:\n"
            f"ğŸ‘¤ ID: {uid}\n"
            f"ğŸ”— File ID: {file_id or 'no-file'}\n"
            f"ğŸ’¬ Username: @{username if username else 'N/A'}"
            )
    except Exception:
        logging.exception("notify admin failed")


from datetime import datetime, timedelta

usage_count = {}
last_feedback_time = {}

def maybe_send_feedback_request(uid: int, chat_id: int):
    now = datetime.utcnow()
    usage_count[uid] = usage_count.get(uid, 0) + 1

    # Ø§Ù„Ø´Ø±ÙˆØ·
    send_feedback = False
    if usage_count[uid] == 1:  # Ø£ÙˆÙ„ Ø§Ø³ØªØ®Ø¯Ø§Ù…
        send_feedback = True
    elif usage_count[uid] % 5 == 0:  # ÙƒÙ„ 5 Ù…Ø±Ø§Øª
        # Ø´Ø±Ø· Ù…Ø±ÙˆØ± ÙŠÙˆÙ… ÙƒØ§Ù…Ù„
        last_time = last_feedback_time.get(uid)
        if not last_time or (now - last_time) > timedelta(days=1):
            send_feedback = True

    if send_feedback:
        rating_markup = types.InlineKeyboardMarkup()
        rating_markup.row(
            types.InlineKeyboardButton("â­ 1", callback_data="rate_1"),
            types.InlineKeyboardButton("â­ 2", callback_data="rate_2"),
            types.InlineKeyboardButton("â­ 3", callback_data="rate_3")
        )
        rating_markup.row(
            types.InlineKeyboardButton("â­ 4", callback_data="rate_4"),
            types.InlineKeyboardButton("â­ 5", callback_data="rate_5"),
            types.InlineKeyboardButton("ØªØ¬Ø§Ù‡Ù„", callback_data="rate_ignore")
        )

        bot.send_message(
            chat_id,
            "âœ¨ ÙƒÙŠÙ ÙƒØ§Ù†Øª ØªØ¬Ø±Ø¨ØªÙƒ Ù…Ø¹ TestGenie ØŸ\n\nØ§Ø®ØªØ± Ø¹Ø¯Ø¯ Ø§Ù„Ù†Ø¬ÙˆÙ… Ù„Ù„ØªÙ‚ÙŠÙŠÙ…:",
            reply_markup=rating_markup
        )

        last_feedback_time[uid] = now


def send_daily_report():
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            today = datetime.date.today().isoformat()

            cursor.execute("""
                SELECT tests_generated, files_processed, new_users, channel_users, external_users
                FROM daily_stats
                WHERE date=?
            """, (today,))
            row = cursor.fetchone()
            if row:
                tests, files, new_users, channel_users, external_users = row
                msg = f"ğŸ“Š ØªÙ‚Ø±ÙŠØ± Ø§Ù„ÙŠÙˆÙ… ({today}):\n" \
                      f"ğŸ“ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ø§Ù„Ù…ÙˆÙ„Ø¯Ø©: {tests}\n" \
                      f"ğŸ“‚ Ø§Ù„Ù…Ù„ÙØ§Øª Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø©: {files}\n" \
                      f"ğŸ‘¥ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙˆÙ† Ø§Ù„Ø¬Ø¯Ø¯: {new_users} (Ù‚Ù†ÙˆØ§Øª: {channel_users} | Ø®Ø§Ø±Ø¬ÙŠ: {external_users})"
                bot3.send_message(ADMIN_ID, msg)
    except Exception as e:
        logging.error(f"âŒ Error sending daily report: {e}")


def send_top_users_report(top_n: int = 5):
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT user_id, tests_generated, files_processed FROM top_users ORDER BY tests_generated DESC LIMIT ?", (top_n,))
            rows = cursor.fetchall()
            msg = "ğŸ† Ø£ÙØ¶Ù„ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ†:\n"
            for idx, (uid, tests, files) in enumerate(rows, 1):
                msg += f"{idx}. UserID: {uid} | Tests: {tests} | Files: {files}\n"
            bot3.send_message(ADMIN_ID, msg)
    except Exception as e:
        logging.error(f"âŒ Error sending top users report: {e}")


# -------------------------------------------------------------------
# -------- Ù…Ø¹Ø§Ù„Ø¬Ø§Øª Ø§Ù„Ø·Ù„Ø¨Ø§Øª -------------------------------------------------------------------

# Ø¥Ø¹Ø¯Ø§Ø¯ Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± ÙˆØ§Ù„ØªØ­ÙƒÙ…
request_queue = queue.Queue(maxsize=200)
semaphore = threading.Semaphore(5)
num_workers = 5







#!/usr/bin/env python3
# -*- coding: utf-8 -*-


# --- Ø¥Ø¹Ø¯Ø§Ø¯ Ø§Ù„Ù…ÙØ§ØªÙŠØ­ ÙˆØ§Ù„Ø¹Ù…Ù„

session = {}  # <--- Ø£Ø¶Ù Ù‡Ø°Ø§ Ø§Ù„Ø³Ø·Ø±
user_files = {}
# 1. Ø¥Ø¹Ø¯Ø§Ø¯ Google Gemini
gemini_model = None
if GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel('gemini-1.5-flash')
        logging.info("âœ… 1. Gemini configured successfully")
    except Exception as e:
        logging.warning(f"âš ï¸ Could not configure Gemini: {e}")

# 2. Ø¥Ø¹Ø¯Ø§Ø¯ Groq
groq_client = None
if GROQ_API_KEY:
    try:
        groq_client = Groq(api_key=GROQ_API_KEY)
        logging.info("âœ… 2. Groq configured successfully")
    except Exception as e:
        logging.warning(f"âš ï¸ Could not configure Groq: {e}")

# 3. Ø¥Ø¹Ø¯Ø§Ø¯ OpenRouter (Ø³ÙŠØªÙ… Ø§Ø³ØªØ®Ø¯Ø§Ù…Ù‡ Ù„Ù†Ù…ÙˆØ°Ø¬ÙŠÙ† Ù…Ø®ØªÙ„ÙÙŠÙ†)
if OPENROUTER_API_KEY:
    logging.info("âœ… 3. OpenRouter is ready")

# 4. Ø¥Ø¹Ø¯Ø§Ø¯ Cohere
cohere_client = None
if COHERE_API_KEY:
    try:
        cohere_client = cohere.Client(COHERE_API_KEY)
        logging.info("âœ… 4. Cohere configured successfully")
    except Exception as e:
        logging.warning(f"âš ï¸ Could not configure Cohere: {e}")


# --- Ø§Ù„Ø¯Ø§Ù„Ø© Ø§Ù„Ù…ÙˆØ­Ø¯Ø© Ù„ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø±Ø¯ÙˆØ¯ ---

def generate_gemini_response(prompt: str) -> str:
    """
    Tries to generate a response by attempting a chain of services silently.
    It logs errors for the developer but does not send progress messages to the user.
    """
    timeout_seconds = 45

    # 1ï¸âƒ£ OpenRouter - Nous Hermes 2 (Ø£ÙØ¶Ù„ Ø¯Ø¹Ù… Ù„Ù„Ø¹Ø±Ø¨ÙŠØ©)
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 1. OpenRouter (Nous Hermes 2)...")
            headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # â† ØºÙŠÙ‘Ø± Ù‡Ø°Ø§ Ø¥Ù„Ù‰ Ø±Ø§Ø¨Ø· Ø§Ù„Ø¨ÙˆØª
            "X-Title": "AI Quiz Bot"
            }
            model_identifier = "nousresearch/nous-hermes-2-mistral:free"
            response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
                json={
                "model": model_identifier,
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("âœ… Success with OpenRouter (Nous Hermes 2).")
            return result_text
        except Exception as e:
            logging.warning(f"âŒ OpenRouter (Nous Hermes 2) failed: {e}")

    # 2ï¸âƒ£ Groq (LLaMA 3)
    if groq_client:
        try:
            logging.info("Attempting request with: 2. Groq (LLaMA 3)...")
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-8b-8192",
                temperature=0.7,
                timeout=timeout_seconds
            )
            if chat_completion.choices[0].message.content:
                logging.info("âœ… Success with Groq.")
                return chat_completion.choices[0].message.content
            else:
                logging.warning("âŒ Groq returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"âŒ Groq failed: {e}")

    # 3ï¸âƒ£ OpenRouter - Gemma
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 3. OpenRouter (Gemma)...")
            headers = {
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # Replace with your bot's link
                "X-Title": "AI Quiz Bot"
            }
            model_identifier = "google/gemma-7b-it:free"
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                json={"model": model_identifier, "messages": [{"role": "user", "content": prompt}]},
                timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("âœ… Success with OpenRouter (Gemma).")
            return result_text
        except Exception as e:
            logging.warning(f"âŒ OpenRouter (Gemma) failed: {e}")

    # 4ï¸âƒ£ Google Gemini
    if gemini_model:
        try:
            logging.info("Attempting request with: 4. Google Gemini...")
            request_options = {"timeout": timeout_seconds}
            response = gemini_model.generate_content(prompt, request_options=request_options)
            if response.text:
                logging.info("âœ… Success with Gemini.")
                return response.text
            else:
                logging.warning("âŒ Gemini returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"âŒ Gemini failed: {e}")

    # 5ï¸âƒ£ Cohere
    if cohere_client:
        try:
            logging.info("Attempting request with: 5. Cohere...")
            response = cohere_client.chat(model='command-r', message=prompt)
            logging.info("âœ… Success with Cohere.")
            return response.text
        except Exception as e:
            logging.warning(f"âŒ Cohere failed: {e}")

    # ğŸš« All models failed
    logging.error("âŒ All API providers failed. Returning empty string.")
    return ""


def generate_smart_response(prompt: str) -> str:
    """
    Tries to generate a response by attempting a chain of services silently.
    It logs errors for the developer but does not send progress messages to the user.
    """
    timeout_seconds = 45


    #  1ï¸âƒ£ Cohere
    if cohere_client:
        try:
            logging.info("Attempting request with: 5. Cohere...")
            response = cohere_client.chat(model='command-r', message=prompt, temperature=0.8)
            logging.info("âœ… Success with Cohere.")
            return response.text
        except Exception as e:
            logging.warning(f"âŒ Cohere failed: {e}")



    # 2ï¸âƒ£ Google Gemini
    if gemini_model:
        try:
            logging.info("Attempting request with: 4. Google Gemini...")
            request_options = {"timeout": timeout_seconds}
            response = gemini_model.generate_content(prompt, request_options=request_options, temperature=0.8)
            if response.text:
                logging.info("âœ… Success with Gemini.")
                return response.text
            else:
                logging.warning("âŒ Gemini returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"âŒ Gemini failed: {e}")


    #  3ï¸âƒ£  Groq (LLaMA 3)
    if groq_client:
        try:
            logging.info("Attempting request with: 2. Groq (LLaMA 3)...")
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-8b-8192",
                temperature=0.8,
                timeout=timeout_seconds
            )
            if chat_completion.choices[0].message.content:
                logging.info("âœ… Success with Groq.")
                return chat_completion.choices[0].message.content
            else:
                logging.warning("âŒ Groq returned no text. Trying fallback...")
        except Exception as e:
            logging.warning(f"âŒ Groq failed: {e}")

    # 4ï¸âƒ£# 5ï¸âƒ£ OpenRouter - Gemma
    if OPENROUTER_API_KEY:
        try:
            logging.info("Attempting request with: 3. OpenRouter (Gemma)...")
            headers = {
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer": "https://t.me/Oiuhelper_bot",  # Replace with your bot's link
                "X-Title": "AI Quiz Bot"
            }
            model_identifier = "google/gemma-7b-it:free"
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                json={"model": model_identifier, "messages": [{"role": "user", "content": prompt}]},
                timeout=timeout_seconds
            )
            response.raise_for_status()
            result_text = response.json()['choices'][0]['message']['content']
            logging.info("âœ… Success with OpenRouter (Gemma).")
            return result_text
        except Exception as e:
            logging.warning(f"âŒ OpenRouter (Gemma) failed: {e}")

    # ğŸš« All models failed
    logging.error("âŒ All API providers failed. Returning empty string.")
    return ""

# -------------------------------------------------------------------
#                 OCR + language detection & translation 
# -------------------------------------------------------------------

def translate_text(text, source='en', target='ar'):
    url = 'https://libretranslate.de/translate'
    payload = {
        'q': text,
        'source': source,
        'target': target,
        'format': 'text'
    }
    try:
        response = requests.post(url, data=payload)
        return response.json()['translatedText']
    except Exception as e:
        print("ØªØ±Ø¬Ù…Ø© ÙØ´Ù„Øª:", e)
        return text  # fallback


from flask import Flask, render_template, session, request, redirect, url_for

def save_user_major(user_id, major):
    with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO users (user_id, major)
            VALUES (?, ?)
            ON CONFLICT(user_id) DO UPDATE SET major=excluded.major
        """, (user_id, major))
        conn.commit()


from langdetect import detect

def detect_language(text: str) -> str:
    try:
        lang = detect(text)
        return lang
    except:
        return "unknown"

def detect_language_from_filename(filename: str) -> str:
    """
    ÙŠØ­Ø§ÙˆÙ„ ØªØ­Ø¯ÙŠØ¯ Ø§Ù„Ù„ØºØ© Ø§Ù„Ù…Ù†Ø§Ø³Ø¨Ø© Ù…Ù† Ø§Ø³Ù… Ø§Ù„Ù…Ù„Ù.
    Ø¥Ø°Ø§ Ø§Ø­ØªÙˆÙ‰ Ø¹Ù„Ù‰ Ø­Ø±ÙˆÙ Ø¹Ø±Ø¨ÙŠØ© â†’ ÙŠØ±Ø¬Ù‘Ø­ Ø§Ù„Ø¹Ø±Ø¨ÙŠØ©.
    Ø®Ù„Ø§Ù Ø°Ù„Ùƒ â†’ Ø§Ù„Ø¥Ù†Ø¬Ù„ÙŠØ²ÙŠØ©.
    """
    for char in filename:
        if '\u0600' <= char <= '\u06FF':  # Ù†Ø·Ø§Ù‚ Ø§Ù„Ø­Ø±ÙˆÙ Ø§Ù„Ø¹Ø±Ø¨ÙŠØ©
            return "ara"
    return "eng"

import os
import fitz  # PyMuPDF
import logging
import requests
from pptx import Presentation
from tempfile import NamedTemporaryFile

# ---- OCR Space Integration ----
def extract_text_with_ocr_space(file_path: str, api_key="helloworld", language="eng") -> tuple:
    """
    Uses OCR.Space API to extract text from an image or scanned PDF.
    Returns: (text, debug_info)
    """
    url = 'https://api.ocr.space/parse/image'
    with open(file_path, 'rb') as f:
        response = requests.post(
            url,
            files={"file": f},
            data={
                "apikey": api_key,
                "language": language,
                "isOverlayRequired": False,
                "OCREngine": 2
            },
        )

    try:
        result = response.json()
        if result.get("IsErroredOnProcessing"):
            error_msg = result.get("ErrorMessage", "Unknown OCR error")
            return "", f"[OCR ERROR] {error_msg}"
        
        parsed = result.get("ParsedResults")
        if not parsed:
            return "", "[OCR ERROR] No ParsedResults returned."

        text = parsed[0].get("ParsedText", "").strip()
        return text, f"[OCR DEBUG] Length: {len(text)} | Excerpt: {text[:100]}"
    
    except Exception as e:
        return "", f"[OCR EXCEPTION] {e}"


# ---- PDF Split + OCR ----
def extract_text_from_pdf_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Splits a PDF into chunks of 3 pages (OCR.Space free limit),
    sends each chunk separately, and concatenates the extracted text.
    """
    try:
        doc = fitz.open(path)
        all_text = []
        # ØªÙ‚Ø³ÙŠÙ… ÙƒÙ„ 3 ØµÙØ­Ø§Øª ÙÙŠ Ù…Ù„Ù Ù…Ø¤Ù‚Øª
        for i in range(0, len(doc), 3):
            subdoc = fitz.open()  # Ù…Ù„Ù Ø¬Ø¯ÙŠØ¯ Ù…Ø¤Ù‚Øª
            for j in range(i, min(i+3, len(doc))):
                subdoc.insert_pdf(doc, from_page=j, to_page=j)
            
            with NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                temp_path = tmp.name
                subdoc.save(temp_path)
                subdoc.close()
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)
            
            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PDF with OCR: {e}")
        return ""


# ---- PPTX Split + OCR ----
def extract_text_from_pptx_with_ocr(path: str, api_key="helloworld", language="eng") -> str:
    """
    Converts PPTX slides into smaller chunks (3 slides per file),
    sends each chunk separately to OCR.Space, and concatenates the text.
    """
    try:
        prs = Presentation(path)
        all_text = []

        # ØªÙ‚Ø³ÙŠÙ… Ø§Ù„Ø¹Ø±Ø¶ ÙƒÙ„ 3 Ø´Ø±Ø§Ø¦Ø­
        for i in range(0, len(prs.slides), 3):
            new_ppt = Presentation()
            # Ø¥Ø¶Ø§ÙØ© ØªØ®Ø·ÙŠØ· ÙØ§Ø±Øº (Ù…Ø·Ù„ÙˆØ¨ Ù„Ø¹Ù…Ù„ Ù†Ø³Ø® Ø§Ù„Ø´Ø±Ø§Ø¦Ø­)
            blank_layout = new_ppt.slide_layouts[6]

            for j in range(i, min(i+3, len(prs.slides))):
                slide = prs.slides[j]
                new_slide = new_ppt.slides.add_slide(blank_layout)
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        textbox = new_slide.shapes.add_textbox(left=0, top=0, width=new_ppt.slide_width, height=100)
                        textbox.text = shape.text
            
            with NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                temp_path = tmp.name
                new_ppt.save(temp_path)
            
            text, debug = extract_text_with_ocr_space(temp_path, api_key=api_key, language=language)
            logging.info(f"OCR PPTX chunk [{i}-{i+2}]: {debug}")
            all_text.append(text)

            os.remove(temp_path)
        
        return "\n".join(all_text).strip()
    
    except Exception as e:
        logging.error(f"Error extracting PPTX with OCR: {e}")
        return ""
# -------------------------------------------------------------------
#                  Logging & Database Setup
# -------------------------------------------------------------------

logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s [%(levelname)s] %(message)s")


import sqlite3

def init_medical_db(db_path='medical_quizzes.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()

    cursor.execute('''
    CREATE TABLE IF NOT EXISTS quizzes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT NOT NULL,
        questions_json TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        user_id INTEGER,
        is_active BOOLEAN DEFAULT 1
    )
    ''')
    conn.commit()
    conn.close()


def init_user_quiz_db(db_path='quiz_users.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()

    
    # Ø¬Ø¯ÙˆÙ„ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ†
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
        user_id     INTEGER PRIMARY KEY,
        major       TEXT,
        native_lang TEXT DEFAULT 'ar',
        quiz_count  INTEGER DEFAULT 0,
        last_reset  TEXT
    )
    """)

    # Ø¬Ø¯ÙˆÙ„ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ø§Ù„Ù…Ù‚ØªØ±Ø­Ø© Ù…Ù† Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ† Ù„Ù„Ø¹Ø¨Ø© Ø§Ù„Ø§Ø³ØªÙ†ØªØ§Ø¬
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS inference_questions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        question TEXT NOT NULL,
        options TEXT NOT NULL,         -- Ø³ÙŠØªÙ… ØªØ®Ø²ÙŠÙ†Ù‡Ø§ ÙƒØ³Ù„Ø³Ù„Ø© JSON
        correct_index INTEGER NOT NULL,
        submitted_by INTEGER,
        approved INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS game_attempts (
        user_id INTEGER,
        game_type TEXT,
        date TEXT,
        PRIMARY KEY (user_id, game_type)
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS recent_questions (
        user_id INTEGER,
        game_type TEXT,
        question TEXT,
        added_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS user_quizzes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        quiz_data TEXT NOT NULL,
        quiz_code TEXT UNIQUE NOT NULL,
        created_at TEXT NOT NULL,
        is_active BOOLEAN DEFAULT 1
    )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS quiz_shares (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        quiz_code TEXT NOT NULL,
        shared_by_user_id INTEGER NOT NULL,
        shared_by_name TEXT,
        shared_at TEXT NOT NULL
    )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS sample_quizzes (
        quiz_code TEXT PRIMARY KEY,
        quiz_data TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)

    # Ù…Ø«Ø§Ù„ Ø¥Ø¶Ø§ÙØ© Ø£Ø¹Ù…Ø¯Ø© Ø¬Ø¯ÙŠØ¯Ø© Ø¨Ø­Ø°Ø±
    # Ø§Ù„Ø£Ø¹Ù…Ø¯Ø© Ø§Ù„Ø¬Ø¯ÙŠØ¯Ø©
    new_columns = [
        ("score", "INTEGER"),
        ("total", "INTEGER"),
        ("timestamp", "TEXT"),
        ("owner_name", "TEXT")  # Ù„ØªØ®Ø²ÙŠÙ† Ø§Ø³Ù… Ø§Ù„Ù…Ø§Ù„Ùƒ
        ]
# Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† ÙˆØ¬ÙˆØ¯ Ø§Ù„Ø¹Ù…ÙˆØ¯ Ù‚Ø¨Ù„ Ø¥Ø¶Ø§ÙØªÙ‡
    for col_name, col_type in new_columns:
        cursor.execute(f"PRAGMA table_info(user_quizzes)")
        existing_cols = [row[1] for row in cursor.fetchall()]
        if col_name not in existing_cols:
            cursor.execute(f"ALTER TABLE user_quizzes ADD COLUMN {col_name} {col_type}")

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS bot_users (
        user_id INTEGER PRIMARY KEY,
        is_channel_user BOOLEAN DEFAULT 0, -- Ø£Ø¹Ø¶Ø§Ø¡ Ø§Ù„Ù‚Ù†ÙˆØ§Øª
        is_external_user BOOLEAN DEFAULT 0 -- Ù…Ø³ØªØ®Ø¯Ù… Ø¬Ø¯ÙŠØ¯ Ø®Ø§Ø±Ø¬ÙŠ
    )
    """)
    
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS stat (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0,
        total_users INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS daily_stats (
        date TEXT PRIMARY KEY,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0,
        new_users INTEGER DEFAULT 0
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS top_users (
        user_id INTEGER PRIMARY KEY,
        tests_generated INTEGER DEFAULT 0,
        files_processed INTEGER DEFAULT 0
    )
    """)
    conn.commit()
    conn.close()


def init_request_db(db_path='requests.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cur = conn.cursor()
    cur.execute('''
        CREATE TABLE IF NOT EXISTS requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            username TEXT,
            file_id TEXT,
            status TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()


import sqlite3, logging

def fetch_user_major(uid, db_path="quiz_users.db"):
    try:
        with sqlite3.connect(db_path, check_same_thread=False) as conn:
            cur = conn.cursor()
            cur.execute("SELECT major FROM users WHERE user_id=?", (uid,))
            row = cur.fetchone()
        return row[0] if row else "General"
    except Exception:
        logging.exception("fetch_user_major failed")
        return "General"



def init_all_dbs():
    init_medical_db()
    init_user_quiz_db()
    init_request_db()


# ---------------------------
# ---- get basic statics -----
# -------------    ---------

import sqlite3

DB_NAME = "quiz_users.db"

def get_tests_generated():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT tests_generated FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


def get_files_processed():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT files_processed FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


def get_total_users():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    cursor.execute("SELECT total_users FROM stat LIMIT 1")
    result = cursor.fetchone()
    conn.close()
    return result[0] if result else 0


# ----------------------------
# ----    flask config  ----------------------------

# ÙˆØ§Ø¬Ù‡Ø© Flask Ù„Ù„ÙØ­Øµ
app = Flask(__name__)

@app.route('/')
def home():
    return render_template('index.html')
    

@app.route('/quiz/<int:quiz_id>')
def show_quiz(quiz_id):
    conn = sqlite3.connect('medical_quizzes.db')
    cursor = conn.cursor()
    
    cursor.execute('SELECT title, questions_json FROM quizzes WHERE id = ?', (quiz_id,))
    quiz = cursor.fetchone()
    conn.close()
    
    if quiz:
        quiz_data = {
            'title': quiz[0],
            'questions': json.loads(quiz[1])
        }
        return render_template('quiz.html', quiz=quiz_data)
    else:
        return "Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± ØºÙŠØ± Ù…ÙˆØ¬ÙˆØ¯", 404




@app.route('/supportme')
def supportme():
    # Ù‡Ù†Ø§ ØªÙ‚ÙˆÙ… Ø¨Ø¬Ù„Ø¨ Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ø§Ù„Ø­Ù‚ÙŠÙ‚ÙŠØ© Ù…Ù† Ù‚Ø§Ø¹Ø¯Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ø£Ùˆ Ø£ÙŠ Ù…ØµØ¯Ø± Ø¢Ø®Ø±
    # Ù‡Ø°Ù‡ Ù…Ø¬Ø±Ø¯ Ù‚ÙŠÙ… Ø§ÙØªØ±Ø§Ø¶ÙŠØ©
    tests_generated = get_tests_generated()
    files_processed = get_files_processed()
    total_users = get_total_users()

    return render_template(
        'supportme.html',
        tests_generated=tests_generated,
        files_processed=files_processed,
        total_users=total_users
    )








# track temporary state for custom-major input
user_states = {}
usage_count = {}
import sqlite3, logging




# quene

def save_request(msg, db_path='requests.db'):
    try:
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        file_id = getattr(getattr(msg, "document", None), "file_id", None)
        cur.execute(
            'INSERT INTO requests (user_id, username, file_id, status) VALUES (?, ?, ?, ?)',
            (msg.from_user.id, msg.from_user.username or "", file_id, 'pending')
        )
        conn.commit()
    except Exception:
        logging.exception("save_request failed")
    finally:
        try: conn.close()
        except: pass

def update_request_status(file_id, new_status, db_path='requests.db'):
    try:
        if not file_id:
            return
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        cur.execute('UPDATE requests SET status=? WHERE file_id=?', (new_status, file_id))
        conn.commit()
    except Exception:
        logging.exception("update_request_status failed")
    finally:
        try: conn.close()
        except: pass



# Ø¯Ø§Ù„Ø© Ø§Ù„Ø¹Ø§Ù…Ù„
def worker():
    while True:
        item = request_queue.get()
        try:
            if isinstance(item, tuple) and len(item) == 2:
                msg, sent_msg = item
                process_message(msg, message_id=sent_msg.message_id, chat_id=sent_msg.chat.id)
            else:
                # Ø¹Ù†ØµØ± ÙˆØ§Ø­Ø¯: ÙÙ‚Ø· Ø§Ù„Ø±Ø³Ø§Ù„Ø©
                msg = item
                process_message(msg)
        except Exception:
            logging.exception("[WORKER ERROR]")
        finally:
            request_queue.task_done()
# ØªØ´ØºÙŠÙ„ Ø§Ù„Ø¹Ù…Ø§Ù„ ÙÙŠ Ø§Ù„Ø®Ù„ÙÙŠØ©
def start_workers():
    for _ in range(num_workers):
        threading.Thread(target=worker, daemon=True).start()
    logging.info("Workers started: %s", num_workers)




def safe_edit_or_send(text, chat_id, message_id, parse_mode="HTML"):
    try:
        if chat_id and message_id:
            return bot.edit_message_text(
                text, chat_id=chat_id, message_id=message_id, parse_mode=parse_mode
            )
    except Exception as e:
        logging.warning("edit_message_text failed (%s), fallback to send_message", e)
    return bot.send_message(chat_id, text, parse_mode=parse_mode)



def is_request_already_queued(file_id=None, user_id=None, message_id=None, db_path='requests.db'):
    try:
        conn = sqlite3.connect(db_path, check_same_thread=False)
        cur = conn.cursor()
        if file_id:
            cur.execute("SELECT status FROM requests WHERE file_id=? ORDER BY id DESC LIMIT 1", (file_id,))
        else:
            cur.execute("SELECT status FROM requests WHERE user_id=? AND message_id=? ORDER BY id DESC LIMIT 1", (user_id, message_id))
        row = cur.fetchone()
        conn.close()
        if row and row[0] in ('pending','processing'):
            return True
        return False
    except Exception:
        logging.exception("is_request_already_queued failed")
        return False


# ------------------------------------
# ------------------------------------------------- Stat Management----
# ----------------------------------------------------------------------

def add_external_user(uid: int):
    with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (uid,))
        if cursor.fetchone()[0] == 0:
            cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (uid, 1))
            conn.commit()


def update_files_and_users(uid: int = None, files_count: int = 1):
    """
    - files_count: Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ù„ÙØ§Øª Ø§Ù„Ø¬Ø¯ÙŠØ¯Ø©
    - uid: user_id Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… Ø§Ù„Ø¬Ø¯ÙŠØ¯ Ø§Ù„Ø®Ø§Ø±Ø¬ÙŠ (Ø§Ø®ØªÙŠØ§Ø±ÙŠ)
    """
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # 1) ØªØ­Ø¯ÙŠØ« Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ø§Ù„Ø¬Ø¯ÙŠØ¯ Ø¥Ø°Ø§ ØªÙ… ØªÙ…Ø±ÙŠØ± uid
            if uid:
                cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (uid,))
                if cursor.fetchone()[0] == 0:
                    # Ø¥Ø¶Ø§ÙØ© Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ø§Ù„Ø®Ø§Ø±Ø¬ÙŠ Ø§Ù„Ø¬Ø¯ÙŠØ¯
                    cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (uid, 1))

            # 2) Ø­Ø³Ø§Ø¨ Ø¹Ø¯Ø¯ Ù…Ø³ØªØ®Ø¯Ù…ÙŠ Ø§Ù„Ù‚Ù†ÙˆØ§Øª
            env_channels = os.getenv("ALLOWED_CHANNELS", "")
            channel_users_count = 0
            if env_channels.strip():
                allowed_channels = set(map(int, env_channels.split(",")))
                for channel_id in allowed_channels:
                    try:
                        chat = bot.get_chat(channel_id)
                        channel_users_count += chat.get_members_count()
                    except Exception as e:
                        logging.warning(f"âš ï¸ Failed to fetch channel members for {channel_id}: {e}")

            # 3) Ø­Ø³Ø§Ø¨ Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ† Ø§Ù„Ø®Ø§Ø±Ø¬ÙŠÙŠÙ† Ø§Ù„Ù…Ø³Ø¬Ù„ÙŠÙ† ÙÙŠ Ø§Ù„Ø¬Ø¯ÙˆÙ„
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_external_user=1")
            external_users_count = cursor.fetchone()[0]

            total_users = channel_users_count + external_users_count

            # 4) ØªØ­Ø¯ÙŠØ« Ø§Ù„Ø¥Ø­ØµØ§Ø¦ÙŠØ§Øª Ø§Ù„Ø¹Ø§Ù…Ø© ÙÙŠ Ø¬Ø¯ÙˆÙ„ stat
            cursor.execute("""
                UPDATE stat
                SET files_processed = files_processed + ?,
                    total_users = ?
                WHERE id = 1
            """, (files_count, total_users))

            conn.commit()
    except Exception as e:
        logging.error(f"âŒ Error updating files and users: {e}")
        

def update_daily_stats(date_str: str = None, tests: int = 0, files: int = 0):
    """
    ØªØ­Ø¯ÙŠØ« Ø§Ù„Ø¥Ø­ØµØ§Ø¦ÙŠØ§Øª Ø§Ù„ÙŠÙˆÙ…ÙŠØ©:
    - ÙŠØ­Ø³Ø¨ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ† Ø§Ù„Ø¬Ø¯Ø¯ ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ Ù…Ù† Ø¬Ø¯ÙˆÙ„ users
    """
    import datetime
    date_str = date_str or datetime.date.today().isoformat()
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # Ø­Ø³Ø§Ø¨ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ† Ø§Ù„Ø¬Ø¯Ø¯ Ø­Ø³Ø¨ Ø§Ù„Ù†ÙˆØ¹
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_channel_user=1")
            channel_users = cursor.fetchone()[0]

            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE is_external_user=1")
            external_users = cursor.fetchone()[0]

            total_new_users = channel_users + external_users

            # Ø¥Ø¯Ø±Ø§Ø¬ Ø£Ùˆ ØªØ­Ø¯ÙŠØ« Ø³Ø¬Ù„ Ø§Ù„ÙŠÙˆÙ…
            cursor.execute("SELECT COUNT(*) FROM daily_stats WHERE date=?", (date_str,))
            if cursor.fetchone()[0] == 0:
                cursor.execute("""
                    INSERT INTO daily_stats (date, tests_generated, files_processed, new_users, channel_users, external_users)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (date_str, tests, files, total_new_users, channel_users, external_users))
            else:
                cursor.execute("""
                    UPDATE daily_stats
                    SET tests_generated = tests_generated + ?,
                        files_processed = files_processed + ?,
                        new_users = ?,
                        channel_users = ?,
                        external_users = ?
                    WHERE date = ?
                """, (tests, files, total_new_users, channel_users, external_users, date_str))

            conn.commit()
    except Exception as e:
        logging.error(f"âŒ Error updating daily stats: {e}")
    

def update_top_user(user_id: int, tests: int = 0, files: int = 0):
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()

            # ØªØ£ÙƒØ¯ Ø£Ù† Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ù…ÙˆØ¬ÙˆØ¯ ÙÙŠ Ø¬Ø¯ÙˆÙ„ users
            cursor.execute("SELECT COUNT(*) FROM bot_users WHERE user_id=?", (user_id,))
            if cursor.fetchone()[0] == 0:
                # Ø¥Ø¶Ø§ÙØ© Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ ÙƒÙ…Ø³ØªØ®Ø¯Ù… Ø®Ø§Ø±Ø¬ÙŠ
                cursor.execute("INSERT INTO bot_users (user_id, is_external_user) VALUES (?, ?)", (user_id, 1))

            # ØªØ­Ø¯ÙŠØ« Ø¬Ø¯ÙˆÙ„ top_users
            cursor.execute("SELECT COUNT(*) FROM top_users WHERE user_id=?", (user_id,))
            if cursor.fetchone()[0] == 0:
                cursor.execute("INSERT INTO top_users (user_id, tests_generated, files_processed) VALUES (?, ?, ?)",
                               (user_id, tests, files))
            else:
                cursor.execute("""
                    UPDATE top_users
                    SET tests_generated = tests_generated + ?,
                        files_processed = files_processed + ?
                    WHERE user_id = ?
                """, (tests, files, user_id))

            conn.commit()
    except Exception as e:
        logging.error(f"âŒ Error updating top user: {e}")
# -------------------------------------------------------------------
#                     Text Extraction & OCR
# -------------------------------------------------------------------

def extract_text_from_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        text = "\n".join([page.get_text() for page in doc])
        return text.strip()
    except Exception as e:
        logging.error(f"Error extracting PDF text: {e}")
        return ""
    # fallback to PyMuPDF text extraction
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])
# Ø£Ø¶Ù Ù‡Ø°Ù‡ Ø§Ù„Ø¯Ø§Ù„Ø© ÙÙŠ Ù‚Ø³Ù… Text Extraction & OCR
def extract_text_from_docx(path: str) -> str:
    try:
        doc = docx.Document(path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting DOCX text: {e}")
        return ""

# ÙˆÙŠØ¬Ø¨ Ø£ÙŠØ¶Ø§Ù‹ ØªØ¹Ø±ÙŠÙ Ø¯Ø§Ù„Ø© Ù„Ù…Ù„ÙØ§Øª txt
def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logging.error(f"Error extracting TXT text: {e}")
        return ""
        
def is_text_empty(text: str) -> bool:
    return not text or len(text.strip()) < 30  # ÙŠÙ…ÙƒÙ† ØªØ¹Ø¯ÙŠÙ„ Ø§Ù„Ø­Ø¯ Ø­Ø³Ø¨ ØªØ¬Ø±Ø¨ØªÙƒ


def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text).strip()
    except Exception as e:
        logging.error(f"Error extracting PPTX text: {e}")
        return ""

def split_text(content, chunk_size=3500):
    return [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]

def summarize_long_text(content: str) -> str:
    """
    Summarize the content in its original language (Arabic or English) using educational style.
    """
    lang = detect_language(content[:1000])  # Ù†ÙƒØªÙÙŠ Ø¨Ø£ÙˆÙ„ 1000 Ø­Ø±Ù Ù„Ù„ØªØ­Ù„ÙŠÙ„
    print(f"[DEBUG] Detected language: {lang}")

    if lang.startswith("ar"):
        summary_prompt = (
            "Ø£Ù†Øª Ù…Ø³Ø§Ø¹Ø¯ ØªØ¹Ù„ÙŠÙ…ÙŠ Ù…Ø­ØªØ±Ù. Ù‚Ù… Ø¨ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø§Ù„ØªØ§Ù„ÙŠ Ø¨Ø£Ø³Ù„ÙˆØ¨ ØªØ¹Ù„ÙŠÙ…ÙŠ Ù…Ù†Ø¸Ù… ÙˆÙˆØ§Ø¶Ø­ Ø¨Ø§Ù„Ù„ØºØ© Ø§Ù„Ø¹Ø±Ø¨ÙŠØ© ÙÙ‚Ø·ØŒ"
            " Ù…Ø¹ Ø§Ù„Ø­ÙØ§Ø¸ Ø¹Ù„Ù‰ Ø§Ù„Ù†Ù‚Ø§Ø· Ø§Ù„Ù…ÙÙŠØ¯Ø© ÙˆØ§Ù„ØªÙØ§ØµÙŠÙ„ Ø§Ù„ØªÙŠ ÙŠÙ…ÙƒÙ† Ø§Ø³ØªØ®Ø¯Ø§Ù…Ù‡Ø§ Ù„Ø§Ø­Ù‚Ù‹Ø§ Ù„ØµÙ†Ø¹ Ø£Ø³Ø¦Ù„Ø© Ø£Ùˆ Ø¨Ø·Ø§Ù‚Ø§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ©."
            "\n\nØ§Ù„Ù…Ø­ØªÙˆÙ‰:\n{chunk}"
        )
        merge_prompt = (
            "ÙÙŠÙ…Ø§ ÙŠÙ„ÙŠ Ù…Ø¬Ù…ÙˆØ¹Ø© Ù…Ù† Ø§Ù„Ù…Ù„Ø®ØµØ§Øª Ø§Ù„Ø¬Ø²Ø¦ÙŠØ© Ù„Ù…Ø­ØªÙˆÙ‰ ØªØ¹Ù„ÙŠÙ…ÙŠ. Ù‚Ù… Ø¨Ø¯Ù…Ø¬Ù‡Ø§ ÙÙŠ Ù…Ù„Ø®Øµ Ù†Ù‡Ø§Ø¦ÙŠ Ø´Ø§Ù…Ù„ ÙˆÙ…ØªØ±Ø§Ø¨Ø· ÙˆÙˆØ§Ø¶Ø­"
            " Ø¨Ø§Ù„Ù„ØºØ© Ø§Ù„Ø¹Ø±Ø¨ÙŠØ©ØŒ Ù…Ø¹ Ø§Ù„Ø­ÙØ§Ø¸ Ø¹Ù„Ù‰ Ø§Ù„ØªÙØ§ØµÙŠÙ„ Ø§Ù„Ù…ÙÙŠØ¯Ø© Ø§Ù„ØªÙŠ ØªØ³Ø§Ø¹Ø¯ Ø¹Ù„Ù‰ ÙÙ‡Ù… Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø£Ùˆ Ø¥Ø¹Ø¯Ø§Ø¯ Ø§Ø®ØªØ¨Ø§Ø± Ù…Ù†Ù‡.\n\n{merged}"
        )
    else:
        summary_prompt = (
            "You are a professional educational assistant. Summarize the following content into a clear and concise educational explanation in **English only**.\n\n"
            "- Preserve factual details and key definitions.\n"
            "- Avoid vague sentences or repetition.\n"
            "- Keep the original language (do not translate).\n\nContent:\n{chunk}"
        )
        merge_prompt = (
            "You are an educational summarizer. Merge the following partial summaries into one final, well-structured summary in **English**, preserving all useful learning content.\n\n{merged}"
        )

    chunks = split_text(content)
    partial_summaries = []

    for i, chunk in enumerate(chunks):
        prompt = summary_prompt.format(chunk=chunk)
        summary = generate_smart_response(prompt.strip())
        partial_summaries.append(summary)

    merged_summary = "\n".join(partial_summaries)
    final_prompt = merge_prompt.format(merged=merged_summary)
    return generate_smart_response(final_prompt.strip())
    

def parse_ai_json(raw_text: str) -> dict | None:

    # 1. ÙÙƒÙ‘Ù Ù‡Ø§Ø±Ø¨Ø§Øª Unicode (\u0627 â†’ Ø§)
    def _unescape(match):
        code = match.group(1)
        return chr(int(code, 16))
    text = re.sub(r'\\u([0-9A-Fa-f]{4})', _unescape, raw_text)

    # 2. Ø§Ø¬ØªØ²Ø¡ Ø£ÙˆÙ„ ÙƒØªÙ„Ø© JSON (Ù…Ù† { Ø¥Ù„Ù‰ })
    m = re.search(r'\{[\s\S]*\}', text)
    json_text = m.group(0) if m else text

    # 3. Ø­Ø§ÙˆÙ„ Ø§Ù„ØªØ­Ù…ÙŠÙ„
    for attempt in (json_text, text):
        try:
            data = json.loads(attempt)
            break
        except json.JSONDecodeError:
            data = None
    if not data:
        return None

    # 4. Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø¨Ù†ÙŠØ© Ø§Ù„Ù€ dict
    if not all(k in data for k in ("question", "options", "correct_index")):
        return None

    # 5. Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† Ø£Ù† options Ù‚Ø§Ø¦Ù…Ø© ÙˆØµØ§Ù„Ø­Ø©
    if not isinstance(data["options"], list) or len(data["options"]) < 2:
        return None

    # 6. Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† correct_index
    ci = data["correct_index"]
    if not isinstance(ci, int) or ci < 0 or ci >= len(data["options"]):
        return None

    return data

def generate_game(prompt, user_id=0, translate_all=False, translate_question=False):
    if user_id == ADMIN_ID or can_generate(user_id):  # <-- Ø§Ù„ØªØ­Ù‚Ù‚ Ù‡Ù†Ø§
        raw_response = generate_smart_response(prompt)
    else:
        raw_response = generate_gemini_response(prompt)
        
    game_data = parse_ai_json(raw_response)

    if not game_data:
        raise ValueError("ÙØ´Ù„ Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø¨ÙŠØ§Ù†Ø§Øª Ø§Ù„Ù„Ø¹Ø¨Ø©")

    if translate_all:
        # ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø³Ø¤Ø§Ù„
        if 'question' in game_data:
            game_data['question'] = translate_text(game_data['question'], source='en', target='ar')

        # ØªØ±Ø¬Ù…Ø© ÙƒÙ„ Ø§Ù„Ø®ÙŠØ§Ø±Ø§Øª
        if 'options' in game_data and isinstance(game_data['options'], list):
            game_data['options'] = [
                translate_text(option, source='en', target='ar') for option in game_data['options']
            ]

    elif translate_question:
        # ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø³Ø¤Ø§Ù„ ÙÙ‚Ø·
        if 'question' in game_data:
            game_data['question'] = translate_text(game_data['question'], source='en', target='ar')

    return game_data

import genanki
import uuid
import requests
import hashlib
import os
import logging
import tempfile
import re
from typing import List, Dict, Tuple
from PIL import Image

# -----------------------
# Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø±Ø§Ø¨Ø· Ø£ÙˆÙ„ ØµÙˆØ±Ø© Ù…Ù† HTML
# -----------------------
IMG_TAG_RE = re.compile(r'<img [^>]*src=[\'"]([^\'"]+)[\'"][^>]*>', re.IGNORECASE)

def _extract_first_img_url_from_html(html: str) -> Tuple[str, str]:
    """
    Returns (image_url, cleaned_html_without_img)
    """
    match = IMG_TAG_RE.search(html)
    if not match:
        return "", html
    url = match.group(1)
    # Ø¥Ø²Ø§Ù„Ø© Ø¬Ù…ÙŠØ¹ ÙˆØ³ÙˆÙ… <img> Ù„ØªØ¬Ù†Ø¨ Ø§Ù„ØªÙƒØ±Ø§Ø±
    cleaned = IMG_TAG_RE.sub("", html)
    return url, cleaned.strip()


# -----------------------
# ØªÙ†Ø²ÙŠÙ„ Ø§Ù„ØµÙˆØ±Ø© Ø¨Ø§Ø³ØªØ®Ø¯Ø§Ù… Pillow
# -----------------------
def _download_image_to_dir(url: str, dest_dir: str) -> Tuple[str, str]:
    """
    Download image from URL into dest_dir.
    Returns (basename, full_path) or (None, None) on failure.
    Ensures filename is ASCII (md5 hash + extension).
    """
    try:
        r = requests.get(url, timeout=15)
        r.raise_for_status()
        content = r.content

        # Ø­ÙØ¸ Ù…Ø¤Ù‚Øª Ù„Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù†ÙˆØ¹
        temp_path = os.path.join(dest_dir, "temp_image")
        with open(temp_path, "wb") as f:
            f.write(content)

        # Ø§Ø³ØªØ®Ø¯Ø§Ù… Pillow Ù„Ø§ÙƒØªØ´Ø§Ù Ù†ÙˆØ¹ Ø§Ù„ØµÙˆØ±Ø©
        try:
            with Image.open(temp_path) as img:
                ext = img.format.lower()  # jpeg, png, gif, ...
        except Exception as e:
            logging.warning(f"âš ï¸ Could not identify image type for {url}: {e}")
            os.remove(temp_path)
            return None, None

        # Ø¥Ù†Ø´Ø§Ø¡ Ø§Ø³Ù… ÙØ±ÙŠØ¯ Ù„Ù„Ù…Ù„Ù Ø¨Ø§Ù„Ù€ md5 + Ø§Ù„Ø§Ù…ØªØ¯Ø§Ø¯ Ø§Ù„ØµØ­ÙŠØ­
        fname = hashlib.md5(url.encode()).hexdigest() + f".{ext}"
        final_path = os.path.join(dest_dir, fname)
        os.rename(temp_path, final_path)

        return fname, final_path

    except Exception as e:
        logging.warning(f"âš ï¸ Failed to download image {url}: {e}")
        return None, None



import tempfile # <--- ØªØ£ÙƒØ¯ Ù…Ù† Ø§Ø³ØªÙŠØ±Ø§Ø¯ Ù‡Ø°Ù‡ Ø§Ù„Ù…ÙƒØªØ¨Ø© ÙÙŠ Ø£Ø¹Ù„Ù‰ Ø§Ù„Ù…Ù„Ù

# -----------------------
# Ø­ÙØ¸ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª ÙÙŠ Ù…Ù„Ù Anki Ù…Ø¹ Ø¯Ø¹Ù… Ø§Ù„ØµÙˆØ± (Ø§Ù„Ù†Ø³Ø®Ø© Ø§Ù„Ù…ØµØ­Ø­Ø©)
# -----------------------
def save_cards_to_apkg(cards: List[Dict], filename: str = 'anki_flashcards.apkg', deck_name: str = "My Flashcards"):
    model = genanki.Model(
        1607392319,
        'Simple Model with Tags',
        fields=[
            {'name': 'Front'},
            {'name': 'Back'},
            {'name': 'Tag'}
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Front}}<br><small style="color:gray">{{Tag}}</small>',
                'afmt': '{{FrontSide}}<hr id="answer">{{Back}}',
            },
        ]
    )

    deck = genanki.Deck(
        deck_id=int(str(uuid.uuid4().int)[:9]),
        name=deck_name
    )

    seen = set()
    media_files = []

    # âœ¨ Ø§Ù„ØªØºÙŠÙŠØ± Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠ: Ø§Ø³ØªØ®Ø¯Ø§Ù… Ù…Ø¬Ù„Ø¯ Ù…Ø¤Ù‚Øª Ø®Ø§Øµ ÙˆÙ…Ù†Ø¹Ø²Ù„ Ù„ÙƒÙ„ Ø¹Ù…Ù„ÙŠØ©
    # Ø³ÙŠØªÙ… Ø­Ø°Ù Ù‡Ø°Ø§ Ø§Ù„Ù…Ø¬Ù„Ø¯ ÙˆÙ…Ø­ØªÙˆÙŠØ§ØªÙ‡ ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ Ø¨Ø¹Ø¯ Ø§Ù„Ø®Ø±ÙˆØ¬ Ù…Ù† Ù‡Ø°Ø§ Ø§Ù„Ø¨Ù„ÙˆÙƒ
    with tempfile.TemporaryDirectory() as temp_dir:
        logging.info(f"Created temporary directory for media: {temp_dir}")

        for card in cards:
            front = card.get('front', '').strip()
            back = card.get('back', '').strip()
            tag = card.get('tag', '').strip()
            image_url = card.get('image_url', '')

            if front and back and front not in seen:
                # Ø¥Ø°Ø§ ÙƒØ§Ù†Øª Ù‡Ù†Ø§Ùƒ ØµÙˆØ±Ø©ØŒ Ù‚Ù… Ø¨ØªÙ†Ø²ÙŠÙ„Ù‡Ø§ Ø¥Ù„Ù‰ Ø§Ù„Ù…Ø¬Ù„Ø¯ Ø§Ù„Ù…Ø¤Ù‚Øª Ø§Ù„Ø®Ø§Øµ Ø¨Ù†Ø§
                if image_url:
                    # Ù†Ù…Ø±Ø± Ø§Ù„Ù…Ø¬Ù„Ø¯ Ø§Ù„Ù…Ø¤Ù‚Øª temp_dir Ù„Ù„Ø¯Ø§Ù„Ø©
                    fname, path = _download_image_to_dir(image_url, temp_dir)
                    if fname and path:
                        media_files.append(path)
                        back += f"<br><img src='{fname}' style='max-height:220px; display:block; margin:12px auto;'>"

                note = genanki.Note(model=model, fields=[front, back, tag])
                deck.add_note(note)
                seen.add(front)

        package = genanki.Package(deck)
        if media_files:
            package.media_files = media_files

        package.write_to_file(filename)

    # Ù„Ø§ Ø­Ø§Ø¬Ø© Ù„Ø­Ø°Ù Ø§Ù„Ù…Ù„ÙØ§Øª ÙŠØ¯ÙˆÙŠÙ‹Ø§ØŒ Ø³ÙŠØªÙ… Ø­Ø°Ù temp_dir ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§
    logging.info(f"Successfully created Anki package: {filename}")
    return filename



def parse_manual_anki_input(text):
    cards = []
    lines = [line.strip() for line in text.strip().split('\n')]
    current_card = []

    for line in lines:
        if line == "":
            if current_card:
                if len(current_card) >= 2:
                    card = {
                        "front": current_card[0],
                        "back": current_card[1],
                        "tag": current_card[2] if len(current_card) > 2 else ""
                    }
                    cards.append(card)
                current_card = []
        else:
            current_card.append(line)

    if current_card and len(current_card) >= 2:
        card = {
            "front": current_card[0],
            "back": current_card[1],
            "tag": current_card[2] if len(current_card) > 2 else ""
        }
        cards.append(card)

    return cards
    
# -------------------------------------------------------------------
#                     Quota Management
# -------------------------------------------------------------------
def add_recent_question(user_id, game_type, question):
    with sqlite3.connect("quiz_users.db") as conn:
        cursor = conn.cursor()
        
        # Ø¥Ø¯Ø®Ø§Ù„ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
        cursor.execute("""
        INSERT INTO recent_questions (user_id, game_type, question) 
        VALUES (?, ?, ?)
        """, (user_id, game_type, question))
        
        # Ø­Ø°Ù Ø§Ù„Ø£Ù‚Ø¯Ù… Ø¥Ø°Ø§ ØªØ¬Ø§ÙˆØ² 10 Ø£Ø³Ø¦Ù„Ø©
        cursor.execute("""
        DELETE FROM recent_questions
        WHERE user_id = ? AND game_type = ?
        AND question NOT IN (
            SELECT question FROM recent_questions
            WHERE user_id = ? AND game_type = ?
            ORDER BY added_at DESC
            LIMIT 10
        )
        """, (user_id, game_type, user_id, game_type))

        conn.commit()

def get_recent_questions(user_id, game_type):
    with sqlite3.connect("quiz_users.db") as conn:
        cursor = conn.cursor()
        cursor.execute("""
        SELECT question FROM recent_questions
        WHERE user_id = ? AND game_type = ?
        ORDER BY added_at DESC
        LIMIT 10
        """, (user_id, game_type))
        rows = cursor.fetchall()
        return [row[0] for row in rows]


def reset_if_needed(user_id: int):
    this_month = datetime.now().strftime("%Y-%m")
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            
            cursor.execute("SELECT last_reset FROM users WHERE user_id = ?", (user_id,))
            row = cursor.fetchone()
            
            if not row or row[0] != this_month:
                cursor.execute("""
                    INSERT OR REPLACE INTO users(user_id, major, quiz_count, last_reset)
                    VALUES (?, COALESCE((SELECT major FROM users WHERE user_id=?), ''), 0, ?)
                """, (user_id, user_id, this_month))
                conn.commit()
    except Exception as e:
        logging.error(f"ğŸš« Ø®Ø·Ø£ ÙÙŠ reset_if_needed Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… {user_id}: {e}")






MAX_FREE_ATTEMPTS = 3  # ğŸ‘ˆ Ø¹Ø¯Ù„Ù‡Ø§ Ø­Ø³Ø¨ Ù…Ø§ ØªØ±ÙŠØ¯
def can_generate(user_id: int) -> bool:
    # Ø§Ù„Ø³Ù…Ø§Ø­ Ù„Ù„Ø£Ø¯Ù…Ù† Ø¯Ø§Ø¦Ù…Ø§Ù‹
    if user_id == ADMIN_ID:
        return True
    
    # Ø§Ø³ØªØ®Ø¯Ø§Ù… 'with' ÙŠØ¶Ù…Ù† ÙØªØ­ ÙˆØ¥ØºÙ„Ø§Ù‚ Ø§Ù„Ø§ØªØµØ§Ù„ Ø¨Ø´ÙƒÙ„ Ø¢Ù…Ù†
    try:
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            
            today = datetime.utcnow().date()
            cursor.execute("SELECT quiz_count, last_reset FROM users WHERE user_id = ?", (user_id,))
            row = cursor.fetchone()
            
            quiz_count = 0  # Ù‚ÙŠÙ…Ø© Ø§ÙØªØ±Ø§Ø¶ÙŠØ©

            if not row:
                # Ù…Ø³ØªØ®Ø¯Ù… Ø¬Ø¯ÙŠØ¯
                cursor.execute("INSERT INTO users (user_id, quiz_count, last_reset) VALUES (?, ?, ?)", (user_id, 0, today.isoformat()))
                conn.commit()
                quiz_count = 0
            else:
                quiz_count, last_reset = row
                last_reset_date = None
                try:
                    if last_reset:
                        last_reset_date = datetime.fromisoformat(last_reset).date()
                except Exception as e:
                    logging.warning(f"âš ï¸ last_reset ØºÙŠØ± ØµØ§Ù„Ø­Ø© Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… {user_id}: {e}")
                    last_reset_date = None
                
                # ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ø¯Ø®ÙˆÙ„ ÙÙŠ Ø´Ù‡Ø± Ø¬Ø¯ÙŠØ¯
                if not last_reset_date or last_reset_date.month != today.month or last_reset_date.year != today.year:
                    cursor.execute("UPDATE users SET quiz_count = 0, last_reset = ? WHERE user_id = ?", (today.isoformat(), user_id))
                    conn.commit()
                    quiz_count = 0
            
            # Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø§Øª
            if quiz_count >= MAX_FREE_ATTEMPTS:
                # Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù‚Ù†ÙˆØ§Øª Ø§Ù„Ù…Ø³Ù…ÙˆØ­ Ø¨Ù‡Ø§ Ø¥Ø°Ø§ ØªØ¬Ø§ÙˆØ² Ø§Ù„Ø­Ø¯
                try:
                    raw = os.getenv("ALLOWED_CHANNELS", "")
                    if not raw.strip(): return False # Ø¥Ø°Ø§ Ù„Ù… ØªÙƒÙ† Ù‡Ù†Ø§Ùƒ Ù‚Ù†ÙˆØ§ØªØŒ Ù„Ø§ ØªØ³Ù…Ø­
                    
                    allowed_channels = set(int(cid) for cid in raw.split(",") if cid.strip())
                    for channel_id in allowed_channels:
                        try:
                            member = bot.get_chat_member(chat_id=channel_id, user_id=user_id)
                            if member.status in ['member', 'administrator', 'creator']:
                                return True  # Ù…Ø³ØªØ®Ø¯Ù… Ù…Ù…ÙŠØ²
                        except Exception as e:
                            logging.warning(f"âš ï¸ ÙØ´Ù„ Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù‚Ù†Ø§Ø© {channel_id} Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… {user_id}: {e}")
                    return False  # Ù„ÙŠØ³ Ø¹Ø¶ÙˆÙ‹Ø§ ÙÙŠ Ø£ÙŠ Ù‚Ù†Ø§Ø© Ù…Ø³Ù…ÙˆØ­ Ø¨Ù‡Ø§
                except Exception as e:
                    logging.error(f"ğŸš« Ø®Ø·Ø£ ÙÙŠ Ù‚Ø±Ø§Ø¡Ø© Ø§Ù„Ù‚Ù†ÙˆØ§Øª Ø§Ù„Ù…Ø³Ù…ÙˆØ­ Ø¨Ù‡Ø§: {e}")
                    return False
            else:
                return True  # Ø¶Ù…Ù† Ø§Ù„Ø­Ø¯ Ø§Ù„Ù…Ø³Ù…ÙˆØ­ Ø¨Ù‡
    
    except Exception as e:
        logging.error(f"ğŸš« Ø®Ø·Ø£ ÙØ§Ø¯Ø­ ÙÙŠ Ø¯Ø§Ù„Ø© can_generate: {e}")
        return False
        


from datetime import datetime


def increment_count(user_id: int):
    # Ù„Ø§ ØªÙ‚Ù… Ø¨Ø²ÙŠØ§Ø¯Ø© Ø§Ù„Ø¹Ø¯Ø§Ø¯ Ø¥Ø°Ø§ ÙƒØ§Ù† Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ù‡Ùˆ Ø§Ù„Ø£Ø¯Ù…Ù† Ø£Ùˆ Ù…Ø³ØªØ®Ø¯Ù… Ù…Ù…ÙŠØ²
    if user_id == ADMIN_ID:
        bot.send_message(ADMIN_ID, "âœ¨ (ÙˆØ¶Ø¹ Ø§Ù„Ø£Ø¯Ù…Ù†: Ù„Ù… ÙŠØªÙ… Ø§Ø­ØªØ³Ø§Ø¨ Ù‡Ø°Ù‡ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø©)")
        return
    
    try:
        # Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù‚Ù†ÙˆØ§Øª Ø§Ù„Ù…Ù…ÙŠØ²Ø©
        raw = os.getenv("ALLOWED_CHANNELS", "")
        allowed_channels = set(int(cid) for cid in raw.split(",") if cid.strip())
        for channel_id in allowed_channels:
            try:
                member = bot.get_chat_member(chat_id=channel_id, user_id=user_id)
                if member.status in ['member', 'administrator', 'creator']:
                    return  # Ù…Ø³ØªØ®Ø¯Ù… Ù…Ù…ÙŠØ²
            except Exception as e:
                logging.warning(f"âš ï¸ ÙØ´Ù„ Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù‚Ù†Ø§Ø© {channel_id} Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… {user_id}: {e}")
        
        # Ø¥Ø°Ø§ Ù„Ù… ÙŠÙƒÙ† Ù…Ù…ÙŠØ² â†’ Ø²ÙŠØ§Ø¯Ø© Ø§Ù„Ø¹Ø¯Ø§Ø¯
        with sqlite3.connect("quiz_users.db", check_same_thread=False) as conn:
            cursor = conn.cursor()
            cursor.execute("UPDATE users SET quiz_count = quiz_count + 1 WHERE user_id = ?", (user_id,))
            conn.commit()
    
    except Exception as e:
        logging.error(f"ğŸš« Ø®Ø·Ø£ ÙÙŠ increment_count Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… {user_id}: {e}")

def can_play_game_today(user_id: int, game_type: str) -> bool:
    if str(user_id) == str(ADMIN_ID):  # Ù…Ù‚Ø§Ø±Ù†Ø© Ø¢Ù…Ù†Ø© Ù„Ø£Ù† ADMIN_ID Ø£Ø­ÙŠØ§Ù†Ù‹Ø§ ÙŠÙƒÙˆÙ† str
        return True

    today = str(date.today())
    cursor.execute(
        "SELECT 1 FROM game_attempts WHERE user_id = ? AND game_type = ? AND date = ?",
        (user_id, game_type, today)
    )
    return cursor.fetchone() is None

def record_game_attempt(user_id: int, game_type: str):
    if str(user_id) == str(ADMIN_ID):
        return  # Ù„Ø§ ØªØ³Ø¬Ù„ Ù„Ù„Ø£Ø¯Ù…Ù†

    today = str(date.today())
    cursor.execute(
        "INSERT OR REPLACE INTO game_attempts(user_id, game_type, date) VALUES (?, ?, ?)",
        (user_id, game_type, today)
    )
    conn.commit()
from collections import defaultdict

# ØªØ®Ø²ÙŠÙ† Ù…Ø¤Ù‚Øª ÙÙŠ Ø§Ù„Ø°Ø§ÙƒØ±Ø©
game_states = defaultdict(dict)  # {user_id: {game_type: count}}

def get_question_count(user_id, game_type):
    return game_states.get(user_id, {}).get(game_type, 0)

def increment_question_count(user_id, game_type):
    game_states[user_id][game_type] = game_states.get(user_id, {}).get(game_type, 0) + 1
    
# -------------------------------------------------------------------
#                 Quiz Generation & Formatting
# -------------------------------------------------------------------

def extract_json_from_string(text: str) -> str:
    """
    Extracts a JSON string from a text that might contain markdown code blocks or other text.
    """
    # Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† Ø¨Ù„ÙˆÙƒ JSON Ø¯Ø§Ø®Ù„ ```json ... ```
    match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
    if match:
        return match.group(1).strip()

    # Ø¥Ø°Ø§ Ù„Ù… ÙŠØ¬Ø¯ Ø¨Ù„ÙˆÙƒØŒ Ø§Ø¨Ø­Ø« Ø¹Ù† Ø£ÙˆÙ„ '{' Ø£Ùˆ '[' ÙˆØ¢Ø®Ø± '}' Ø£Ùˆ ']'
    start = -1
    end = -1
    
    # Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† Ø¨Ø¯Ø§ÙŠØ© Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø£Ùˆ Ø§Ù„ÙƒØ§Ø¦Ù†
    first_brace = text.find('{')
    first_bracket = text.find('[')
    
    if first_brace == -1:
        start = first_bracket
    elif first_bracket == -1:
        start = first_brace
    else:
        start = min(first_brace, first_bracket)

    # Ø¥Ø°Ø§ Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ø¨Ø¯Ø§ÙŠØ©ØŒ Ø£Ø±Ø¬Ø¹ Ø§Ù„Ù†Øµ Ø§Ù„Ø£ØµÙ„ÙŠ
    if start == -1:
        return text

    # Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† Ù†Ù‡Ø§ÙŠØ© Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø£Ùˆ Ø§Ù„ÙƒØ§Ø¦Ù†
    last_brace = text.rfind('}')
    last_bracket = text.rfind(']')
    end = max(last_brace, last_bracket)

    # Ø¥Ø°Ø§ ØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ø¨Ø¯Ø§ÙŠØ© ÙˆÙ†Ù‡Ø§ÙŠØ©ØŒ Ø£Ø±Ø¬Ø¹ Ù…Ø§ Ø¨ÙŠÙ†Ù‡Ù…Ø§
    if end > start:
        return text[start:end+1].strip()
        
    # ÙƒØ®ÙŠØ§Ø± Ø£Ø®ÙŠØ±ØŒ Ø£Ø±Ø¬Ø¹ Ø§Ù„Ù†Øµ ÙƒÙ…Ø§ Ù‡Ùˆ
    return text
    
def generate_quizzes_from_text(content: str, major: str, user_id: int, num_quizzes: int = 10):
    prompt = (
        f"You are a strict AI quiz generator. Your only task is to generate a JSON array of {num_quizzes} quiz questions "
        f"that are based **strictly and only** on the information explicitly stated in the following content.\n\n"
        "â—ï¸Important Rules:\n"
        "- DO NOT invent, infer, or assume any information not clearly mentioned in the text.\n"
        "- If a concept is not explained or mentioned clearly in the content, DO NOT create a question about it.\n"
        "- Stay fully inside the boundaries of the content.\n"
        "- Every question must test **recall** or **recognition** from the provided text only, not general knowledge.\n"
        "- Questions must be varied: some fill-in-the-blank, some multiple-choice.\n"
        "- Include at most one True/False question.\n"
        "- All questions and answers must be in the same language as the content.\n"
        "- if the content language is arabic give the questions and answers in arabic.\n\n"
        "Each question must be an object with:\n"
        "- 'question': the question string\n"
        "- 'options': a list of exactly 4 answer options\n"
        "- 'correct_index': the index (0-3) of the correct answer in the options list\n"
        "- 'explanation': short sentence to explain **why this is the correct answer**, max 2 lines\n\n"
        "âš ï¸ Format Instructions:\n"
        "- ONLY return a raw JSON array. No markdown, no explanation, no formatting.\n"
        "- Do not include any introductory or closing text.\n"
        "- Ensure the JSON is valid and parsable.\n\n"
        f"Content:\n{content}"
    )

    # ØªØ­Ø¯ÙŠØ¯ Ø§Ù„Ø¯Ø§Ù„Ø© Ø¨Ù†Ø§Ø¡Ù‹ Ø¹Ù„Ù‰ ØµÙ„Ø§Ø­ÙŠØ© Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…
    if user_id == ADMIN_ID or can_generate(user_id):  # <-- Ø§Ù„ØªØ­Ù‚Ù‚ Ù‡Ù†Ø§
        raw_response = generate_smart_response(prompt)
    else:
        raw_response = generate_gemini_response(prompt)
    
    # --- Ø§Ù„ØªØ¹Ø¯ÙŠÙ„ ÙŠØ¨Ø¯Ø£ Ù‡Ù†Ø§ ---
    # 1. ØªÙ†Ø¸ÙŠÙ Ø§Ù„Ø§Ø³ØªØ¬Ø§Ø¨Ø© Ù„Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ù€ JSON
    clean_json_str = extract_json_from_string(raw_response)
    
    # 2. Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù…Ø§ Ø¥Ø°Ø§ ÙƒØ§Ù†Øª Ø§Ù„Ø§Ø³ØªØ¬Ø§Ø¨Ø© ÙØ§Ø±ØºØ© Ø¨Ø¹Ø¯ Ø§Ù„ØªÙ†Ø¸ÙŠÙ
    if not clean_json_str:
        logging.error(f"âŒ JSON extraction failed. Raw output was:\n{raw_response}")
        return [] # Ø£Ø±Ø¬Ø¹ Ù‚Ø§Ø¦Ù…Ø© ÙØ§Ø±ØºØ© Ø¨Ø¯Ù„Ø§Ù‹ Ù…Ù† Ø±Ø³Ø§Ù„Ø© Ø®Ø·Ø£

    try:
        # 3. Ù…Ø­Ø§ÙˆÙ„Ø© ØªØ­Ù„ÙŠÙ„ Ø§Ù„Ø³Ù„Ø³Ù„Ø© Ø§Ù„Ù†Ø¸ÙŠÙØ©
        quizzes_json = json.loads(clean_json_str)
        quizzes = []

        for item in quizzes_json:
            q = item.get("question", "").strip()
            opts = item.get("options", [])
            corr = item.get("correct_index", -1)
            expl = item.get("explanation", "").strip()

            if isinstance(q, str) and q and isinstance(opts, list) and len(opts) == 4 and isinstance(corr, int) and 0 <= corr < 4:
                quizzes.append((q, [str(opt).strip() for opt in opts], corr, expl))
            else:
                logging.warning(f"âŒ Skipping invalid question structure: {item}")

        return quizzes

    except json.JSONDecodeError as e:
        logging.error(f"âŒ JSON parsing failed: {e}\nCleaned string was:\n{clean_json_str}\nRaw output was:\n{raw_response}")
        return [] # Ø£Ø±Ø¬Ø¹ Ù‚Ø§Ø¦Ù…Ø© ÙØ§Ø±ØºØ© Ø¹Ù†Ø¯ Ø§Ù„ÙØ´Ù„
    # --- Ø§Ù„ØªØ¹Ø¯ÙŠÙ„ ÙŠÙ†ØªÙ‡ÙŠ Ù‡Ù†Ø§ ---


def save_quiz_to_db(quiz_data, user_id):
    conn = sqlite3.connect('medical_quizzes.db')
    cursor = conn.cursor()
    
    cursor.execute('''
    INSERT INTO quizzes (title, questions_json, user_id)
    VALUES (?, ?, ?)
    ''', (quiz_data['title'], json.dumps(quiz_data['questions']), user_id))
    
    quiz_id = cursor.lastrowid
    conn.commit()
    conn.close()
    return quiz_id



import json
import logging
import requests

WIKIMEDIA_API = "https://commons.wikimedia.org/w/api.php"

def search_image_on_wikimedia(query: str) -> str:
    """
    Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† ØµÙˆØ±Ø© Ù…Ù† Wikimedia Ø¨Ù†Ø§Ø¡Ù‹ Ø¹Ù„Ù‰ Ø§Ù„ÙˆØµÙ
    """
    params = {
        "action": "query",
        "format": "json",
        "prop": "imageinfo",
        "generator": "search",
        "gsrsearch": query + " filetype:bitmap OR filetype:jpeg OR filetype:png",
        "gsrlimit": 1,
        "iiprop": "url",
    }
    try:
        r = requests.get(WIKIMEDIA_API, params=params)
        r.raise_for_status()
        data = r.json()

        if "query" in data and "pages" in data["query"]:
            page = next(iter(data["query"]["pages"].values()))
            if "imageinfo" in page:
                return page["imageinfo"][0]["url"]
        return ""
    except Exception as e:
        logging.error(f"âŒ ÙØ´Ù„ Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† Ø§Ù„ØµÙˆØ±Ø©: {e}")
        return ""

def generate_Medical_quizzes(content: str, major: str, user_id: int, num_quizzes: int = 10):
    # (Ø§Ù„Ø¨Ø±ÙˆÙ…Ø¨Øª Ø§Ù„Ù…Ø­Ø³Ù† Ù…Ù† Ø§Ù„Ø®Ø·ÙˆØ© 2 ÙŠØ¬Ø¨ ÙˆØ¶Ø¹Ù‡ Ù‡Ù†Ø§)
    prompt = (
        f"You are a medical education expert. Your task is to create a JSON-formatted quiz for {major} "
        "medical students (Year 3-4) based ONLY on the provided reference text.\n\n"
        "## EXTREMELY STRICT RULES:\n"
        f"1. You MUST STRICTLY generate {num_quizzes} questions. No more, no less.\n"
        "2. 70% multiple-choice (basic sciences), 30% problem-solving (clinical cases).\n"
        "3. Use only information from the reference text.\n"
        "4. Clinical questions must have realistic short scenarios (2-3 sentences).\n"
        "5. For any question that would benefit from an image, add an 'image_prompt' field.\n"
        "6. Language: English.\n"
        "7. CRITICAL: The 'questions' array MUST NOT be empty. If you cannot generate questions from the text, return an empty JSON object {} and nothing else.\n\n"
        "## JSON OUTPUT STRUCTURE:\n"
        "{\n"
        "  \"title\": \"Medical Quiz in [major]\",\n"
        "  \"questions\": [\n"
        "    {\n"
        "      \"id\": 1,\n"
        "      \"type\": \"multiple_choice\",\n"
        "      \"question\": \"...\",\n"
        "      \"options\": [\"...\", \"...\", \"...\", \"...\"],\n"
        "      \"correct_index\": 0,\n"
        "      \"image_prompt\": \"Description of the image (optional)\"\n"
        "    }\n"
        "  ]\n"
        "}\n\n"
        f"## Reference text:\n{content}"
    )

    # --- Ø¢Ù„ÙŠØ© Ø¥Ø¹Ø§Ø¯Ø© Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© (Ø§Ù„Ø®Ø·ÙˆØ© 3) ---
    for attempt in range(3):  # Ø³ÙŠØ­Ø§ÙˆÙ„ Ø­ØªÙ‰ 3 Ù…Ø±Ø§Øª
        logging.info(f"Attempt {attempt + 1} to generate medical quiz...")
        
        # Ø§Ø®ØªÙŠØ§Ø± Ø§Ù„Ù†Ù…ÙˆØ°Ø¬
        if user_id == ADMIN_ID or can_generate(user_id):
            raw_response = generate_smart_response(prompt)
        else:
            raw_response = generate_gemini_response(prompt)
        
        clean_json_str = extract_json_from_string(raw_response)
        if not clean_json_str or clean_json_str == "{}":
            logging.warning(f"Attempt {attempt + 1} failed: AI returned empty response.")
            continue # Ø§Ù†ØªÙ‚Ù„ Ù„Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ø§Ù„ØªØ§Ù„ÙŠØ©

        try:
            quiz_data = json.loads(clean_json_str)

            # --- Ø§Ù„ØªØ­Ù‚Ù‚ Ø§Ù„Ù…Ø´Ø¯Ø¯ (Ø§Ù„Ø®Ø·ÙˆØ© 1) ---
            if "title" not in quiz_data or not quiz_data.get("questions"):
                raise ValueError("Invalid JSON: missing title or questions array is empty.")

            # --- Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù†Ø§Ø¬Ø­Ø© ---
            for i, q in enumerate(quiz_data["questions"]):
                if "image_prompt" in q and q["image_prompt"].strip():
                    image_url = search_image_on_wikimedia(q["image_prompt"])
                    q["image_url"] = image_url if image_url else ""
                else:
                    q["image_url"] = ""
                if "id" not in q:
                    q["id"] = i + 1

            if major not in quiz_data["title"]:
                quiz_data["title"] = f"Medical Quiz in {major}"

            quiz_id = save_quiz_to_db(quiz_data, user_id)
            quiz_data["db_id"] = quiz_id

            logging.info(f"Successfully generated quiz on attempt {attempt + 1}.")
            return quiz_data # Ø¥Ø±Ø¬Ø§Ø¹ Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ø§Ù„Ù†Ø§Ø¬Ø­Ø© ÙˆØ¥Ù†Ù‡Ø§Ø¡ Ø§Ù„Ø¯Ø§Ù„Ø©

        except (json.JSONDecodeError, ValueError) as e:
            # --- ØªØ³Ø¬ÙŠÙ„ Ø§Ù„Ø®Ø·Ø£ Ù„Ù„ØªØ´Ø®ÙŠØµ (Ø§Ù„Ø®Ø·ÙˆØ© 4) ---
            logging.error(f"Attempt {attempt + 1} failed during JSON processing: {e}")
            logging.error(f"RAW AI RESPONSE WAS:\n{raw_response}\n")
            continue # Ø§Ù†ØªÙ‚Ù„ Ù„Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ø§Ù„ØªØ§Ù„ÙŠØ©

    # Ø¥Ø°Ø§ ÙØ´Ù„Øª ÙƒÙ„ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø§Øª
    logging.error("All 3 attempts to generate medical quiz failed.")
    return None
    


def generate_anki_cards_from_text(content: str, major: str = "General", user_id: int = 0, num_cards: int = 15) -> tuple:
    for attempt in range(3):  # ØªØ¬Ø±Ø¨Ø© Ø­ØªÙ‰ 3 Ù…Ø±Ø§Øª
        prompt = f"""
You are an AI assistant specialized in creating study flashcards.

ğŸ¯ Task:
Extract the most important {num_cards} points from the following content, and convert each into an **Anki-style flashcard**.

ğŸ”¹ Rules:
- Each flashcard must include:
  - "front": a short question or hint.
  - "back": the detailed answer or explanation.
  - "tag": (optional) topic label like Grammar, Biology, Logic, etc.
- The front must be phrased to encourage recall (e.g. "What is...", "Define...", "How does...").
- Don't use Markdown, just clean plain text.
- Keep the cards diverse and helpful.
- Output must be a valid JSON **object** with two keys: "title" and "cards".

ğŸš« Important:
- Do NOT generate multiple choice or true/false questions.
- Only generate flashcards suitable for Anki with a front and a back.
- The flashcards must be written in the same language as the input content. If the content is in Arabic, answer in Arabic. If English, answer in English.

ğŸ“˜ Content to process (field: {major}):
{content}

âœ… Example output format:
{{
  "title": "Basics of Organic Chemistry",
  "cards": [
    {{
      "front": "What is the function of mitochondria?",
      "back": "It is the powerhouse of the cell.",
      "tag": "Biology"
    }},
    {{
      "front": "Ù…Ø§ Ù‡ÙŠ Ø§Ù„Ø§Ø³ØªØ¹Ø§Ø±Ø©ØŸ",
      "back": "Ø§Ù„Ø§Ø³ØªØ¹Ø§Ø±Ø© Ù‡ÙŠ Ø§Ø³ØªØ®Ø¯Ø§Ù… Ø§Ù„ÙƒÙ„Ù…Ø© ÙÙŠ ØºÙŠØ± Ù…Ø¹Ù†Ø§Ù‡Ø§ Ø§Ù„Ø­Ù‚ÙŠÙ‚ÙŠ Ù„Ø¹Ù„Ø§Ù‚Ø© Ù…Ø¹ Ù‚Ø±ÙŠÙ†Ø© Ù…Ø§Ù†Ø¹Ø©.",
      "tag": "Literature"
    }}
  ]
}}
"""
        if user_id == ADMIN_ID or can_generate(user_id):  # <-- Ø§Ù„ØªØ­Ù‚Ù‚ Ù‡Ù†Ø§
            raw_output = generate_smart_response(prompt)
        else:
            raw_output = generate_gemini_response(prompt)
            
        clean_json = extract_json_from_string(raw_output)

        try:
            data = json.loads(clean_json)
            title = data.get("title", "Ø¨Ø·Ø§Ù‚Ø§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ©")
            card_list = data.get("cards", [])

            cards = []
            for item in card_list:
                front = item.get("front") or item.get("question")
                back = item.get("back") or item.get("answer")

                if isinstance(front, str) and isinstance(back, str) and front.strip() and back.strip():
                    cards.append({"front": front.strip(), "back": back.strip()})
                else:
                    logging.warning(f"âŒ Skipping invalid card: {item}")

            if len(cards) >= 5:
                return cards, title

        except json.JSONDecodeError as e:
            logging.error(f"âŒ Failed to parse Anki cards: {e}\nClean JSON:\n{clean_json}\nRaw:\n{raw_output}")

    return [], "Ø¨Ø·Ø§Ù‚Ø§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ©"   



import json
import logging
import requests
from pptx import Presentation

WIKIMEDIA_API = "https://commons.wikimedia.org/w/api.php"
UNSPLASH_API = "https://api.unsplash.com/search/photos"
PEXELS_API = "https://api.pexels.com/v1/search"

# --- Ø¨Ø­Ø« Ø§Ù„ØµÙˆØ± Ù…Ù† Ø£ÙƒØ«Ø± Ù…Ù† Ù…ØµØ¯Ø± ---
def search_image(query: str) -> str:
    """
    ÙŠØ¨Ø­Ø« Ø¹Ù† ØµÙˆØ±Ø© Ù…Ù†Ø§Ø³Ø¨Ø© Ù…Ù† Wikimedia Ø«Ù… Unsplash Ø«Ù… Pexels.
    ÙŠØ¹ÙŠØ¯ Ø£ÙˆÙ„ Ø±Ø§Ø¨Ø· ØµÙˆØ±Ø© ØµØ§Ù„Ø­Ø©.
    """
    # 1. Wikimedia
    params = {
        "action": "query",
        "format": "json",
        "prop": "imageinfo",
        "generator": "search",
        "gsrsearch": query + " filetype:bitmap OR filetype:jpeg OR filetype:png",
        "gsrlimit": 1,
        "iiprop": "url",
    }
    try:
        r = requests.get(WIKIMEDIA_API, params=params, timeout=10)
        r.raise_for_status()
        data = r.json()
        if "query" in data and "pages" in data["query"]:
            page = next(iter(data["query"]["pages"].values()))
            if "imageinfo" in page:
                return page["imageinfo"][0]["url"]
    except Exception as e:
        logging.warning(f"âš ï¸ ÙØ´Ù„ Wikimedia: {e}")

    # 2. Unsplash (ÙŠØªØ·Ù„Ø¨ API key ÙÙŠ Ù…ØªØºÙŠØ± Ø§Ù„Ø¨ÙŠØ¦Ø© UNSPLASH_KEY)
    unsplash_key = os.getenv("UNSPLASH_KEY")
    if unsplash_key:
        try:
            r = requests.get(
                UNSPLASH_API,
                params={"query": query, "per_page": 1},
                headers={"Authorization": f"Client-ID {unsplash_key}"},
                timeout=10,
            )
            r.raise_for_status()
            data = r.json()
            if data.get("results"):
                return data["results"][0]["urls"]["regular"]
        except Exception as e:
            logging.warning(f"âš ï¸ ÙØ´Ù„ Unsplash: {e}")

    # 3. Pexels (ÙŠØªØ·Ù„Ø¨ API key ÙÙŠ Ù…ØªØºÙŠØ± Ø§Ù„Ø¨ÙŠØ¦Ø© PEXELS_KEY)
    pexels_key = os.getenv("PEXELS_KEY")
    if pexels_key:
        try:
            r = requests.get(
                PEXELS_API,
                params={"query": query, "per_page": 1},
                headers={"Authorization": pexels_key},
                timeout=10,
            )
            r.raise_for_status()
            data = r.json()
            if data.get("photos"):
                return data["photos"][0]["src"]["medium"]
        except Exception as e:
            logging.warning(f"âš ï¸ ÙØ´Ù„ Pexels: {e}")

    return ""


def generate_special_anki_cards_from_text(content: str, major: str = "General", user_id: int = 0, num_cards: int = 15) -> tuple:
    for attempt in range(3):  # ØªØ¬Ø±Ø¨Ø© Ø­ØªÙ‰ 3 Ù…Ø±Ø§Øª
        prompt = f"""
You are an AI assistant specialized in creating study flashcards.

ğŸ¯ Task:
Extract the most important {num_cards} points from the following content, and convert each into an **Anki-style flashcard**.

ğŸ”¹ Rules:

Each flashcard must include:

"front": a short question or hint.

"back": the detailed answer or explanation.

"tag": (optional) topic label like Grammar, Biology, Logic, etc.

"image_hint": (optional) a short description of an image that would help illustrate the card (only if relevant).

The front must be phrased to encourage recall (e.g. "What is...", "Define...", "How does...").

Don't use Markdown, just clean plain text.

Keep the cards diverse and helpful.

Output must be a valid JSON object with two keys: "title" and "cards".

ğŸš« Important:

Do NOT generate multiple choice or true/false questions.

Only generate flashcards suitable for Anki with a front and a back.

The flashcards must be written in the same language as the input content. If the content is in Arabic, answer in Arabic. If English, answer in English.

ğŸ“˜ Content to process (field: {major}):
{content}

âœ… Example output format:
{{
"title": "Basics of Organic Chemistry",
"cards": [
{{
"front": "What is the function of mitochondria?",
"back": "It is the powerhouse of the cell.",
"tag": "Biology",
"image_hint": "microscopic image of mitochondria"
}},
{{
"front": "Ù…Ø§ Ù‡ÙŠ Ø§Ù„Ø§Ø³ØªØ¹Ø§Ø±Ø©ØŸ",
"back": "Ø§Ù„Ø§Ø³ØªØ¹Ø§Ø±Ø© Ù‡ÙŠ Ø§Ø³ØªØ®Ø¯Ø§Ù… Ø§Ù„ÙƒÙ„Ù…Ø© ÙÙŠ ØºÙŠØ± Ù…Ø¹Ù†Ø§Ù‡Ø§ Ø§Ù„Ø­Ù‚ÙŠÙ‚ÙŠ Ù„Ø¹Ù„Ø§Ù‚Ø© Ù…Ø¹ Ù‚Ø±ÙŠÙ†Ø© Ù…Ø§Ù†Ø¹Ø©.",
"tag": "Literature",
"image_hint": ""
}}
]
}}
"""
        # ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ù…Ø®Ø±Ø¬Ø§Øª
        if user_id == ADMIN_ID or can_generate(user_id):
            raw_output = generate_smart_response(prompt)
        else:
            raw_output = generate_gemini_response(prompt)

        clean_json = extract_json_from_string(raw_output)

        try:
            data = json.loads(clean_json)
            title = data.get("title", "Ø¨Ø·Ø§Ù‚Ø§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ©")
            card_list = data.get("cards", [])

            cards = []
            for item in card_list:
                front = item.get("front") or item.get("question")
                back = item.get("back") or item.get("answer")
                tag = item.get("tag", "")
                image_hint = item.get("image_hint", "").strip()
                image_url = ""

                if isinstance(front, str) and isinstance(back, str) and front.strip() and back.strip():
                    # Ø§Ù„Ø¨Ø­Ø« Ø¹Ù† ØµÙˆØ±Ø© Ø¥Ø°Ø§ ÙƒØ§Ù† ÙÙŠÙ‡ image_hint
                    if image_hint:
                        image_url = search_image(image_hint)

                    cards.append({
                        "front": front.strip(),
                        "back": back.strip(),
                        "tag": tag.strip(),
                        "image_url": image_url  # <-- Ø±Ø§Ø¨Ø· Ø§Ù„ØµÙˆØ±Ø© ÙÙ‚Ø·
                    })
                else:
                    logging.warning(f"âŒ Skipping invalid card: {item}")

            if len(cards) >= 5:
                return cards, title

        except json.JSONDecodeError as e:
            logging.error(f"âŒ Failed to parse Anki cards: {e}\nClean JSON:\n{clean_json}\nRaw:\n{raw_output}")

    return [], "Ø¨Ø·Ø§Ù‚Ø§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ©"


# -------------------------------------------------------------------
#                 games
# -------------------------------------------------------------------
import random

topics = [
    "Ø­ÙŠØ§Ø© Ø§Ù„Ø·Ø§Ù„Ø¨", "ØªØ®Ø·ÙŠØ· Ø§Ù„Ø³ÙØ±", "Ù…Ø´Ø§Ø±ÙŠØ¹ Ø¬Ù…Ø§Ø¹ÙŠØ©", "Ù…Ù‚Ø§Ø¨Ù„Ø§Øª Ø§Ù„Ø¹Ù…Ù„",
    "Ø§Ù„Ø¶ØºØ· Ø§Ù„Ø²Ù…Ù†ÙŠ", "Ù…ÙˆØ§Ù‚Ù Ø¹Ø§Ø·ÙÙŠØ©", "Ø§Ø³ØªØ®Ø¯Ø§Ù… Ø§Ù„ØªÙƒÙ†ÙˆÙ„ÙˆØ¬ÙŠØ§", "Ù‚Ø±Ø§Ø±Ø§Øª Ù…Ø§Ù„ÙŠØ©",
    "ØµØ±Ø§Ø¹Ø§Øª Ø§Ù„ÙØ±ÙŠÙ‚", "ØªØ­Ø¯ÙŠØ¯ Ø§Ù„Ø£Ù‡Ø¯Ø§Ù"
]

def generate_vocabulary_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "vocab")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""  
You are an AI vocabulary quiz creator.  
Generate one vocabulary question for a student majoring in {major}.
- Vocabulary should be relevant to real life or academic use and not an uncommon Vocabulary.
- Show the meaning of an English word in English 
- Provide 4 English words as options  
- Only ONE option should be correct.  
- Don't explain anything. Just give raw JSON.

Example:
{{
  "question": "Question",
  "options": ["Option", "Option", "Option", "Option"],
  "correct_index": 0
}}

Use this seed to diversify the question: {rand}
âŒ Avoid repeating or paraphrasing these questions:
{recent_prompt}
"""
    q = generate_game(prompt)

    # Ø­ÙØ¸ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
    add_recent_question(user_id, "speed", q["question"])
    return q

def generate_speed_challenge(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "speed")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""
You are a quiz bot.

Generate a **fun, fast-answer quiz** for a student in {major}.

Requirements:
- The question must be in English.
- The 4 options must be in English.
- Use fun and fast general knowledge topics (e.g. logic, daily life trivia, or language puzzles). Avoid repeating the same categories.
- Keep it simple and not too academic.
- Return raw JSON only.
- No explanation.
- Use this seed to increase randomness: {rand}
âŒ Avoid repeating or paraphrasing these questions:
{recent_prompt}

Example output:
{{
  "question": "Question?",
  "options": ["Option", "Option", "Option", "Option"],
  "correct_index": 0
}}
"""
    q = generate_game(prompt, translate_question=True)

    # Ø­ÙØ¸ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
    add_recent_question(user_id, "speed", q["question"])
    return q
    

# â˜… Ù„Ø¹Ø¨Ø© Ø§Ù„Ø§Ø®Ø·Ø§Ø¡ Ø§Ù„Ø´Ø§Ø¦Ø¹Ø©
def generate_common_mistakes_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "mistakes")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    prompt = f"""
You are an educational game generator.

Your task:
- Generate a multiple-choice question highlighting a **common mistake** in the field of {major}.
- The question must be in English.
- The **options must be in English**.
- Provide **4 options** only, with one correct.
- Don't explain.
- Return only raw JSON.

âŒ Avoid repeating or paraphrasing these questions:
{recent_prompt}
Use this random seed to diversify the question: {rand}

Example output:
{{
  "question": "Which sentence is grammatically incorrect?",
  "options": ["He go to school every day.", "She plays the piano.", "They are studying now.", "I have finished my homework."],
  "correct_index": 0
}}
"""
    q = generate_game(prompt, translate_question=True)

    # Ø­ÙØ¸ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
    add_recent_question(user_id, "speed", q["question"])
    return q


def generate_inference_game(user_id, major, native_lang="Arabic"):
    rand = random.randint(1000, 9999)
    recent = get_recent_questions(user_id, "inference")
    recent_prompt = "\n".join(f"- {q}" for q in recent)
    
    random_topic = random.choice(topics)
    prompt = f"""
You are an AI-powered life skills test creator.

Generate a **new and unique** question that develops one of the following skills:  
- Critical thinking  
- Emotional intelligence  
- Time management  
- Self-awareness  
- Decision making  
- Problem solving  
- Logic  
- Pattern recognition  
- Mental map understanding  

ğŸ”¹ **Requirements**:  
- Write the **question in Arabic**  
- Write **all options in Arabic**  
- Use a realistic scenario or student-life context related to: **{random_topic}**  
- Provide **exactly 4 options**, with **one correct answer**  
- **Never repeat** past examples or add explanations  
- Make the question **engaging and clever**  
- Incorporate variability using this random number: **{rand}**  
- the options should be as short as possible but understandable
âŒ Avoid repeating or paraphrasing these questions:
{recent_prompt}
ğŸ”¸ Return **JSON-only output** (no additional text).  

Example (Johnsonâ€™s format):  
{{
  "question": "Question",  
  "options": ["Options", "Option", "Option", "Option"],  
  "correct_index": 2  
}}  
"""
    q = generate_game(prompt, translate_question=True)

    # Ø­ÙØ¸ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
    add_recent_question(user_id, "speed", q["question"])
    return q
    
# ----------------------------------
# ------------- inference review -------------------------------------------------------------------


def review_inference_question_with_ai(question_text: str, options: list[str], correct_index: int) -> bool:
    prompt = f"""
You are an AI educational assistant.

A student submitted the following inference question. Review it and decide if it's valid:
- Is the question clear and meaningful?
- Are the 4 options distinct and related to the question?
- Is there **one and only one correct answer**?

Respond only with YES or NO.

Question: {question_text}
Options: {options}
Correct index: {correct_index}
"""
    response = generate_smart_response(prompt).strip().lower()
    return "yes" in response

# Ø¹Ø¯Ù‘Ù„ Ù‡Ø°Ù‡ Ø§Ù„Ø¯Ø§Ù„Ø©
def send_quiz_to_user(chat_id, quiz_data, message_id=None): # message_id ÙŠØµØ¨Ø­ Ø§Ø®ØªÙŠØ§Ø±ÙŠØ§Ù‹
    markup = InlineKeyboardMarkup()
    quiz_url = f"{WEBHOOK_URL}/quiz/{quiz_data['db_id']}"
    btn = InlineKeyboardButton("ğŸš€ ÙØªØ­ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±", url=quiz_url)
    markup.add(btn)
    
    message = f"""
    ğŸ† ØªÙ… Ø¥Ù†Ø´Ø§Ø¡ Ø§Ø®ØªØ¨Ø§Ø±Ùƒ Ø§Ù„Ø·Ø¨ÙŠ Ø¨Ù†Ø¬Ø§Ø­!
    Ø§Ù„Ø¹Ù†ÙˆØ§Ù†: {quiz_data['title']}
    Ø¹Ø¯Ø¯ Ø§Ù„Ø£Ø³Ø¦Ù„Ø©: {len(quiz_data['questions'])}
    """
    # Ø§Ø­Ø°Ù Ø§Ù„Ø´Ø±Ø· ÙˆØ£Ø±Ø³Ù„ Ø±Ø³Ø§Ù„Ø© Ø¬Ø¯ÙŠØ¯Ø© Ø¯Ø§Ø¦Ù…Ù‹Ø§
    bot.send_message(
        chat_id,
        message,
        reply_markup=markup
    )
    
def process_pending_inference_questions():
    cursor.execute("SELECT id, question, options, correct_index FROM inference_questions WHERE approved = 0")
    pending = cursor.fetchall()

    for row in pending:
        qid, qtext, options_json, correct_index = row
        try:
            options = json.loads(options_json)
        except:
            continue  # ØªØ¬Ø§Ù‡Ù„ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ø°Ø§Øª Ø§Ù„ØªÙ†Ø³ÙŠÙ‚ Ø§Ù„Ø®Ø§Ø·Ø¦

        if review_inference_question_with_ai(qtext, options, correct_index):
            cursor.execute("UPDATE inference_questions SET approved = 1 WHERE id = ?", (qid,))
        else:
            cursor.execute("DELETE FROM inference_questions WHERE id = ?", (qid,))

    conn.commit()


import sqlite3
import time
import json
import uuid
import threading
from datetime import datetime
from telebot import types
import uuid

from datetime import datetime

def log_quiz_share(quiz_code, shared_by_user_id, shared_by_name):
    conn = sqlite3.connect("quiz_users.db")
    c = conn.cursor()

    shared_at = datetime.now().isoformat()

    c.execute("""
        INSERT INTO quiz_shares (quiz_code, shared_by_user_id, shared_by_name, shared_at)
        VALUES (?, ?, ?, ?)
    """, (quiz_code, shared_by_user_id, shared_by_name, shared_at))

    conn.commit()
    conn.close()



# ØªÙˆÙ„ÙŠØ¯ ÙƒÙˆØ¯ ÙØ±ÙŠØ¯ Ù…Ø¹ Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„ØªÙƒØ±Ø§Ø±
def generate_unique_quiz_code():
    while True:
        code = f"QC_{uuid.uuid4().hex[:6]}"
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()
        c.execute("SELECT 1 FROM user_quizzes WHERE quiz_code = ?", (code,))
        if not c.fetchone():
            conn.close()
            return code
        conn.close()

# ØªØ®Ø²ÙŠÙ† Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± (Ù†Ø³Ø®Ø© Ù…Ø­Ø³Ù†Ø©)
def store_quiz(user_id, quizzes, bot):
    try:
    
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()

        quiz_code = generate_unique_quiz_code()

        # Ø¬Ù„Ø¨ Ø§Ø³Ù… Ø§Ù„Ù…Ø§Ù„Ùƒ
        try:
            owner_chat = bot.get_chat(user_id)
            owner_name = owner_chat.first_name or owner_chat.username or f"user_{user_id}"
        except Exception:
            owner_name = "ØµØ¯ÙŠÙ‚Ùƒ"

        created_at = datetime.now().isoformat()

        c.execute("""
            INSERT INTO user_quizzes (user_id, quiz_data, quiz_code, created_at, is_active, owner_name)
            VALUES (?, ?, ?, ?, 1, ?)
        """, (user_id, json.dumps(quizzes), quiz_code, created_at, owner_name))

        conn.commit()
        conn.close()

        return quiz_code
    except Exception as e:
        print(f"Error storing quiz: {e}")
        return None



def start_quiz(chat_id, quiz_code, bot):
    if quiz_code == "sample":
        # Ø§Ø³ØªØ¹Ù„Ø§Ù… Ù…Ù† Ø¬Ø¯ÙˆÙ„ Ø§Ù„Ø¹ÙŠÙ†Ø©
        cursor.execute("SELECT quiz_data FROM sample_quizzes WHERE quiz_code = ?", (quiz_code,))
        row = cursor.fetchone()
        if not row:
            return False
        quiz_data_json = row[0]
    else:
        # Ø§Ø³ØªØ¹Ù„Ø§Ù… Ù…Ù† Ø¬Ø¯ÙˆÙ„ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…ÙŠÙ†
        cursor.execute("SELECT quiz_data FROM user_quizzes WHERE quiz_code = ?", (quiz_code,))
        row = cursor.fetchone()
        if not row:
            return False
        quiz_data_json = row[0]

    # ØªØ­ÙˆÙŠÙ„ JSON Ø¥Ù„Ù‰ Ø¨ÙŠØ§Ù†Ø§ØªØŒ Ø«Ù… Ù…ØªØ§Ø¨Ø¹Ø© ØªØ´ØºÙŠÙ„ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
    quiz_data = json.loads(quiz_data_json)
    # ØªØ§Ø¨Ø¹ Ø¹Ø±Ø¶ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© ÙˆØ§Ù„ØªØ¹Ø§Ù…Ù„ Ù…Ø¹ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ù‡Ù†Ø§...
    # ...

    return True
    
    
    
logger = logging.getLogger(__name__)


# Ù†Ø¸Ø§Ù… Ø¥Ø¯Ø§Ø±Ø© Ø§Ù„Ø­Ø§Ù„Ø© Ø§Ù„Ù…Ø­Ø³Ù†
class QuizManager:

    def __init__(self):
        self.active_quizzes = {}

    def start_quiz(self, chat_id, quiz_code, bot2, message_id=None):
        try:
            # ----- Ø¬Ù„Ø¨ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ø­Ø³Ø¨ Ù†ÙˆØ¹ Ø§Ù„ÙƒÙˆØ¯ -----
            if quiz_code == "sample":
                # Ø§ÙØªØ­ Ø§ØªØµØ§Ù„Ù‹Ø§ Ø¨Ù‚Ø§Ø¹Ø¯Ø© quiz_users.db Ø­ÙŠØ« Ù…Ø®Ø²Ù† sample_quizzes
                conn = sqlite3.connect('quiz_users.db')
                cur = conn.cursor()
                cur.execute("SELECT quiz_data FROM sample_quizzes WHERE quiz_code = ?", (quiz_code,))
                row = cur.fetchone()
                conn.close()

                if not row:
                    logger.info("sample quiz not found in sample_quizzes")
                    return False

                # row[0] ÙŠÙØªØ±Ø¶ Ø£Ù†Ù‡ JSON Ù†ØµÙŠ
                quizzes_raw = json.loads(row[0]) if isinstance(row[0], str) else row[0]
                # sample Ø¹Ø§Ø¯Ø© Ù„Ø§ ÙŠÙ…Ù„Ùƒ owner
                user_id = None
                score = 0
                total = len(quizzes_raw)

                shared_by_name = "TestGenie"   # Ø§Ø³Ù… Ø§ÙØªØ±Ø§Ø¶ÙŠ Ù„Ù„Ø¹Ø±Ø¶
                owner_name = "Ø§Ù„Ù…Ø±Ø³Ù„"

            else:
                # Ø§Ø³ØªØ¹Ù„Ø§Ù… Ù…Ù† user_quizzes
                conn = sqlite3.connect("quiz_users.db")
                c = conn.cursor()
                c.execute("""
                    SELECT user_id, score, total, quiz_data, owner_name
                    FROM user_quizzes
                    WHERE quiz_code = ?
                """, (quiz_code,))
                row = c.fetchone()
                conn.close()

                if not row:
                    logger.info("user quiz not found for code %s", quiz_code)
                    return False

                # ØªÙÙƒÙŠÙƒ row Ø¨Ø£Ù…Ø§Ù†
                user_id = row[0]
                score = row[1] or 0
                total = row[2] or 0
                quiz_data_field = row[3]
                owner_name = row[4] or "Ø§Ù„Ù…Ø§Ù„Ùƒ"

                # quiz_data_field ÙŠÙ…ÙƒÙ† Ø£Ù† ÙŠÙƒÙˆÙ† Ù†Øµ JSON Ø£Ùˆ Ù„Ù‚Ø§Ø¦Ù…Ø© Ø¨Ø§Ù„ÙØ¹Ù„
                quizzes_raw = json.loads(quiz_data_field) if isinstance(quiz_data_field, str) else quiz_data_field
                shared_by_name = owner_name

            # ----- ØªØ­ÙˆÙŠÙ„ Ø£ÙŠ Ø¨Ù†ÙŠØ© (dict Ø£Ùˆ list) Ø¥Ù„Ù‰ Ù‚Ø§Ø¦Ù…Ø© Ù…ÙˆØ­Ø¯Ø© Ù…Ù† Ø§Ù„Ù‚ÙˆØ§Ù…ÙŠØ³ -----
            formatted_quizzes = []
            for q in quizzes_raw:
                if isinstance(q, dict):
                # Ø¨Ù†ÙŠØ© Ø§Ù„Ù‚Ø§Ù…ÙˆØ³ Ø§Ù„Ù…ØªÙˆÙ‚Ø¹Ø©
                    question = q.get('question')
                    options = q.get('options')
                    correct_idx = q.get('correct_index')
                    explanation = q.get('explanation', '')
                elif isinstance(q, (list, tuple)) and len(q) >= 4:
                    # Ø¨Ù†ÙŠØ© Ù‚Ø§Ø¦Ù…Ø© Ù‚Ø¯ ØªÙƒÙˆÙ† [question, options, correct_idx, explanation]
                    question, options, correct_idx, explanation = q[0], q[1], q[2], q[3]
                else:
                    logger.warning("Unknown question format: %r", q)
                    continue

                # ØªØ­Ù‚Ù‚ Ù…Ù† ØµØ­Ø© Ø§Ù„Ø®ÙŠØ§Ø±Ø§Øª
                if not (isinstance(options, list) and len(options) == 4):
                    logger.warning("Skipping question with invalid options: %r", question)
                    continue
                if not isinstance(correct_idx, int) or not (0 <= correct_idx < 4):
                    logger.warning("Skipping question with invalid correct_index: %r", question)
                    continue

                formatted_quizzes.append({
                    'question': question,
                    'options': options,
                    'answer': options[correct_idx],
                    'explanation': explanation
                })

            if not formatted_quizzes:
                logger.info("No valid questions found for quiz %s", quiz_code)
                return False

            # ----- Ø§Ø­ÙØ¸ Ø§Ù„Ø­Ø§Ù„Ø© Ø§Ù„Ù†Ø´Ø·Ø© Ù„Ù„Ø§Ù…ØªØ­Ø§Ù† -----
            self.active_quizzes[chat_id] = {
                'quizzes': formatted_quizzes,
                'current_index': 0,
                'score': 0,
                'quiz_code': quiz_code,
                'start_time': datetime.now(),
                'owner_id': user_id,
                'owner_score': score,
                'owner_total': total
            }

            # ----- Ø±Ø³Ø§Ù„Ø© ØªÙ‚Ø¯ÙŠÙ…ÙŠØ© -----
            estimated_time = round((total or len(formatted_quizzes)) * 0.5)  # Ù†ØµÙ Ø¯Ù‚ÙŠÙ‚Ø© Ù„ÙƒÙ„ Ø³Ø¤Ø§Ù„ ØªÙ‚Ø±ÙŠØ¨Ù‹Ø§

            if user_id is not None and chat_id == user_id:
                msg = f"ØªÙ… ØªÙˆÙ„ÙŠØ¯ {len(formatted_quizzes)} Ø³Ø¤Ø§Ù„Ù‹Ø§ØŒ Ø§Ø³ØªØ¹Ø¯ Ù„Ù„Ø§Ø®ØªØ¨Ø§Ø±."
            else:
                percent = round((score / total) * 100) if total else 0
                msg = (
                    f"ğŸ¯ <b>{shared_by_name} Ø£Ø±Ø³Ù„ Ù„Ùƒ ØªØ­Ø¯ÙŠÙ‹Ø§!</b> ğŸ¤\n\n"
                    f"ğŸ“‹ <b>Ø¹Ø¯Ø¯ Ø§Ù„Ø£Ø³Ø¦Ù„Ø©:</b> {total or len(formatted_quizzes)}\n"
                    f"ğŸ† <b>Ù†ØªÙŠØ¬Ø© {owner_name}:</b> {score}/{total if total else len(formatted_quizzes)} â€” ({percent}%)\n"
                    f"â³ <b>Ø§Ù„ÙˆÙ‚Øª Ø§Ù„Ù…Ù‚Ø¯Ø±:</b> Ø­ÙˆØ§Ù„ÙŠ {estimated_time} Ø¯Ù‚ÙŠÙ‚Ø©\n\n"
                    f"ğŸ”¥ <b>Ù‡Ù„ ØªØ³ØªØ·ÙŠØ¹ Ø§Ù„ØªÙÙˆÙ‚ Ø¹Ù„ÙŠÙ‡ØŸ</b>"
                )

            # Ø§Ø±Ø³Ù„ Ø£Ùˆ Ø¹Ø¯Ù‘Ù„ Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªØ­Ù…ÙŠÙ„
            if message_id:
                try:
                    bot2.edit_message_text(chat_id=chat_id, message_id=message_id, text=msg, parse_mode="HTML")
                except Exception as e:
                    logger.exception("Failed to edit message: %s", e)
            else:
                bot2.send_message(chat_id, msg, parse_mode="HTML")
                time.sleep(1)

            # Ø£Ø±Ø³Ù„ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø£ÙˆÙ„
            self.send_question(chat_id, bot2)
            return True

        except Exception as e:
            logger.exception("start_quiz failed: %s", e)
            return False
    
    
    def get_quiz_info(self, quiz_code):
        """Ø§Ù„Ø­ØµÙˆÙ„ Ø¹Ù„Ù‰ Ù…Ø¹Ù„ÙˆÙ…Ø§Øª Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù…Ù† Ù‚Ø§Ø¹Ø¯Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª"""
        conn = sqlite3.connect("quiz_users.db")
        c = conn.cursor()
        c.execute("SELECT user_id, quiz_data FROM user_quizzes WHERE quiz_code = ?", (quiz_code,))
        result = c.fetchone()
        conn.close()
        
        if not result:
            return None
            
        return {
            'user_id': result[0],
            'quizzes': json.loads(result[1])
        }
    
    def can_access_quiz(self, user_id, quiz_code):
        """Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† ØµÙ„Ø§Ø­ÙŠØ© Ø§Ù„ÙˆØµÙˆÙ„ Ù„Ù„Ø§Ø®ØªØ¨Ø§Ø±"""
        info = self.get_quiz_info(quiz_code)
        if not info:
            return False
            
            
    def send_question(self, chat_id, bot2):
        state = self.active_quizzes.get(chat_id)
        if not state:
            return
            
        quiz = state['quizzes'][state['current_index']]
        
        try:
            poll = bot2.send_poll(
                chat_id=chat_id,
                question=quiz['question'],
                options=quiz['options'],
                type='quiz',
                correct_option_id=quiz['options'].index(quiz['answer']),
                explanation=quiz['explanation'],
                is_anonymous=False,
                open_period=30
            )
            
            state['last_poll_id'] = poll.message_id
        except Exception as e:
            print(f"Error sending poll: {e}")
            self.handle_quiz_end(chat_id, bot2, error=True)


    def handle_answer(self, poll_answer, bot2):
        chat_id = poll_answer.user.id
        state = self.active_quizzes.get(chat_id)
        if not state:
            return
            
        current_quiz = state['quizzes'][state['current_index']]
        is_correct = poll_answer.option_ids[0] == current_quiz['options'].index(current_quiz['answer'])
        
        if is_correct:
            state['score'] += 1
            feedback = "âœ… Ø¥Ø¬Ø§Ø¨Ø© ØµØ­ÙŠØ­Ø©!"
        else:
            feedback = f"âŒ Ø¥Ø¬Ø§Ø¨Ø© Ø®Ø§Ø·Ø¦Ø©! Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø© Ù‡ÙŠ: {current_quiz['answer']}"
            
        if current_quiz['explanation']:
            feedback += f"\n\nğŸ’¡ Ø§Ù„ØªÙØ³ÙŠØ±: {current_quiz['explanation']}"
            
        bot2.send_message(chat_id, feedback)
        
        # Ø§Ù„Ø§Ù†ØªÙ‚Ø§Ù„ Ù„Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„ØªØ§Ù„ÙŠ
        state['current_index'] += 1
        if state['current_index'] < len(state['quizzes']):
            self.send_question(chat_id, bot2)
        else:
            self.handle_quiz_end(chat_id, bot2)

    def handle_quiz_end(self, chat_id, bot2, error=False):
        state = self.active_quizzes.pop(chat_id, None)
        if not state:
            return
            
        if error:
            bot2.send_message(chat_id, "âš ï¸ Ø­Ø¯Ø« Ø®Ø·Ø£ ÙÙŠ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±. ÙŠØ±Ø¬Ù‰ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù„Ø§Ø­Ù‚Ù‹Ø§")
            return
            
        total = len(state['quizzes'])
        score = state['score']
        quiz_code = state['quiz_code']

                # Ø­ÙØ¸ Ø§Ù„Ù†ØªÙŠØ¬Ø© ÙÙŠ DB
        conn = sqlite3.connect("quiz_users.db")
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE user_quizzes
            SET score = ?, total = ?, timestamp = ?
            WHERE quiz_code = ?
        """, (score, total, datetime.now().isoformat(), state['quiz_code']))
        conn.commit()
        conn.close()

        keyboard = types.InlineKeyboardMarkup()
        keyboard.add(types.InlineKeyboardButton("ğŸ”„ Ø¥Ø¹Ø§Ø¯Ø© Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±", callback_data=f"retry:{quiz_code}"))
        keyboard.add(types.InlineKeyboardButton("ğŸ“¤ Ù…Ø´Ø§Ø±ÙƒØ© Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±", callback_data=f"share_quiz:{quiz_code}"))
        keyboard.add(types.InlineKeyboardButton("â¡ï¸ Ø§Ù„Ø¹ÙˆØ¯Ø© Ø§Ù„Ù‰ TestGenie âœ¨", url="https://t.me/Oiuhelper_bot"))
        
        end_msg = bot2.send_message(
            chat_id,
            f"ğŸ‰ Ø§Ù†ØªÙ‡Ù‰ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±!\n\nÙ†ØªÙŠØ¬ØªÙƒ: {score}/{total}\n\nÙ…Ø§Ø°Ø§ ØªØ±ÙŠØ¯ Ø£Ù† ØªÙØ¹Ù„ Ø§Ù„Ø¢Ù†ØŸ",
            reply_markup=keyboard
        )

# ØªÙ‡ÙŠØ¦Ø© Ø§Ù„Ù…Ø¯ÙŠØ±
quiz_manager = QuizManager()


@bot2.poll_answer_handler()
def handle_poll_answer(poll_answer):
    quiz_manager.handle_answer(poll_answer, bot2)


# Ø¯Ø§Ù„Ø© Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª (Ø§Ù„Ù†Ø³Ø®Ø© Ø§Ù„Ù…Ø­Ø³Ù†Ø©)
def send_quizzes(chat_id, quizzes, message_id=None):
    try:
        # Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªØ­Ù…ÙŠÙ„
        msg = bot2.send_message(chat_id, "â³ Ø¬Ø§Ø±ÙŠ ØªØ­Ø¶ÙŠØ± Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±...")
        
        # ØªØ®Ø²ÙŠÙ† Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
        quiz_code = store_quiz(chat_id, quizzes)
        if not quiz_code:
            raise Exception("Failed to store quiz")
        
        # Ø¨Ø¯Ø¡ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
        if quiz_manager.start_quiz(chat_id, quiz_code, bot2):
            bot.delete_message(chat_id, msg.message_id)
        else:
            bot2.edit_message_text("âŒ ÙØ´Ù„ Ø¨Ø¯Ø¡ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±", chat_id, msg.message_id)
            
    except Exception as e:
        print(f"Error in send_quizzes: {e}")
        bot2.send_message(chat_id, "Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø¨Ø¯Ø¡ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±")







import logging
import threading

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

state_lock = threading.Lock()
user_states = {}  # global



# -------------------------------------------------------------------
#                  Telegram Bot Handlers
# -------------------------------------------------------------------


@bot2.message_handler(commands=['start'])
def unified_start_handler(message):
    if message.chat.type != "private":
        return
    
    chat_id = message.chat.id
    args = message.text.split()
    
    if len(args) == 1:
        bot2.send_message(chat_id, """
        ğŸ‘‹ Ø£Ù‡Ù„Ø§Ù‹ Ø¨Ùƒ ÙÙŠ Quizzy! ğŸ˜Š
        Ø£Ù†Ø§ Ù‡Ù†Ø§ Ù„Ù…Ø³Ø§Ø¹Ø¯ØªÙƒ ÙÙŠ Ø¥Ø¬Ø±Ø§Ø¡ Ø¥Ø®ØªØ¨Ø§Ø±Ø§Øª Ø³Ø±ÙŠØ¹Ø© ÙˆØ§Ù„Ø§Ø³ØªÙ…ØªØ§Ø¹ Ø¨Ø§Ù„ØªØ¹Ù„Ù….
        """)
        return
    
    # Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„ÙƒÙˆØ¯
    param = args[1]
    quiz_code = param[5:] if param.startswith("quiz_") else param

    # Ø­Ø§Ù„Ø© Ø®Ø§ØµØ© Ù„Ø§Ø®ØªØ¨Ø§Ø± sample
    if quiz_code == "sample":
        loading_msg = bot2.send_message(chat_id, "ğŸ§  Ø¬Ø§Ø±ÙŠ ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø§Ù„ØªØ¬Ø±ÙŠØ¨ÙŠ...")
        if not quiz_manager.start_quiz(chat_id, quiz_code, bot2, loading_msg.message_id):
            bot2.edit_message_text(
                chat_id=chat_id,
                message_id=loading_msg.message_id,
                text="âŒ Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø§Ù„ØªØ¬Ø±ÙŠØ¨ÙŠ."
            )
        return
    
    # Ø£ÙŠ ÙƒÙˆØ¯ Ø§Ø®ØªØ¨Ø§Ø± Ø¢Ø®Ø±
    loading_msg = bot2.send_message(chat_id, "ğŸ§  Ø¬Ø§Ø±ÙŠ ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±...")
    if not quiz_manager.start_quiz(chat_id, quiz_code, bot2, loading_msg.message_id):
        bot2.edit_message_text(
            chat_id=chat_id,
            message_id=loading_msg.message_id,
            text="âŒ Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø£Ùˆ Ù‚Ø¯ ØªÙƒÙˆÙ† ØµÙ„Ø§Ø­ÙŠØªÙ‡ Ø§Ù†ØªÙ‡Øª."
            )
    # Ø¥Ø°Ø§ ÙƒØ§Ù† ÙƒÙˆØ¯ Ù…Ø®ØªÙ„ÙØŒ ØªØªØµØ±Ù Ø­Ø³Ø¨ Ø§Ù„ÙƒÙˆØ¯ Ø¹Ø§Ø¯ÙŠ
    # ... Ø¨Ø§Ù‚ÙŠ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© ...
    
@bot2.callback_query_handler(func=lambda c: True)
def handle_main_menu(c):
    
    if c.message.chat.type != "private":
        return
    uid = c.from_user.id
    data = c.data
    chat_id = c.message.chat.id
    message_id = c.message.message_id

    if data.startswith("retry:"):
        quiz_code = data[6:]
        quiz_manager.start_quiz(chat_id, quiz_code, bot2)
            
        
    elif data.startswith("share_quiz:"):
        quiz_code = data[6:]
        chat_id = c.message.chat.id  # â† ØªØ£ÙƒØ¯ Ù…Ù† ØªØ¹ÙŠÙŠÙ† chat_id Ù‡Ù†Ø§

        try:
            user_chat = bot2.get_chat(uid)
            shared_by_name = user_chat.first_name or user_chat.username or f"user_{uid}"
        except Exception:
            shared_by_name = "ØµØ¯ÙŠÙ‚Ùƒ"

        log_quiz_share(quiz_code, uid, shared_by_name)
        file_path = user_files[uid]
        
        share_link = f"https://t.me/QuizzyAI_bot?start=quiz_{quiz_code}"
        
        msg_text_share = f"""ğŸ“¢ {shared_by_name} Ø£Ø±Ø³Ù„ Ù„Ùƒ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±!  

ğŸ“‚ Ø§Ù„Ù…Ù„Ù: {msg.document.file_name}

Ø¬Ø±Ø¨Ù‡ ÙˆØ§Ø®ØªØ¨Ø± Ù…Ø¹Ù„ÙˆÙ…Ø§ØªÙƒ ğŸ‘‡  
{share_link}
"""
        msg_text = f"""<b>ğŸ‰ Ø´Ø§Ø±Ùƒ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù…Ø¹ Ø²Ù…Ù„Ø§Ø¦Ùƒ!</b>

    Ø§Ù†Ø³Ø® Ø§Ù„Ø±Ø§Ø¨Ø· Ø£Ø¯Ù†Ø§Ù‡ Ø£Ùˆ Ø§Ø¶ØºØ· Ù„ÙØªØ­Ù‡ Ù…Ø¨Ø§Ø´Ø±Ø©:
    ğŸ”— <a href="{share_link}">{share_link}</a>

    ğŸ“ Ø¹Ù†Ø¯ ÙØªØ­ Ø§Ù„Ø±Ø§Ø¨Ø·ØŒ Ø³ÙŠØ¨Ø¯Ø£ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ Ø¨Ø¥Ø°Ù† Ø§Ù„Ù„Ù‡.  
    ğŸ“¢ Ø¨Ù…Ø´Ø§Ø±ÙƒØªÙƒ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù‚Ø¯ ÙŠØµÙŠØ± Ø¹Ø§Ù…Ù‹Ø§.
    """

        keyboard = types.InlineKeyboardMarkup()
        keyboard.add(
            types.InlineKeyboardButton("ğŸ”— Ù†Ø³Ø® Ø§Ù„Ø±Ø§Ø¨Ø·", switch_inline_query=msg_text_share),
            types.InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home")
        )

        bot2.edit_message_text(msg_text, chat_id=chat_id, message_id=message_id, parse_mode="HTML", reply_markup=keyboard)


@bot.message_handler(commands=['start'])
def unified_start_handler(message):
    # âœ… ØªØ¬Ø§Ù‡Ù„ Ø§Ù„Ø±Ø³Ø§Ø¦Ù„ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø§Øª
    if message.chat.type != "private":
        return

    chat_id = message.chat.id
    args = message.text.split()
    uid = message.from_user.id
    
    if not can_generate(uid):
        add_external_user(uid)
    
    # âœ… Ø¥Ø°Ø§ ÙˆÙØ¬Ø¯ Ø¨Ø§Ø±Ø§Ù…ÙŠØªØ± (Ù…Ø«Ù„ quiz_ab12cd Ø£Ùˆ anki_sample)
    if len(args) > 1:
        param = args[1]

        # âœ… Ø¥Ø°Ø§ ÙƒØ§Ù† Ø¨Ø§Ø±Ø§Ù…ÙŠØªØ± anki_sample
        if param == "anki_sample":
            user_states[chat_id] = "awaiting_anki_file_ai"  # Ø­ÙØ¸ Ø§Ù„Ø­Ø§Ù„Ø©
            bot.send_message(
                "ğŸ“ Ø¯Ø¹Ù†Ø§ Ù†Ø¨Ø¯Ø£ Ø¨Ø¥Ù†Ø´Ø§Ø¡ **Ù…Ù„Ù Ø¨Ø·Ø§Ù‚Ø§ØªÙƒ Ø§Ù„Ø£ÙˆÙ„**!\n"
                "ğŸ“‚ Ø£Ø±Ø³Ù„ Ù…Ù„Ù **PDF** Ø£Ùˆ **DOCX** Ø£Ùˆ **PPTX**ØŒ Ø£Ùˆ Ø­ØªÙ‰ Ù†ØµÙ‹Ø§ Ù…Ø¨Ø§Ø´Ø±Ù‹Ø§ ğŸ“œ.\n"
                "Ø³ÙŠØªÙ… ØªÙˆÙ„ÙŠØ¯ Ù…Ù„Ù **Ø£Ù†ÙƒÙŠ** Ù…Ø®ØµØµ Ù„Ùƒ ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ ğŸ¯",
                chat_id=chat_id,
                parse_mode="Markdown"
            )

            return

        # âœ… Ù…Ø¹Ø§Ù„Ø¬Ø© Ø±ÙˆØ§Ø¨Ø· Ø§Ù„Ù…Ø´Ø§Ø±ÙƒØ© Ù…Ø«Ù„: ?start=quiz_ab12cd
        quiz_code = param[5:] if param.startswith("quiz_") else param

        loading_msg = bot.send_message(chat_id, "ğŸ§  Ø¬Ø§Ø±ÙŠ ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±...")

        # âœ… Ù…Ø­Ø§ÙˆÙ„Ø© Ø¨Ø¯Ø¡ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
        if not quiz_manager.start_quiz(chat_id, quiz_code, bot):
            bot.edit_message_text(
                chat_id=chat_id,
                message_id=loading_msg.message_id,
                text="âŒ Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø£Ùˆ Ø§Ù†ØªÙ‡Øª ØµÙ„Ø§Ø­ÙŠØªÙ‡."
            )
        return

    # âœ… Ø¥Ø°Ø§ Ù„Ù… ÙŠÙˆØ¬Ø¯ Ø¨Ø§Ø±Ø§Ù…ÙŠØªØ± â†’ Ø¹Ø±Ø¶ Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©
    send_main_menu(chat_id)


def send_main_menu(chat_id, message_id=None):
    keyboard = InlineKeyboardMarkup(row_width=2)
    buttons = [
        InlineKeyboardButton("ğŸ“ ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±", callback_data="go_generate"),
        InlineKeyboardButton("ğŸ“š Ù…Ø±Ø§Ø¬Ø¹Ø© Ø³Ø±ÙŠØ¹Ø©", callback_data="soon_review"),
        InlineKeyboardButton("ğŸ“„ Ù…Ù„Ø®Øµ PDF", callback_data="soon_summary"),
        InlineKeyboardButton("ğŸ§  Ø¨Ø·Ø§Ù‚Ø§Øª Anki", callback_data="anki"),
        InlineKeyboardButton("ğŸ® Ø£Ù„Ø¹Ø§Ø¨ ØªØ¹Ù„ÙŠÙ…ÙŠØ©", callback_data="go_games"),
        InlineKeyboardButton("âš™ï¸ Ø­Ø³Ø§Ø¨ÙŠ", callback_data="go_account_settings"),
    ]
    keyboard.add(*buttons)
    keyboard.add(InlineKeyboardButton("â• Ø£Ø¶ÙÙ†ÙŠ Ø¥Ù„Ù‰ Ù…Ø¬Ù…ÙˆØ¹Ø©", url=f"https://t.me/{bot.get_me().username}?startgroup=true"))

    text = (
        "ğŸ‘‹ <b>Ø£Ù‡Ù„Ø§Ù‹ Ø¨Ùƒ ÙÙŠ TestGenie!</b> âœ¨\n\n"
        "ğŸ¯ Ø£Ø¯ÙˆØ§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ© Ø°ÙƒÙŠØ© Ø¨ÙŠÙ† ÙŠØ¯ÙŠÙƒ:\n"
        "- Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ù…Ù† Ù…Ù„ÙØ§ØªÙƒ\n"
        "- Ø¨Ø·Ø§Ù‚Ø§Øª Ù…Ø±Ø§Ø¬Ø¹Ø© (Anki)\n"
        "- Ù…Ù„Ø®ØµØ§Øª PDF/Word <i>(Ù‚Ø±ÙŠØ¨Ø§Ù‹)</i>\n"
        "- Ø£Ù„Ø¹Ø§Ø¨ ØªØ¹Ù„ÙŠÙ…ÙŠØ© Ù…Ù…ØªØ¹Ø©\n\n"
        "ğŸ“Œ ÙƒÙ„ Ù…Ø§ ØªØ­ØªØ§Ø¬Ù‡ Ù„ØªØªØ¹Ù„Ù‘Ù… Ø¨Ø°ÙƒØ§Ø¡... Ø¨ÙŠÙ† ÙŠØ¯ÙŠÙƒ Ø§Ù„Ø¢Ù†.\n\n"
        "ğŸ‘‡ Ø§Ø®ØªØ± Ù…Ø§ ÙŠÙ†Ø§Ø³Ø¨Ùƒ ÙˆØ§Ø¨Ø¯Ø£ Ø§Ù„Ø¢Ù†:"
    )

    if message_id:
        bot.edit_message_text(
            text,
            chat_id=chat_id,
            message_id=message_id,
            reply_markup=keyboard,
            parse_mode="HTML"
        )
    else:
        bot.send_message(
            chat_id,
            text,
            reply_markup=keyboard,
            parse_mode="HTML"
        )



@bot.callback_query_handler(func=lambda call: call.data.startswith("rate_"))
def handle_rating(call):
    uid = call.from_user.id
    username = call.from_user.username or "Ù…Ø³ØªØ®Ø¯Ù…"
    rating = call.data.replace("rate_", "")

    if rating == "ignore":
        bot.answer_callback_query(call.id, "âœ… ØªÙ… ØªØ¬Ø§Ù‡Ù„ Ø§Ù„ØªÙ‚ÙŠÙŠÙ….")
        return

    # Ø¥Ø´Ø¹Ø§Ø± Ø§Ù„Ø£Ø¯Ù…Ù† Ø¹Ø¨Ø± bot3
    bot3.send_message(
        ADMIN_ID,
        f"â­ ØªÙ‚ÙŠÙŠÙ… Ø¬Ø¯ÙŠØ¯ Ù…Ù† @{username} (UID: {uid})\n\nØ§Ù„ØªÙ‚ÙŠÙŠÙ…: {rating} Ù†Ø¬ÙˆÙ…"
    )

    # Ø±Ø³Ø§Ù„Ø© Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù…
    bot.send_message(
        call.message.chat.id,
        f"ğŸ’Œ Ø´ÙƒØ±Ø§Ù‹ Ù„Ùƒ @{username}! ØªÙ… Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ø§Ø­Ø¸Ø§ØªÙƒ ({rating}â­) Ø¥Ù„Ù‰ ÙØ±ÙŠÙ‚ TestGenie ğŸ™"
    )

    bot.answer_callback_query(call.id, "âœ… ØªÙ… ØªØ³Ø¬ÙŠÙ„ ØªÙ‚ÙŠÙŠÙ…ÙƒØŒ Ø´ÙƒØ±Ø§Ù‹ Ù„Ùƒ!")


@bot.callback_query_handler(func=lambda c: True)
def handle_main_menu(c):
    try:
        bot.answer_callback_query(c.id)
    except:
        pass

    if c.message.chat.type != "private":
        return
    try: 
    
        data = c.data
        chat_id = c.message.chat.id
        message_id = c.message.message_id
        uid = c.from_user.id
        logging.info("Callback received: uid=%s data=%s", uid, data)


        # Ø±Ø¯ÙˆØ¯ Ø®Ø§Ø·Ø¦Ø© Ø¹Ø´ÙˆØ§Ø¦ÙŠØ© ØªØ¸Ù‡Ø± Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù…
        wrong_responses = [
            "âŒ Ø®Ø·Ø£! Ø¬Ø±Ø¨ Ù…Ø¬Ø¯Ø¯Ù‹Ø§ ğŸ˜‰\nâœ… Ø§Ù„ØµØ­ÙŠØ­: {correct}",
            "ğŸš« Ù„Ù„Ø£Ø³ÙØŒ Ù„ÙŠØ³Øª Ø§Ù„ØµØ­ÙŠØ­Ø©!\nâœ… Ø§Ù„Ø¬ÙˆØ§Ø¨: {correct}",
            "ğŸ˜… Ù„ÙŠØ³Øª Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø©ØŒ Ø§Ù„Ø¬ÙˆØ§Ø¨ Ù‡Ùˆ: {correct}",
            "âŒ Ù„Ø§ØŒ Ø­Ø§ÙˆÙ„ Ù…Ø±Ø© Ø£Ø®Ø±Ù‰!\nâœ”ï¸ Ø§Ù„ØµØ­ÙŠØ­ Ù‡Ùˆ: {correct}"
        ]




    # ---------- ØµÙØ­Ø© Ø§Ù„ÙØ¦Ø§Øª Ø§Ù„Ø£ÙˆÙ„Ù‰ ----------
        if data == "go_generate":
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("ğŸ©º Ø§Ù„Ø·Ø¨ ÙˆØ§Ù„ØµØ­Ø©", "category:health:page1"),
                ("ğŸ› ï¸ Ø§Ù„Ù‡Ù†Ø¯Ø³Ø©", "category:engineering:page1"),
                ("ğŸ’» Ø§Ù„Ø­Ø§Ø³ÙˆØ¨", "category:computer:page1"),
                ("ğŸ“Š Ø§Ù„Ø¥Ø¯Ø§Ø±Ø©", "category:business:page1"),
                ("ğŸ—£ï¸ Ø§Ù„Ù„ØºØ§Øª", "category:languages:page1"),
            ]
            # Ø¨Ù†Ø§Ø¡ Ø§Ù„ØµÙÙˆÙ Ø¨Ø£Ù…Ø§Ù† (Ù„Ø§ Ù†Ù…Ø±Ø± None)
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            keyboard.add(InlineKeyboardButton("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("â¡ï¸ Ø§Ù„Ù…Ø²ÙŠØ¯ Ù…Ù† Ø§Ù„ØªØ®ØµØµØ§Øª", callback_data="go_next"))
            keyboard.add(InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "ğŸ¯ *Ø§Ø®ØªØ± Ù…Ø¬Ø§Ù„ ØªØ®ØµØµÙƒ* (1/2)\n\nØ­Ø¯Ø¯ Ø§Ù„ÙØ¦Ø© Ø§Ù„Ø£Ù‚Ø±Ø¨ Ù„ØªØ®ØµØµÙƒ Ù…Ù† Ø§Ù„Ù‚Ø§Ø¦Ù…Ø©:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
            )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø¹Ø±Ø¶ Ø§Ù„ØµÙØ­Ø© Ø§Ù„Ø£ÙˆÙ„Ù‰:", e)
                # Ù…Ø­Ø§ÙˆÙ„Ø© Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© Ø¬Ø¯ÙŠØ¯Ø©
                try:
                    bot.send_message(chat_id, "ğŸ¯ Ø§Ø®ØªØ± Ù…Ø¬Ø§Ù„ ØªØ®ØµØµÙƒ (1/2).", reply_markup=keyboard)
                except Exception as e2:
                    print("ÙØ´Ù„ Ø§Ù„Ø¥Ø±Ø³Ø§Ù„:", e2)
            return

    # ---------- ØµÙØ­Ø© Ø§Ù„ÙØ¦Ø§Øª Ø§Ù„Ø«Ø§Ù†ÙŠØ© ----------
        if data == "go_next":
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("ğŸ“¿ Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø¥Ø³Ù„Ø§Ù…ÙŠØ©", "category:islamic:page2"),
                ("âš–ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† ÙˆØ§Ù„Ø³ÙŠØ§Ø³Ø©", "category:law:page2"),
                ("ğŸ§ª Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø·Ø¨ÙŠØ¹ÙŠØ©", "category:science:page2"),
                ("ğŸ¨ Ø§Ù„ÙÙ†ÙˆÙ†", "category:arts:page2"),
                ("ğŸ‘©â€ğŸ« Ø§Ù„ØªØ±Ø¨ÙŠØ©", "category:education:page2"),
            ]
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            keyboard.add(InlineKeyboardButton("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("â¬…ï¸ Ø§Ù„Ø¹ÙˆØ¯Ø©", callback_data="go_prev"))
            keyboard.add(InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "ğŸ“ *Ø§Ù„Ù…Ø²ÙŠØ¯ Ù…Ù† Ø§Ù„ØªØ®ØµØµØ§Øª* (2/2)\n\nØ§Ø®ØªØ± Ù…Ù† Ø§Ù„ÙØ¦Ø§Øª Ø§Ù„ØªØ§Ù„ÙŠØ©:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø¹Ø±Ø¶ Ø§Ù„ØµÙØ­Ø© Ø§Ù„Ø«Ø§Ù†ÙŠØ©:", e)
            return

    # ---------- Ø§Ù„Ø¹ÙˆØ¯Ø© Ù„Ù„ØµÙØ­Ø© Ø§Ù„Ø£ÙˆÙ„Ù‰ ----------
        if data == "go_prev":
        # Ø¥Ø¹Ø§Ø¯Ø© ØªÙˆØ¬ÙŠÙ‡ Ø¥Ù„Ù‰ Ù†ÙØ³ ØªÙ†ÙÙŠØ° go_generate
        # Ù†Ø³ØªØ¯Ø¹ÙŠ Ù†ÙØ³ Ø§Ù„ÙƒÙˆØ¯ Ø£Ùˆ Ù†Ø¹ÙŠØ¯ Ø§Ù„ØªÙˆØ¬ÙŠÙ‡:
        # Ù„Ø¥Ø¹Ø§Ø¯Ø© Ø§Ù„Ø§Ø³ØªØ®Ø¯Ø§Ù… ÙŠÙ…ÙƒÙ† Ø¹Ù…Ù„ Ø¯Ø§Ù„Ø© build_go_generate_keyboard()
            bot.answer_callback_query(c.id)
        # Ù‡Ù†Ø§ Ù†Ø¹ÙŠØ¯ Ù†ÙØ³ Ø§Ù„ÙƒÙˆØ¯ Ø£Ø¹Ù„Ø§Ù‡ Ø£Ùˆ Ù†Ù†Ø§Ø¯ÙŠ ØªÙ†ÙÙŠØ° go_generate:
        # Ø£Ø³Ù‡Ù„: Ù†Ø¹ÙŠØ¯ ØªØ¹ÙŠÙŠÙ† data ÙˆÙ†ÙƒÙˆÙ„Ù‘Ø¯ Ø§Ù„ØµÙØ­Ø©:
        # ÙˆÙ„ÙƒÙ† Ù„ØªØ¨Ø³ÙŠØ·ØŒ Ù†Ø¹ÙŠØ¯ ÙˆØ§Ø¬Ù‡Ø© go_generate Ù…Ø¨Ø§Ø´Ø±Ø©:
        # (ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ø³ØªØ¯Ø¹Ø§Ø¡ Ø±Ø³Ø§Ù„Ø© Ù…ÙˆØ¬ÙˆØ¯Ø© Ù…Ø³Ø¨Ù‚Ù‹Ø§)
        # Ù†Ø¹ÙŠØ¯ Ù†ÙØ³ Ù„ÙˆØ­Ø© go_generate ÙƒÙ…Ø§ ÙÙŠ Ø§Ù„Ø£Ø¹Ù„Ù‰:
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                ("ğŸ©º Ø§Ù„Ø·Ø¨ ÙˆØ§Ù„ØµØ­Ø©", "category:health:page1"),
                ("ğŸ› ï¸ Ø§Ù„Ù‡Ù†Ø¯Ø³Ø©", "category:engineering:page1"),
                ("ğŸ’» Ø§Ù„Ø­Ø§Ø³ÙˆØ¨", "category:computer:page1"),
                ("ğŸ“Š Ø§Ù„Ø¥Ø¯Ø§Ø±Ø©", "category:business:page1"),
                ("ğŸ—£ï¸ Ø§Ù„Ù„ØºØ§Øª", "category:languages:page1"),
            ]
            for i in range(0, len(buttons), 2):
                row = buttons[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)
            keyboard.add(InlineKeyboardButton("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("â¡ï¸ Ø§Ù„Ù…Ø²ÙŠØ¯ Ù…Ù† Ø§Ù„ØªØ®ØµØµØ§Øª", callback_data="go_next"))
            keyboard.add(InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home"))

            try:
                bot.edit_message_text(
                    "ğŸ¯ *Ø§Ø®ØªØ± Ù…Ø¬Ø§Ù„ ØªØ®ØµØµÙƒ* (1/2)\n\nØ­Ø¯Ø¯ Ø§Ù„ÙØ¦Ø© Ø§Ù„Ø£Ù‚Ø±Ø¨ Ù„ØªØ®ØµØµÙƒ Ù…Ù† Ø§Ù„Ù‚Ø§Ø¦Ù…Ø©:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø¹Ø±Ø¶ Ø§Ù„ØµÙØ­Ø© Ø§Ù„Ø£ÙˆÙ„Ù‰ (go_prev):", e)
            return

    # ---------- Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ø®ØªÙŠØ§Ø± Ø§Ù„ÙØ¦Ø© (Ø¢Ù…Ù†Ø©) ----------
        if data.startswith("category:"):
            # Ø§Ù„ØµÙŠØºØ© Ø§Ù„Ù…ØªÙˆÙ‚Ø¹Ø©: "category:<key>:<page>"
            parts = data.split(":")
            if len(parts) != 3:
                print("callback category ØºÙŠØ± Ù…ØªÙˆØ§ÙÙ‚:", data)
                return
            _, cat_key, page = parts  # cat_key Ù…Ø«Ù„Ø§Ù‹ 'health' ØŒ page Ù…Ø«Ù„ 'page1'

            # Ø®Ø±ÙŠØ·Ø© Ø§Ù„ØªØ®ØµØµØ§Øª Ø§Ù„ÙØ±Ø¹ÙŠØ© (Ø§Ø³ØªØ®Ø¯Ù… Ø§Ù„Ù…ÙØ§ØªÙŠØ­ Ø§Ù„Ø¥Ù†Ø¬Ù„ÙŠØ²ÙŠØ©)
            SUBS = {
                "health": [
                    ("ğŸ§¬ Ø§Ù„Ø·Ø¨ Ø§Ù„Ø¨Ø´Ø±ÙŠ", "major:Ø·Ø¨_Ø¨Ø´Ø±ÙŠ"),
                    ("ğŸ’Š Ø§Ù„ØµÙŠØ¯Ù„Ø©", "major:ØµÙŠØ¯Ù„Ø©"),
                    ("ğŸ¥ Ø§Ù„ØªÙ…Ø±ÙŠØ¶", "major:ØªÙ…Ø±ÙŠØ¶"),
                    ("ğŸ”¬ Ø¹Ù„ÙˆÙ… Ø§Ù„Ù…Ø®ØªØ¨Ø±Ø§Øª", "major:Ø¹Ù„ÙˆÙ…_Ù…Ø®ØªØ¨Ø±Ø§Øª"),
                    ("ğŸ¦· Ø·Ø¨ Ø§Ù„Ø£Ø³Ù†Ø§Ù†", "major:Ø·Ø¨_Ø£Ø³Ù†Ø§Ù†"),
                ],
                "engineering": [
                    ("âš™ï¸ Ø§Ù„Ù…ÙŠÙƒØ§Ù†ÙŠÙƒØ§", "major:Ù‡Ù†Ø¯Ø³Ø©_Ù…ÙŠÙƒØ§Ù†ÙŠÙƒÙŠØ©"),
                    ("ğŸ§ª Ø§Ù„ÙƒÙŠÙ…ÙŠØ§Ø¦ÙŠØ©", "major:Ù‡Ù†Ø¯Ø³Ø©_ÙƒÙŠÙ…ÙŠØ§Ø¦ÙŠØ©"),
                    ("ğŸ’¡ Ø§Ù„ÙƒÙ‡Ø±Ø¨Ø§Ø¦ÙŠØ©", "major:Ù‡Ù†Ø¯Ø³Ø©_ÙƒÙ‡Ø±Ø¨Ø§Ø¦ÙŠØ©"),
                    ("ğŸ—ï¸ Ø§Ù„Ù…Ø¯Ù†ÙŠØ©", "major:Ù‡Ù†Ø¯Ø³Ø©_Ù…Ø¯Ù†ÙŠØ©"),
                    ("ğŸ”§ Ø§Ù„Ø¨Ø±Ù…Ø¬ÙŠØ§Øª", "major:Ù‡Ù†Ø¯Ø³Ø©_Ø¨Ø±Ù…Ø¬ÙŠØ§Øª"),
                    ("ğŸ“¡ Ø§Ù„Ø§ØªØµØ§Ù„Ø§Øª", "major:Ù‡Ù†Ø¯Ø³Ø©_Ø§ØªØµØ§Ù„Ø§Øª"),
                ],
                "computer": [
                    ("ğŸ’» Ø¹Ù„ÙˆÙ… Ø§Ù„Ø­Ø§Ø³ÙˆØ¨", "major:Ø¹Ù„ÙˆÙ…_Ø­Ø§Ø³ÙˆØ¨"),
                    ("ğŸ“± ØªØ·ÙˆÙŠØ± Ø§Ù„Ø¨Ø±Ù…Ø¬ÙŠØ§Øª", "major:ØªØ·ÙˆÙŠØ±_Ø¨Ø±Ù…Ø¬ÙŠØ§Øª"),
                    ("ğŸ”’ Ø£Ù…Ù† Ø§Ù„Ù…Ø¹Ù„ÙˆÙ…Ø§Øª", "major:Ø§Ù…Ù†_Ù…Ø¹Ù„ÙˆÙ…Ø§Øª"),
                    ("ğŸ¤– Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø§ØµØ·Ù†Ø§Ø¹ÙŠ", "major:Ø°ÙƒØ§Ø¡_Ø§ØµØ·Ù†Ø§Ø¹ÙŠ"),
                    ("ğŸ“Š Ø¹Ù„Ù… Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª", "major:Ø¹Ù„Ù…_Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª"),
                ],
                "business": [
                    ("ğŸ“ˆ Ø¥Ø¯Ø§Ø±Ø© Ø§Ù„Ø£Ø¹Ù…Ø§Ù„", "major:Ø§Ø¯Ø§Ø±Ø©_Ø§Ø¹Ù…Ø§Ù„"),
                    ("ğŸ’¹ Ø§Ù„Ø§Ù‚ØªØµØ§Ø¯", "major:Ø§Ù‚ØªØµØ§Ø¯"),
                    ("ğŸ“Š Ø§Ù„Ù…Ø­Ø§Ø³Ø¨Ø©", "major:Ù…Ø­Ø§Ø³Ø¨Ø©"),
                    ("ğŸ§® Ø§Ù„ØªØ³ÙˆÙŠÙ‚", "major:ØªØ³ÙˆÙŠÙ‚"),
                    ("ğŸ¦ Ø§Ù„ØªÙ…ÙˆÙŠÙ„", "major:ØªÙ…ÙˆÙŠÙ„"),
                ],
                "languages": [
                    ("ğŸŒ Ø§Ù„Ø¥Ù†Ø¬Ù„ÙŠØ²ÙŠØ©", "major:Ù„ØºØ©_Ø§Ù†Ø¬Ù„ÙŠØ²ÙŠØ©"),
                    ("ğŸ‡«ğŸ‡· Ø§Ù„ÙØ±Ù†Ø³ÙŠØ©", "major:Ù„ØºØ©_ÙØ±Ù†Ø³ÙŠØ©"),
                    ("ğŸ‡¸ğŸ‡¦ Ø§Ù„Ø¹Ø±Ø¨ÙŠØ©", "major:Ù„ØºØ©_Ø¹Ø±Ø¨ÙŠØ©"),
                    ("ğŸ“š Ø§Ù„ØªØ±Ø¬Ù…Ø©", "major:ØªØ±Ø¬Ù…Ø©"),
                    ("ğŸ‡©ğŸ‡ª Ø§Ù„Ø£Ù„Ù…Ø§Ù†ÙŠØ©", "major:Ù„ØºØ©_Ø£Ù„Ù…Ø§Ù†ÙŠØ©"),
                ],
                "islamic": [
                    ("ğŸ“œ Ø§Ù„ÙÙ‚Ù‡", "major:ÙÙ‚Ù‡"),
                    ("ğŸ’¡ Ø§Ù„Ø¹Ù‚ÙŠØ¯Ø©", "major:Ø¹Ù‚ÙŠØ¯Ø©"),
                    ("ğŸ“– Ø§Ù„ØªÙØ³ÙŠØ±", "major:ØªÙØ³ÙŠØ±"),
                    ("ğŸ•Œ Ø§Ù„Ø¯Ø±Ø§Ø³Ø§Øª Ø§Ù„Ø¥Ø³Ù„Ø§Ù…ÙŠØ©", "major:Ø¯Ø±Ø§Ø³Ø§Øª_Ø§Ø³Ù„Ø§Ù…ÙŠØ©"),
                    ("ğŸŒ™ Ø§Ù„Ø³ÙŠØ±Ø© Ø§Ù„Ù†Ø¨ÙˆÙŠØ©", "major:Ø³ÙŠØ±Ø©_Ù†Ø¨ÙˆÙŠØ©"),
                ],
                "law": [
                    ("ğŸ“œ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„Ø¯ÙˆÙ„ÙŠ", "major:Ù‚Ø§Ù†ÙˆÙ†_Ø¯ÙˆÙ„ÙŠ"),
                    ("ğŸ›ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„ÙˆØ·Ù†ÙŠ", "major:Ù‚Ø§Ù†ÙˆÙ†_ÙˆØ·Ù†ÙŠ"),
                    ("ğŸ—³ï¸ Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø³ÙŠØ§Ø³ÙŠØ©", "major:Ø¹Ù„ÙˆÙ…_Ø³ÙŠØ§Ø³ÙŠØ©"),
                    ("ğŸ‘®â€â™‚ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„Ø¬Ù†Ø§Ø¦ÙŠ", "major:Ù‚Ø§Ù†ÙˆÙ†_Ø¬Ù†Ø§Ø¦ÙŠ"),
                    ("âš–ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„ØªØ¬Ø§Ø±ÙŠ", "major:Ù‚Ø§Ù†ÙˆÙ†_ØªØ¬Ø§Ø±ÙŠ"),
                ],
                "science": [
                    ("ğŸ§ª Ø§Ù„ÙÙŠØ²ÙŠØ§Ø¡", "major:ÙÙŠØ²ÙŠØ§Ø¡"),
                    ("ğŸ”¬ Ø§Ù„ÙƒÙŠÙ…ÙŠØ§Ø¡", "major:ÙƒÙŠÙ…ÙŠØ§Ø¡"),
                    ("ğŸ”¢ Ø§Ù„Ø±ÙŠØ§Ø¶ÙŠØ§Øª", "major:Ø±ÙŠØ§Ø¶ÙŠØ§Øª"),
                    ("ğŸŒ¿ Ø§Ù„Ø£Ø­ÙŠØ§Ø¡", "major:Ø§Ø­ÙŠØ§Ø¡"),
                ],
                "arts": [
                    ("ğŸ­ Ø§Ù„ÙÙ†ÙˆÙ† Ø§Ù„Ø£Ø¯Ø§Ø¦ÙŠØ©", "major:ÙÙ†ÙˆÙ†_Ø§Ø¯Ø§Ø¦ÙŠØ©"),
                    ("ğŸ–¼ï¸ Ø§Ù„ÙÙ†ÙˆÙ† Ø§Ù„Ø¨ØµØ±ÙŠØ©", "major:ÙÙ†ÙˆÙ†_Ø¨ØµØ±ÙŠØ©"),
                    ("ğŸ“š Ø§Ù„Ø£Ø¯Ø¨", "major:Ø§Ø¯Ø¨"),
                    ("ğŸŒ Ø§Ù„ØªØ§Ø±ÙŠØ®", "major:ØªØ§Ø±ÙŠØ®"),
                    ("ğŸµ Ø§Ù„Ù…ÙˆØ³ÙŠÙ‚Ù‰", "major:Ù…ÙˆØ³ÙŠÙ‚Ù‰"),
                ],
                "education": [
                    ("ğŸ‘©â€ğŸ« Ø¹Ù„Ù… Ø§Ù„Ù†ÙØ³ Ø§Ù„ØªØ±Ø¨ÙˆÙŠ", "major:Ø¹Ù„Ù…_Ø§Ù„Ù†ÙØ³_ØªØ±Ø¨ÙˆÙŠ"),
                    ("ğŸ“˜ Ø§Ù„Ù…Ù†Ø§Ù‡Ø¬", "major:Ù…Ù†Ø§Ù‡Ø¬_ØªØ¯Ø±ÙŠØ³"),
                    ("ğŸ§© Ø§Ù„ØªØ±Ø¨ÙŠØ© Ø§Ù„Ø®Ø§ØµØ©", "major:ØªØ±Ø¨ÙŠØ©_Ø®Ø§ØµØ©"),
                    ("ğŸ“Š Ø¥Ø¯Ø§Ø±Ø© Ø§Ù„ØªØ¹Ù„ÙŠÙ…", "major:Ø§Ø¯Ø§Ø±Ø©_ØªØ¹Ù„ÙŠÙ…"),
                    ("ğŸ‘¶ Ø§Ù„Ø·ÙÙˆÙ„Ø© Ø§Ù„Ù…Ø¨ÙƒØ±Ø©", "major:Ø·ÙÙˆÙ„Ø©_Ù…Ø¨ÙƒØ±Ø©"),
                ],
            }

            sub_btns = SUBS.get(cat_key, [])
            keyboard = InlineKeyboardMarkup(row_width=1)
            for label, cb in sub_btns:
                keyboard.add(InlineKeyboardButton(label, callback_data=cb))

            keyboard.add(InlineKeyboardButton("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", callback_data="major_custom"))

            # Ø²Ø± Ø§Ù„Ø±Ø¬ÙˆØ¹: Ø¥Ø°Ø§ ÙƒØ§Ù†Øª Ø§Ù„ØµÙØ­Ø© page1 Ù†Ø±Ø¬Ø¹ Ø§Ù„Ù‰ go_generate ÙˆØ¥Ù„Ø§ Ù†Ø±Ø¬Ø¹ Ø§Ù„Ù‰ go_next
            back_to = "go_generate" if page == "page1" else "go_next"
            keyboard.add(InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data=back_to))

            try:
                pretty_name = cat_key.replace("_", " ")
                bot.edit_message_text(
                    f"ğŸ“ *{pretty_name}*\n\nğŸ‘‡ Ø§Ø®ØªØ± ØªØ®ØµØµÙƒ Ø§Ù„Ø¯Ù‚ÙŠÙ‚ Ù…Ù† Ø§Ù„Ù‚Ø§Ø¦Ù…Ø©:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ ØªØ¹Ø¯ÙŠÙ„ Ø§Ù„Ø±Ø³Ø§Ù„Ø© (category):", e)
                try:
                    bot.send_message(chat_id, f"ğŸ“ {pretty_name}\nØ§Ø®ØªØ± ØªØ®ØµØµÙƒ:", reply_markup=keyboard)
                except Exception as e2:
                    print("ÙØ´Ù„ Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© Ø§Ù„ÙØ±Ø¹ÙŠØ©:", e2)
            
            return



    # ---------- Ø§Ù„ØªØ®ØµØµ Ø§Ù„Ù…Ø®ØµØµ ----------
        if data == "major_custom":
            user_states[uid] = "awaiting_major"
            bot.edit_message_text("ğŸ“ Ø£Ø±Ø³Ù„ Ø§Ø³Ù… ØªØ®ØµØµÙƒ (Ù…Ø«Ø§Ù„: Ù‡Ù†Ø¯Ø³Ø© Ø·ÙŠØ±Ø§Ù†ØŒ Ø¹Ù„Ù… Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª):", chat_id=chat_id, message_id=message_id)
            
            return



    # ----------------- Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø¹Ù†Ø¯ Ø§Ø®ØªÙŠØ§Ø± Ø§Ù„ØªØ®ØµØµ -----------------


        if data.startswith("major:"):
            major_key = data.split(":", 1)[1]  # 'Ø·Ø¨_Ø¨Ø´Ø±ÙŠ' Ø£Ùˆ 'Ù‡Ù†Ø¯Ø³Ø©_Ù…ÙŠÙƒØ§Ù†ÙŠÙƒÙŠØ©'

            # Ø§Ø­ÙØ¸ Ø§Ù„ØªØ®ØµØµ ÙÙŠ DB (ØªØ£ÙƒØ¯ Ø£Ù† conn,cursor Ù…ÙˆØ¬ÙˆØ¯Ø§Ù†)
            try:
                cursor.execute("INSERT OR REPLACE INTO users(user_id, major) VALUES(?, ?)", (uid, major_key))
                conn.commit()
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø­ÙØ¸ Ø§Ù„ØªØ®ØµØµ:", e)

            # Ø³Ù„ÙˆÙƒ Ø®Ø§Øµ Ù„Ù€ Ø·Ø¨_Ø¨Ø´Ø±ÙŠ
            if major_key == "Ø·Ø¨_Ø¨Ø´Ø±ÙŠ":
                kb2 = InlineKeyboardMarkup(row_width=2)
                kb2.add(InlineKeyboardButton("ğŸ§  Ø§Ù„ÙˆØ¶Ø¹ Ø§Ù„Ù…ØªÙ‚Ø¯Ù…", callback_data="advanced"), InlineKeyboardButton("ğŸ“„ Ø§Ù„ÙˆØ¶Ø¹ Ø§Ù„Ù…Ø¨Ø§Ø´Ø±", callback_data="simple"))
                bot.edit_message_text(
                    "ğŸ§¬ *Ù†Ø¸Ø§Ù… Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ø§Ù„Ù…ØªÙ‚Ø¯Ù… - Ù„Ø·Ù„Ø¨Ø© Ø§Ù„Ø·Ø¨ Ø§Ù„Ø¨Ø´Ø±ÙŠ*\n\n"
                    "ÙŠØ³Ø±Ù‘Ù†Ø§ Ø¥Ø¹Ù„Ø§Ù…Ùƒ Ø¨Ø£Ù† Ù†Ø¸Ø§Ù…Ù†Ø§ Ø§Ù„Ø¬Ø¯ÙŠØ¯ Ù„ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª ÙŠØ¹ØªÙ…Ø¯ Ø¹Ù„Ù‰ Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø§ØµØ·Ù†Ø§Ø¹ÙŠ Ø§Ù„Ù…ØªÙ‚Ø¯Ù…ØŒ "
                    "ÙˆÙŠÙ‚ÙˆÙ… Ø¨Ø¥Ø¹Ø¯Ø§Ø¯ Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª *Ø°ÙƒÙŠØ© ÙˆÙ…ØªÙˆØ§Ø²Ù†Ø©* ØªØ­Ø§ÙƒÙŠ Ù‡ÙŠÙƒÙ„Ø© Ø§Ù„Ø§Ù…ØªØ­Ø§Ù†Ø§Øª Ø§Ù„ÙØ¹Ù„ÙŠØ© Ù…Ù† Ø­ÙŠØ« Ø§Ù„ØªØ¯Ø±Ø¬ ÙˆØ§Ù„ØªÙ†ÙˆØ¹.\n\n"
                    "ğŸ“Œ Ø§Ø®ØªØ± Ø£Ø­Ø¯ Ø§Ù„Ø£ÙˆØ¶Ø§Ø¹ Ø§Ù„ØªØ§Ù„ÙŠØ© Ù„Ø¨Ø¯Ø¡ ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±Ùƒ:",
                    chat_id=chat_id,
                    message_id=message_id,
                    reply_markup=kb2,
                    parse_mode="Markdown"
                )
            else:
                user_states[uid] = "awaiting_simple_test_file"

                sent_msg = bot.edit_message_text(
                    f"âœ… ØªÙ… ØªØ­Ø¯ÙŠØ¯ ØªØ®ØµØµÙƒ: *{major_key.replace('_', ' ')}*\nØ§Ù„Ø¢Ù† Ø£Ø±Ø³Ù„ Ù…Ù„Ù (PDF/DOCX/TXT) Ø£Ùˆ Ù†ØµÙ‹Ø§ Ù…Ø¨Ø§Ø´Ø±Ù‹Ø§ Ù„ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±Ùƒ.",
                    chat_id=chat_id, message_id=message_id, parse_mode="Markdown"
            )
                return
            return

    # ... Ø¨Ù‚ÙŠØ© Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø§Øª Ø§Ù„Ø£Ø®Ø±Ù‰ ...


        if data == "anki":
            bot.answer_callback_query(c.id)
            choice_markup = types.InlineKeyboardMarkup()
            choice_markup.row(
                types.InlineKeyboardButton("ğŸ“ ØªÙˆÙ„ÙŠØ¯ ÙŠØ¯ÙˆÙŠ", callback_data="manual_anki"),
                types.InlineKeyboardButton("ğŸ¤– ØªÙˆÙ„ÙŠØ¯ Ø¢Ù„ÙŠ", callback_data="ai_anki")
            )
            bot.edit_message_text(
                chat_id=chat_id,
                message_id=message_id,
                text="ğŸ”§ Ø­Ø¯Ø¯ Ø·Ø±ÙŠÙ‚Ø© Ø¥Ù†Ø´Ø§Ø¡ Ø¨Ø·Ø§Ù‚Ø§Øª Anki:",
                reply_markup=choice_markup
            )

        elif data == "manual_anki":
            with state_lock:
                user_states[uid] = "awaiting_anki_file_manual"
            bot.answer_callback_query(c.id, "âœ… ØªÙ… ØªÙØ¹ÙŠÙ„ Ø§Ù„ÙˆØ¶Ø¹ Ø§Ù„ÙŠØ¯ÙˆÙŠ â€” Ø£Ø±Ø³Ù„ Ø§Ù„Ø¨Ø·Ø§Ø¦Ù‚ ÙƒÙ†Øµ.")
            bot.edit_message_text(
                chat_id=chat_id, message_id=message_id,
                text=(
                    "âœï¸ *ØµÙŠØ§ØºØ© Ø¨Ø·Ø§Ù‚Ø§Øª Anki ÙŠØ¯ÙˆÙŠÙ‹Ø§* \n\n"
                    "ÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª Ø¨Ù†Ø³Ù‚ Ù…Ø­Ø¯Ø¯ ÙƒÙ…Ø§ ÙŠÙ„ÙŠ:\n\n"
                    "Ø§Ù„Ø³Ø¤Ø§Ù„ Ø£Ùˆ Ø§Ù„Ù…ØµØ·Ù„Ø­.\n"
                    "Ø§Ù„Ø¬ÙˆØ§Ø¨ Ø£Ùˆ Ø§Ù„ØªØ¹Ø±ÙŠÙ.\n"
                    "ØªØ§Ù‚ *(Ø¥Ø®ØªÙŠØ§Ø±ÙŠ)*\n\n"
                    "Ù…Ø«Ø§Ù„:\n"
                    "Ù…Ø§ Ù‡Ùˆ Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø§ØµØ·Ù†Ø§Ø¹ÙŠØŸ\n"
                    "Ù‡Ùˆ Ù‚Ø¯Ø±Ø© Ø§Ù„Ø¢Ù„Ø© Ø¹Ù„Ù‰ Ù…Ø­Ø§ÙƒØ§Ø© Ø§Ù„ØªÙÙƒÙŠØ± Ø§Ù„Ø¨Ø´Ø±ÙŠ.\n"                
                    "#ØªÙ‚Ù†ÙŠØ©\n\n"
                    "Ù…Ù† Ø§ÙƒØªØ´Ù Ø§Ù„Ø¬Ø§Ø°Ø¨ÙŠØ©ØŸ\n"
                    "Ù†ÙŠÙˆØªÙ†\n\n"
                    "Ù‡Ø°Ø§ Ø§Ù„Ù†Ø³Ù‚ ÙŠØ³Ø§Ø¹Ø¯ ÙÙŠ ØªÙˆÙ„ÙŠØ¯ Ù…Ù„Ù Anki Ø¨Ø³Ù‡ÙˆÙ„Ø© ÙˆØ¯Ù‚Ø©. ğŸ’¡"
                ),
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return


        elif data == "ai_anki":
            with state_lock:
                user_states[uid] = "awaiting_anki_file_ai"
            bot.answer_callback_query(c.id, "âœ… ØªÙ… ØªÙØ¹ÙŠÙ„ ÙˆØ¶Ø¹ Anki Ø¢Ù„ÙŠ â€” Ø£Ø±Ø³Ù„ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ø¢Ù†.")
            bot.edit_message_text(
                chat_id=chat_id, message_id=message_id,
                text=(
                    "âœ¨ *ØªÙˆÙ„ÙŠØ¯ Ø¨Ø·Ø§Ù‚Ø§Øª Anki Ø¨ÙˆØ§Ø³Ø·Ø© Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø¥ØµØ·Ù†Ø§Ø¹ÙŠ*\n\n"
                    "Ø§Ø±ÙØ¹ Ù…Ù„ÙÙ‹Ø§ Ø¨ØµÙŠØºØ© PDFØŒ WordØŒ Ø£Ùˆ Ù†Øµ Ø¹Ø§Ø¯ÙŠØŒ ÙˆØ³ÙŠÙ‚ÙˆÙ… Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø¥ØµØ·Ù†Ø§Ø¹ÙŠ Ø¨ØªØ­Ù„ÙŠÙ„ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙˆØªÙˆÙ„ÙŠØ¯ Ø¨Ø·Ø§Ù‚Ø§Øª Anki Ù…ØªÙˆØ§ÙÙ‚Ø© Ø¨Ø´ÙƒÙ„ ØªÙ„Ù‚Ø§Ø¦ÙŠ.\n\n"
                    "Ù…ÙŠØ²Ø© ÙØ±ÙŠØ¯Ø©:\n"
                    "- ØªØ­ÙˆÙŠÙ„ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø¥Ù„Ù‰ Ù…Ù„Ù Anki Ø¬Ø§Ù‡Ø² Ù„Ù„Ø§Ø³ØªØ®Ø¯Ø§Ù….\n"
                    "- ÙŠØ¯Ø¹Ù… Ù…Ø¬Ù…ÙˆØ¹Ø© ÙˆØ§Ø³Ø¹Ø© Ù…Ù† Ø§Ù„Ù…ÙˆØ§Ø¶ÙŠØ¹ ÙˆØ§Ù„Ù…ØµØ§Ø¯Ø±.\n"
                    "- ÙˆÙØ± Ø§Ù„ÙˆÙ‚Øª ÙˆØ§Ù„Ø¬Ù‡Ø¯ Ø¨Ø§Ø³ØªØ®Ø¯Ø§Ù… ØªÙ‚Ù†ÙŠØ© Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø¥ØµØ·Ù†Ø§Ø¹ÙŠ Ø§Ù„Ù…ØªÙ‚Ø¯Ù…Ø©.\n\n"
                    "Ø§Ø¨Ø¯Ø£ Ø§Ù„Ø¢Ù† ÙˆØ£Ø±Ø³Ù„ Ù…Ù„ÙÙƒ Ù„ØªØ­ÙˆÙŠÙ„Ù‡ Ø¥Ù„Ù‰ Ù…Ù„Ù Anki Ù…Ø®ØµØµ! ğŸ“š"
                ),
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return


        elif data == "go_account_settings":
            bot.answer_callback_query(c.id)
            settings_keyboard = InlineKeyboardMarkup()
            settings_keyboard.add(
                InlineKeyboardButton("ğŸ“ ØªØºÙŠÙŠØ± Ø§Ù„ØªØ®ØµØµ", callback_data="change_specialty"),
            )
            settings_keyboard.add(
                InlineKeyboardButton("ğŸ“‰ Ù…Ø³ØªÙˆÙ‰ Ø§Ù„Ø¥Ø®ØªØ¨Ø§Ø±Ø§Øª", callback_data="tests_level"),
            )
            settings_keyboard.add(
                InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data="go_back_home")
            )

            bot.edit_message_text(
                text="âš™ï¸ *Ø¥Ø¹Ø¯Ø§Ø¯Ø§Øª Ø§Ù„Ø­Ø³Ø§Ø¨*\n\n"
                    "ÙŠÙ…ÙƒÙ†Ùƒ ØªØ®ØµÙŠØµ ØªØ¬Ø±Ø¨ØªÙƒ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠØ© Ù‡Ù†Ø§.\n"
                    "Ø§Ø®ØªØ± Ù…Ø§ ØªØ±ØºØ¨ Ø¨ØªØ¹Ø¯ÙŠÙ„Ù‡ ğŸ‘‡",
                chat_id=c.message.chat.id,
                message_id=c.message.message_id,
                reply_markup=settings_keyboard,
                parse_mode="Markdown"
            )
        
        elif data == "go_games":
            raw = fetch_user_major(uid)

            if not row or not row[0]:
                user_states[uid] = "awaiting_major_for_games"
                bot.send_message(uid, "ğŸ§  Ù‚Ø¨Ù„ Ø£Ù† Ù†Ø¨Ø¯Ø£ Ø§Ù„Ù„Ø¹Ø¨ØŒ Ø£Ø®Ø¨Ø±Ù†Ø§ Ø¨ØªØ®ØµØµÙƒ:")
                return

            keyboard = InlineKeyboardMarkup(row_width=1)
            keyboard.add(
                InlineKeyboardButton("ğŸ”’ Ø§Ù„Ø¹Ø¨ ÙÙŠ Ø§Ù„Ø®Ø§Øµ", callback_data="game_private"),
                InlineKeyboardButton("ğŸ‘¥ Ø§Ù„Ø¹Ø¨ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø©", switch_inline_query="game"),
                InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home")
            )
            bot.edit_message_text(
                "ğŸ® Ø§Ø®ØªØ± Ø·Ø±ÙŠÙ‚Ø© Ø§Ù„Ù„Ø¹Ø¨:\n\n"
                "- ğŸ”’ ÙÙŠ Ø§Ù„Ø®Ø§Øµ (Ø£Ù„Ø¹Ø§Ø¨ Ø´Ø®ØµÙŠØ© Ø­Ø³Ø¨ ØªØ®ØµØµÙƒ)\n"
                "- ğŸ‘¥ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø© (Ø´Ø§Ø±Ùƒ Ø§Ù„Ø£ØµØ¯Ù‚Ø§Ø¡ Ø¨Ø§Ù„ØªØ­Ø¯ÙŠ!)",
                chat_id=chat_id,
                message_id=message_id,
                reply_markup=keyboard
            )
    
        elif data == "go_back_home":
            # Ø¥Ø¹Ø§Ø¯Ø© Ø¹Ø±Ø¶ ÙˆØ§Ø¬Ù‡Ø© Ø§Ù„Ø¨Ø¯Ø§ÙŠØ©
            keyboard = InlineKeyboardMarkup(row_width=2)
            buttons = [
                InlineKeyboardButton("ğŸ“ ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±", callback_data="go_generate"),
                InlineKeyboardButton("ğŸ“š Ù…Ø±Ø§Ø¬Ø¹Ø© Ø³Ø±ÙŠØ¹Ø©", callback_data="soon_review"),
                InlineKeyboardButton("ğŸ“„ Ù…Ù„Ø®Øµ PDF", callback_data="soon_summary"),
                InlineKeyboardButton("ğŸ§  Ø¨Ø·Ø§Ù‚Ø§Øª Anki", callback_data="anki"),
                InlineKeyboardButton("ğŸ® Ø£Ù„Ø¹Ø§Ø¨ ØªØ¹Ù„ÙŠÙ…ÙŠØ©", callback_data="go_games"),
                InlineKeyboardButton("âš™ï¸ Ø­Ø³Ø§Ø¨ÙŠ", callback_data="go_account_settings"),
            ]
            keyboard.add(*buttons)

            keyboard.add(InlineKeyboardButton("â• Ø£Ø¶ÙÙ†ÙŠ Ø¥Ù„Ù‰ Ù…Ø¬Ù…ÙˆØ¹Ø©", url=f"https://t.me/{bot.get_me().username}?startgroup=true"))

            bot.edit_message_text(
                "ğŸ‘‹ Ø£Ù‡Ù„Ø§ Ø¨Ùƒ ÙÙŠ *TestGenie* âœ¨\n\n"
                "ğŸ¯ Ø£Ø¯ÙˆØ§Øª ØªØ¹Ù„ÙŠÙ…ÙŠØ© Ø°ÙƒÙŠØ© Ø¨ÙŠÙ† ÙŠØ¯ÙŠÙƒ:\n"
                "- Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ù…Ù† Ù…Ù„ÙØ§ØªÙƒ\n"
                "- Ø¨Ø·Ø§Ù‚Ø§Øª Ù…Ø±Ø§Ø¬Ø¹Ø© (Anki)\n"
                "- Ù…Ù„Ø®ØµØ§Øª PDF/Word _(Ù‚Ø±ÙŠØ¨Ø§Ù‹)_\n"
                "- Ø£Ù„Ø¹Ø§Ø¨ ØªØ¹Ù„ÙŠÙ…ÙŠØ©\n\n"
                "ğŸ“Œ ÙƒÙ„ Ù…Ø§ ØªØ­ØªØ§Ø¬Ù‡ Ù„ØªØªØ¹Ù„Ù‘Ù… Ø¨Ø°ÙƒØ§Ø¡... Ø¨ÙŠÙ† ÙŠØ¯ÙŠÙƒ Ø§Ù„Ø¢Ù†.\n\n"
                "Ø§Ø®ØªØ± Ù…Ø§ ÙŠÙ†Ø§Ø³Ø¨Ùƒ ÙˆØ§Ø¨Ø¯Ø£ Ø§Ù„Ø¢Ù† ğŸ‘‡",
                chat_id=chat_id,
                message_id=message_id,
                reply_markup=keyboard,
                parse_mode="Markdown"
            )
                
        

    # ----------------- ØªØºÙŠÙŠØ± Ø§Ù„ØªØ®ØµØµ (Ø§Ù„ÙˆØ§Ø¬Ù‡Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ© Ù„Ù„ÙØ¦Ø§Øª) -----------------
        if data == "change_specialty":
            keyboard = InlineKeyboardMarkup(row_width=2)
            main_categories = [
                ("ğŸ©º Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„ØµØ­ÙŠØ©", "spec_category:health"),
                ("ğŸ› ï¸ Ø§Ù„Ù‡Ù†Ø¯Ø³Ø©", "spec_category:engineering"),
                ("ğŸ’» Ø¹Ù„ÙˆÙ… Ø§Ù„Ø­Ø§Ø³ÙˆØ¨", "spec_category:computer"),
                ("ğŸ“Š Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø¥Ø¯Ø§Ø±ÙŠØ©", "spec_category:business"),
                ("ğŸ—£ï¸ Ø§Ù„Ù„ØºØ§Øª", "spec_category:languages"),
                ("ğŸ“¿ Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø¥Ø³Ù„Ø§Ù…ÙŠØ©", "spec_category:islamic"),
                ("âš–ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† ÙˆØ§Ù„Ø³ÙŠØ§Ø³Ø©", "spec_category:law"),
                ("ğŸ”¬ Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø·Ø¨ÙŠØ¹ÙŠØ©", "spec_category:science"),
                ("ğŸ¨ Ø§Ù„ÙÙ†ÙˆÙ†", "spec_category:arts"),
                ("ğŸ‘©â€ğŸ« Ø§Ù„ØªØ±Ø¨ÙŠØ©", "spec_category:education"),
                ("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", "major_custom"),  # ØªØ®ØµØµ ÙŠØ¯ÙˆÙŠ
            ]

        # Ø¨Ù†Ø§Ø¡ Ø§Ù„ØµÙÙˆÙ Ø¨Ø´ÙƒÙ„ Ø¢Ù…Ù†
            for i in range(0, len(main_categories), 2):
                row = main_categories[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            # Ø²Ø± Ø±Ø¬ÙˆØ¹
            keyboard.add(InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data="go_account_settings"))

            try:
                bot.edit_message_text(
                    "ğŸ“ *Ø§Ø®ØªØ± Ù…Ø¬Ø§Ù„ ØªØ®ØµØµÙƒ*\n\nØ­Ø¯Ø¯ Ø§Ù„ÙØ¦Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ© Ø£ÙˆÙ„Ø§Ù‹ Ø«Ù… Ø§Ø®ØªØ± ØªØ®ØµØµÙƒ Ø§Ù„Ø¯Ù‚ÙŠÙ‚:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø¹Ø±Ø¶ change_specialty:", e)
                try:
                    bot.send_message(chat_id, "ğŸ“ Ø§Ø®ØªØ± Ù…Ø¬Ø§Ù„ ØªØ®ØµØµÙƒ:", reply_markup=keyboard)
                except Exception as e2:
                    print("Failed to send change_specialty:", e2)


    # ----------------- Ø¹Ø±Ø¶ Ø§Ù„ØªØ®ØµØµØ§Øª Ø§Ù„ÙØ±Ø¹ÙŠØ© Ø¨Ù†Ø§Ø¡Ù‹ Ø¹Ù„Ù‰ Ø§Ù„ÙØ¦Ø© -----------------
        if data.startswith("spec_category:"):
            # ØµÙŠØºØ© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª: spec_category:<key>
            try:
                _, cat_key = data.split(":", 1)
            except ValueError:
                print("spec_category callback ØºÙŠØ± Ù…ØªÙˆØ§ÙÙ‚:", data)
                return

            SUBS = {  
                "health": [  
                    ("ğŸ§¬ Ø§Ù„Ø·Ø¨ Ø§Ù„Ø¨Ø´Ø±ÙŠ", "change_major:Ø·Ø¨_Ø¨Ø´Ø±ÙŠ"),  
                    ("ğŸ’Š Ø§Ù„ØµÙŠØ¯Ù„Ø©", "change_major:ØµÙŠØ¯Ù„Ø©"),  
                    ("ğŸ¥ Ø§Ù„ØªÙ…Ø±ÙŠØ¶", "change_major:ØªÙ…Ø±ÙŠØ¶"),  
                    ("ğŸ”¬ Ø¹Ù„ÙˆÙ… Ø§Ù„Ù…Ø®ØªØ¨Ø±Ø§Øª", "change_major:Ø¹Ù„ÙˆÙ…_Ù…Ø®ØªØ¨Ø±Ø§Øª"),  
                    ("ğŸ¦· Ø·Ø¨ Ø§Ù„Ø£Ø³Ù†Ø§Ù†", "change_major:Ø·Ø¨_Ø£Ø³Ù†Ø§Ù†"),  
                ],  
                "engineering": [  
                    ("âš™ï¸ Ø§Ù„Ù…ÙŠÙƒØ§Ù†ÙŠÙƒØ§", "change_major:Ù‡Ù†Ø¯Ø³Ø©_Ù…ÙŠÙƒØ§Ù†ÙŠÙƒÙŠØ©"),  
                    ("ğŸ§ª Ø§Ù„ÙƒÙŠÙ…ÙŠØ§Ø¦ÙŠØ©", "change_major:Ù‡Ù†Ø¯Ø³Ø©_ÙƒÙŠÙ…ÙŠØ§Ø¦ÙŠØ©"),
                    ("ğŸ’¡ Ø§Ù„ÙƒÙ‡Ø±Ø¨Ø§Ø¦ÙŠØ©", "change_major:Ù‡Ù†Ø¯Ø³Ø©_ÙƒÙ‡Ø±Ø¨Ø§Ø¦ÙŠØ©"),  
                    ("ğŸ—ï¸ Ø§Ù„Ù…Ø¯Ù†ÙŠØ©", "change_major:Ù‡Ù†Ø¯Ø³Ø©_Ù…Ø¯Ù†ÙŠØ©"),  
                    ("ğŸ”§ Ø§Ù„Ø¨Ø±Ù…Ø¬ÙŠØ§Øª", "change_major:Ù‡Ù†Ø¯Ø³Ø©_Ø¨Ø±Ù…Ø¬ÙŠØ§Øª"),  
                    ("ğŸ“¡ Ø§Ù„Ø§ØªØµØ§Ù„Ø§Øª", "change_major:Ù‡Ù†Ø¯Ø³Ø©_Ø§ØªØµØ§Ù„Ø§Øª"),  
                ],  
                "computer": [  
                    ("ğŸ’» Ø¹Ù„ÙˆÙ… Ø§Ù„Ø­Ø§Ø³ÙˆØ¨", "change_major:Ø¹Ù„ÙˆÙ…_Ø­Ø§Ø³ÙˆØ¨"),  
                    ("ğŸ“± ØªØ·ÙˆÙŠØ± Ø§Ù„Ø¨Ø±Ù…Ø¬ÙŠØ§Øª", "change_major:ØªØ·ÙˆÙŠØ±_Ø¨Ø±Ù…Ø¬ÙŠØ§Øª"),  
                    ("ğŸ”’ Ø£Ù…Ù† Ø§Ù„Ù…Ø¹Ù„ÙˆÙ…Ø§Øª", "change_major:Ø§Ù…Ù†_Ù…Ø¹Ù„ÙˆÙ…Ø§Øª"),  
                    ("ğŸ¤– Ø§Ù„Ø°ÙƒØ§Ø¡ Ø§Ù„Ø§ØµØ·Ù†Ø§Ø¹ÙŠ", "change_major:Ø°ÙƒØ§Ø¡_Ø§ØµØ·Ù†Ø§Ø¹ÙŠ"),  
                    ("ğŸ“Š Ø¹Ù„Ù… Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª", "change_major:Ø¹Ù„Ù…_Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª"),  
                ],  
                "business": [  
                    ("ğŸ“ˆ Ø¥Ø¯Ø§Ø±Ø© Ø§Ù„Ø£Ø¹Ù…Ø§Ù„", "change_major:Ø§Ø¯Ø§Ø±Ø©_Ø§Ø¹Ù…Ø§Ù„"),  
                    ("ğŸ’¹ Ø§Ù„Ø§Ù‚ØªØµØ§Ø¯", "change_major:Ø§Ù‚ØªØµØ§Ø¯"),  
                    ("ğŸ“Š Ø§Ù„Ù…Ø­Ø§Ø³Ø¨Ø©", "change_major:Ù…Ø­Ø§Ø³Ø¨Ø©"),  
                    ("ğŸ§® Ø§Ù„ØªØ³ÙˆÙŠÙ‚", "change_major:ØªØ³ÙˆÙŠÙ‚"),  
                    ("ğŸ¦ Ø§Ù„ØªÙ…ÙˆÙŠÙ„", "change_major:ØªÙ…ÙˆÙŠÙ„"),  
                ],  
                "languages": [  
                    ("ğŸŒ Ø§Ù„Ø¥Ù†Ø¬Ù„ÙŠØ²ÙŠØ©", "change_major:Ù„ØºØ©_Ø§Ù†Ø¬Ù„ÙŠØ²ÙŠØ©"),  
                    ("ğŸ‡«ğŸ‡· Ø§Ù„ÙØ±Ù†Ø³ÙŠØ©", "change_major:Ù„ØºØ©_ÙØ±Ù†Ø³ÙŠØ©"),  
                    ("ğŸ‡¸ğŸ‡¦ Ø§Ù„Ø¹Ø±Ø¨ÙŠØ©", "change_major:Ù„ØºØ©_Ø¹Ø±Ø¨ÙŠØ©"),  
                    ("ğŸ“š Ø§Ù„ØªØ±Ø¬Ù…Ø©", "change_major:ØªØ±Ø¬Ù…Ø©"),  
                    ("ğŸ‡©ğŸ‡ª Ø§Ù„Ø£Ù„Ù…Ø§Ù†ÙŠØ©", "change_major:Ù„ØºØ©_Ø£Ù„Ù…Ø§Ù†ÙŠØ©"),  
                ],  
                "islamic": [  
                    ("ğŸ“œ Ø§Ù„ÙÙ‚Ù‡", "change_major:ÙÙ‚Ù‡"),  
                    ("ğŸ’¡ Ø§Ù„Ø¹Ù‚ÙŠØ¯Ø©", "change_major:Ø¹Ù‚ÙŠØ¯Ø©"),  
                    ("ğŸ“– Ø§Ù„ØªÙØ³ÙŠØ±", "change_major:ØªÙØ³ÙŠØ±"),  
                    ("ğŸ•Œ Ø§Ù„Ø¯Ø±Ø§Ø³Ø§Øª Ø§Ù„Ø¥Ø³Ù„Ø§Ù…ÙŠØ©", "change_major:Ø¯Ø±Ø§Ø³Ø§Øª_Ø§Ø³Ù„Ø§Ù…ÙŠØ©"),  
                    ("ğŸŒ™ Ø§Ù„Ø³ÙŠØ±Ø© Ø§Ù„Ù†Ø¨ÙˆÙŠØ©", "change_major:Ø³ÙŠØ±Ø©_Ù†Ø¨ÙˆÙŠØ©"),  
                ],  
                "law": [  
                    ("ğŸ“œ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„Ø¯ÙˆÙ„ÙŠ", "change_major:Ù‚Ø§Ù†ÙˆÙ†_Ø¯ÙˆÙ„ÙŠ"),  
                    ("ğŸ›ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„ÙˆØ·Ù†ÙŠ", "change_major:Ù‚Ø§Ù†ÙˆÙ†_ÙˆØ·Ù†ÙŠ"),  
                    ("ğŸ—³ï¸ Ø§Ù„Ø¹Ù„ÙˆÙ… Ø§Ù„Ø³ÙŠØ§Ø³ÙŠØ©", "change_major:Ø¹Ù„ÙˆÙ…_Ø³ÙŠØ§Ø³ÙŠØ©"),  
                    ("ğŸ‘®â€â™‚ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„Ø¬Ù†Ø§Ø¦ÙŠ", "change_major:Ù‚Ø§Ù†ÙˆÙ†_Ø¬Ù†Ø§Ø¦ÙŠ"),  
                    ("âš–ï¸ Ø§Ù„Ù‚Ø§Ù†ÙˆÙ† Ø§Ù„ØªØ¬Ø§Ø±ÙŠ", "change_major:Ù‚Ø§Ù†ÙˆÙ†_ØªØ¬Ø§Ø±ÙŠ"),  
                ],  
                "science": [  
                    ("ğŸ§ª Ø§Ù„ÙÙŠØ²ÙŠØ§Ø¡", "change_major:ÙÙŠØ²ÙŠØ§Ø¡"),  
                    ("ğŸ”¬ Ø§Ù„ÙƒÙŠÙ…ÙŠØ§Ø¡", "change_major:ÙƒÙŠÙ…ÙŠØ§Ø¡"),  
                    ("ğŸ”¢ Ø§Ù„Ø±ÙŠØ§Ø¶ÙŠØ§Øª", "change_major:Ø±ÙŠØ§Ø¶ÙŠØ§Øª"),  
                    ("ğŸŒ¿ Ø§Ù„Ø£Ø­ÙŠØ§Ø¡", "change_major:Ø§Ø­ÙŠØ§Ø¡"),  
                ],  
                "arts": [  
                    ("ğŸ­ Ø§Ù„ÙÙ†ÙˆÙ† Ø§Ù„Ø£Ø¯Ø§Ø¦ÙŠØ©", "change_major:ÙÙ†ÙˆÙ†_Ø§Ø¯Ø§Ø¦ÙŠØ©"),  
                    ("ğŸ–¼ï¸ Ø§Ù„ÙÙ†ÙˆÙ† Ø§Ù„Ø¨ØµØ±ÙŠØ©", "change_major:ÙÙ†ÙˆÙ†_Ø¨ØµØ±ÙŠØ©"),  
                    ("ğŸ“š Ø§Ù„Ø£Ø¯Ø¨", "change_major:Ø§Ø¯Ø¨"),  
                    ("ğŸŒ Ø§Ù„ØªØ§Ø±ÙŠØ®", "change_major:ØªØ§Ø±ÙŠØ®"),  
                    ("ğŸµ Ø§Ù„Ù…ÙˆØ³ÙŠÙ‚Ù‰", "change_major:Ù…ÙˆØ³ÙŠÙ‚Ù‰"),  
                ],  
                "education": [  
                    ("ğŸ‘©â€ğŸ« Ø¹Ù„Ù… Ø§Ù„Ù†ÙØ³ Ø§Ù„ØªØ±Ø¨ÙˆÙŠ", "change_major:Ø¹Ù„Ù…_Ø§Ù„Ù†ÙØ³_ØªØ±Ø¨ÙˆÙŠ"),  
                    ("ğŸ“˜ Ø§Ù„Ù…Ù†Ø§Ù‡Ø¬", "change_major:Ù…Ù†Ø§Ù‡Ø¬_ØªØ¯Ø±ÙŠØ³"),  
                    ("ğŸ§© Ø§Ù„ØªØ±Ø¨ÙŠØ© Ø§Ù„Ø®Ø§ØµØ©", "change_major:ØªØ±Ø¨ÙŠØ©_Ø®Ø§ØµØ©"),  
                    ("ğŸ“Š Ø¥Ø¯Ø§Ø±Ø© Ø§Ù„ØªØ¹Ù„ÙŠÙ…", "change_major:Ø§Ø¯Ø§Ø±Ø©_ØªØ¹Ù„ÙŠÙ…"),  
                    ("ğŸ‘¶ Ø§Ù„Ø·ÙÙˆÙ„Ø© Ø§Ù„Ù…Ø¨ÙƒØ±Ø©", "change_major:Ø·ÙÙˆÙ„Ø©_Ù…Ø¨ÙƒØ±Ø©"),  
                ],  
            }
            sub_btns = SUBS.get(cat_key, [])
            keyboard = InlineKeyboardMarkup(row_width=2)

            # Ø¨Ù†Ø§Ø¡ Ø§Ù„Ø£Ø²Ø±Ø§Ø± Ø¨ØµÙÙŠÙ†
            for i in range(0, len(sub_btns), 2):
                row = sub_btns[i:i+2]
                btns = [InlineKeyboardButton(label, callback_data=cb) for label, cb in row]
                keyboard.row(*btns)

            # Ø§Ø¶Ø§ÙØ© Ø®ÙŠØ§Ø±Ø§Øª Ø§Ù„Ù†Ù‡Ø§ÙŠØ©
            keyboard.add(InlineKeyboardButton("â“ ØªØ®ØµØµ Ø¢Ø®Ø±", callback_data="major_custom"))
            keyboard.add(InlineKeyboardButton("ğŸ”™ Ø§Ù„Ø¹ÙˆØ¯Ø©", callback_data="change_specialty"))
            keyboard.add(InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home"))

            try:
                pretty = cat_key.replace("_", " ")
                bot.edit_message_text(
                f"ğŸ“ *{pretty}*\n\nğŸ‘‡ Ø§Ø®ØªØ± ØªØ®ØµØµÙƒ Ø§Ù„Ø¯Ù‚ÙŠÙ‚ Ù…Ù† Ø§Ù„Ù‚Ø§Ø¦Ù…Ø©:",
                    chat_id=chat_id, message_id=message_id, reply_markup=keyboard, parse_mode="Markdown"
                )
            except Exception as e:
                print("Ø®Ø·Ø£ ÙÙŠ Ø¹Ø±Ø¶ Ø§Ù„ØªØ®ØµØµØ§Øª Ø§Ù„ÙØ±Ø¹ÙŠØ©:", e)
                try:
                    bot.send_message(chat_id, f"ğŸ“ {pretty}\nØ§Ø®ØªØ± ØªØ®ØµØµÙƒ:", reply_markup=keyboard)
                except Exception as e2:
                    print("ÙØ´Ù„ Ø¥Ø±Ø³Ø§Ù„ ÙØ±Ø¹ÙŠØ©:", e2)


        if data.startswith("change_major:"):
            user_states[uid] = "awaiting_custom_major"
            selected = data.split(":", 1)[1]  # Ù…Ø«Ø§Ù„ 'Ø·Ø¨_Ø¨Ø´Ø±ÙŠ'
            save_user_major(uid, selected)     # Ù†ÙØ³ Ø§Ù„Ø¯Ø§Ù„Ø© Ù„Ø­ÙØ¸ Ø§Ù„ØªØ®ØµØµ ÙÙŠ Ù‚Ø§Ø¹Ø¯Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª

            bot.edit_message_text(
                f"âœ… ØªÙ… ØªØºÙŠÙŠØ± ØªØ®ØµØµÙƒ Ø¥Ù„Ù‰: *{selected.replace('_', ' ')}*",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )

            try:
                send_main_menu(chat_id, message_id)  # Ø¥Ø¸Ù‡Ø§Ø± Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©
            except:
                pass

            return




        if data == "advanced":
            with state_lock:
                user_states[int(uid)] = "awaiting_advanced_test_file"
            bot.answer_callback_query(c.id, "âœ… ØªÙ… ØªÙØ¹ÙŠÙ„ ÙˆØ¶Ø¹ advanced  â€” Ø£Ø±Ø³Ù„ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ø¢Ù†.")
            bot.edit_message_text(
                "ğŸ§  ØªÙ… Ø§Ø®ØªÙŠØ§Ø± *Ø§Ù„ÙˆØ¶Ø¹ Ø§Ù„Ù…ØªÙ‚Ø¯Ù…*.\n\n"
                "Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù…Ù„Ù Ø£Ùˆ Ø§Ù„Ù†Øµ Ø§Ù„Ø°ÙŠ ØªØ±ØºØ¨ Ø¨Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù…Ù†Ù‡.",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return

        elif data == "simple":
            with state_lock:
                user_states[uid] = "awaiting_simple_test_file"
            bot.answer_callback_query(c.id, "âœ… ØªÙ… ØªÙØ¹ÙŠÙ„ ÙˆØ¶Ø¹ simple â€” Ø£Ø±Ø³Ù„ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ø¢Ù†.")
            bot.edit_message_text(
            "ğŸ“„ ØªÙ… Ø§Ø®ØªÙŠØ§Ø± *Ø§Ù„ÙˆØ¶Ø¹ Ø§Ù„Ù…Ø¨Ø§Ø´Ø±*.\n\n"
                "Ø£Ø±Ø³Ù„ Ø§Ù„Ù…Ù„Ù Ø£Ùˆ Ø§Ù„Ù†Øµ Ù„Ù„Ø¨Ø¯Ø¡ ÙÙŠ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±.",
                chat_id=chat_id,
                message_id=message_id,
                parse_mode="Markdown"
            )
            logging.info("State set: %s -> %s", uid, user_states.get(int(uid)))
            return



        elif data == "game_private":
            try:
                cursor.execute("SELECT major FROM users WHERE user_id = ?", (uid,))
                row = cursor.fetchone()
                major = row[0] if row else "Ø¹Ø§Ù…"

                keyboard = InlineKeyboardMarkup(row_width=1)
                keyboard.add(
                    InlineKeyboardButton("ğŸ§© Vocabulary Match", callback_data="game_vocab"),
                    InlineKeyboardButton("â±ï¸ ØªØ­Ø¯ÙŠ Ø§Ù„Ø³Ø±Ø¹Ø©", callback_data="game_speed"),
                    InlineKeyboardButton("âŒ Ø§Ù„Ø£Ø®Ø·Ø§Ø¡ Ø§Ù„Ø´Ø§Ø¦Ø¹Ø©", callback_data="game_mistakes"),
                    InlineKeyboardButton("ğŸ§  Ù„Ø¹Ø¨Ø© Ø§Ù„Ø§Ø³ØªÙ†ØªØ§Ø¬", callback_data="game_inference"),
                    InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data="go_games")
                )
                bot.edit_message_text(
                    f"ğŸ“ ØªØ®ØµØµÙƒ Ø§Ù„Ø­Ø§Ù„ÙŠ: {major}\n"
                    "Ø§Ø®ØªØ± Ù„Ø¹Ø¨Ø© ğŸ‘‡",
                    chat_id=chat_id,
                    message_id=message_id,
                    reply_markup=keyboard
                )
            except Exception as e:
                logging.exception("âŒ Ø­Ø¯Ø« Ø®Ø·Ø£ ÙÙŠ game_private")
                bot.send_message(uid, "âŒ Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø¹Ø±Ø¶ Ø§Ù„Ø£Ù„Ø¹Ø§Ø¨.")

    
        elif data == "back_to_games":
            try:
                bot.delete_message(c.message.chat.id, c.message.message_id)
            except Exception as e:
                logging.warning(f"âŒ ÙØ´Ù„ Ø­Ø°Ù Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ø¹Ù†Ø¯ Ø§Ù„Ø±Ø¬ÙˆØ¹: {e}")
    
    

        elif data in ["game_vocab", "game_speed", "game_mistakes", "game_inference"]:
            game_type = data.split("_", 1)[1]

            # Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø¥Ù…ÙƒØ§Ù†ÙŠØ© Ø§Ù„Ù„Ø¹Ø¨ Ø§Ù„ÙŠÙˆÙ…ÙŠ (6 Ù…Ø±Ø§Øª)
            state = game_states.get(uid, {"count": 0})
            if state["count"] >= 6:
                return bot.send_message(uid, "ğŸ›‘ Ù„Ù‚Ø¯ ÙˆØµÙ„Øª Ø¥Ù„Ù‰ Ø§Ù„Ø­Ø¯ Ø§Ù„Ø£Ù‚ØµÙ‰ Ù„Ù„Ø£Ù„Ø¹Ø§Ø¨ Ø§Ù„Ù…Ø¬Ø§Ù†ÙŠØ© (6 Ù…Ø±Ø§Øª).")

            if not can_play_game_today(uid, game_type):
                bot.answer_callback_query(c.id, "âŒ Ù„Ù‚Ø¯ Ù„Ø¹Ø¨Øª Ù‡Ø°Ù‡ Ø§Ù„Ù„Ø¹Ø¨Ø© Ø§Ù„ÙŠÙˆÙ…!")
                return

            loading_msg = bot.send_message(chat_id, "â³ Ø¬Ø§Ø±ÙŠ ØªØ­Ø¶ÙŠØ± Ø§Ù„Ø³Ø¤Ø§Ù„...")

            try:
                record_game_attempt(uid, game_type)

                # Ø§Ù„ØªØ®ØµØµ
                cursor.execute("SELECT major FROM users WHERE user_id=?", (uid,))
                row = cursor.fetchone()
                major = row[0] if row else "Ø¹Ø§Ù…"

                # ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø­Ø³Ø¨ Ù†ÙˆØ¹ Ø§Ù„Ù„Ø¹Ø¨Ø©
                if game_type == "vocab":
                    raw = generate_vocabulary_game(uid, major, native_lang="Arabic")
                elif game_type == "speed":
                    raw = generate_speed_challenge(uid, major, native_lang="Arabic")
                elif game_type == "mistakes":
                    raw = generate_common_mistakes_game(uid, major, native_lang="Arabic")
                elif game_type == "inference":
                    raw = generate_inference_game(uid, major, native_lang="Arabic")

                question = raw["question"]
                options = raw["options"]
                correct_index = raw["correct_index"]

                if not isinstance(options, list) or len(options) < 2:
                    raise ValueError("Ø¹Ø¯Ø¯ Ø§Ù„Ø®ÙŠØ§Ø±Ø§Øª ØºÙŠØ± ØµØ§Ù„Ø­")

                # Ø­ÙØ¸ Ø®ÙŠØ§Ø±Ø§Øª Ø§Ù„Ø³Ø¤Ø§Ù„ ÙÙŠ Ø§Ù„Ø°Ø§ÙƒØ±Ø© Ø§Ù„Ù…Ø¤Ù‚ØªØ©
                game_states[uid] = {"count": state["count"] + 1, "options": options}

                keyboard = InlineKeyboardMarkup(row_width=2)

                # Ø£Ø²Ø±Ø§Ø± Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø§Øª
                for i, option in enumerate(options):
                    short_option = (option[:50] + "...") if len(option) > 50 else option
                    callback_data = f"ans_{game_type}_{i}_{correct_index}"
                    keyboard.add(InlineKeyboardButton(short_option, callback_data=callback_data))
    
                # Ø£Ø²Ø±Ø§Ø± Ø§Ù„ØªØ­ÙƒÙ…
                keyboard.row(
                    InlineKeyboardButton("ğŸ”„ Ø³Ø¤Ø§Ù„ Ø¬Ø¯ÙŠØ¯", callback_data=f"new_{game_type}"),
                    InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data="back_to_games")
                )
                keyboard.add(
                    InlineKeyboardButton(
                        "ğŸ“¤ Ø´Ø§Ø±Ùƒ Ù‡Ø°Ù‡ Ø§Ù„Ù„Ø¹Ø¨Ø©", 
                        switch_inline_query="Ø¬Ø±Ø¨ Ù‡Ø°Ù‡ Ø§Ù„Ù„Ø¹Ø¨Ø© Ø§Ù„Ø±Ø§Ø¦Ø¹Ø© Ù…Ù† @Oiuhelper_bot ğŸ¯")
                )

                bot.delete_message(chat_id, loading_msg.message_id)
                text = f"ğŸ§  Ø§Ø®ØªØ± Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø©:\n\n{question}"
                bot.send_message(chat_id, text, reply_markup=keyboard)

            except Exception as e:
                try:
                    bot.delete_message(chat_id, loading_msg.message_id)
                except:
                    pass
                logging.error(f"ÙØ´Ù„ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ù„Ø¹Ø¨Ø©: {str(e)}")
                bot.send_message(uid, "âŒ Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ù„Ø¹Ø¨Ø©ØŒ Ø­Ø§ÙˆÙ„ Ù„Ø§Ø­Ù‚Ø§Ù‹")

        # Ù…Ø¹Ø§Ù„Ø¬Ø© Ø·Ù„Ø¨ Ø³Ø¤Ø§Ù„ Ø¬Ø¯ÙŠØ¯
    


        elif data.startswith("new_"):
            game_type = data.split("_", 1)[1]

            # ØªØ­Ù‚Ù‚ Ù…Ù† Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø§Øª (ÙƒÙ…Ø§ ÙÙŠ Ø§Ù„Ù‚Ø³Ù… Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠ)
            state = game_states.get(uid, {"count": 0})
            if state["count"] >= 6:
                msg = random.choice([
                    "ğŸš« ÙˆØµÙ„Øª Ø¥Ù„Ù‰ Ø§Ù„Ø­Ø¯ Ø§Ù„Ø£Ù‚ØµÙ‰ Ù„Ø¹Ø¯Ø¯ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ø§Ù„ÙŠÙˆÙ…!\nâœ¨ Ø¬Ø±Ø¨ ØºØ¯Ù‹Ø§ Ø£Ùˆ Ø´Ø§Ø±Ùƒ Ø§Ù„Ø¨ÙˆØª Ù…Ø¹ Ø£ØµØ¯Ù‚Ø§Ø¦Ùƒ!",
                    "âŒ Ø§Ù†ØªÙ‡Øª Ù…Ø­Ø§ÙˆÙ„Ø§Øª Ø§Ù„ÙŠÙˆÙ…! ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù…Ø¬Ø¯Ø¯Ù‹Ø§ Ù„Ø§Ø­Ù‚Ù‹Ø§.",
                    "ğŸ›‘ Ù„Ø§ Ù…Ø²ÙŠØ¯ Ù…Ù† Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ø§Ù„Ø¢Ù†. Ø¹Ø¯ Ù„Ø§Ø­Ù‚Ù‹Ø§ Ù„ØªÙƒÙ…Ù„ Ø±Ø­Ù„ØªÙƒ!"
                ])
                return bot.answer_callback_query(c.id, msg, show_alert=True)

            loading_msg = bot.send_message(c.message.chat.id, "â³ Ø¬Ø§Ø±ÙŠ ØªØ­Ø¶ÙŠØ± Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„ØªØ§Ù„ÙŠ...")

            try:
                # ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
                cursor.execute("SELECT major FROM users WHERE user_id=?", (uid,))
                row = cursor.fetchone()
                major = row[0] if row else "Ø¹Ø§Ù…"

                game_generators = {
                    "vocab": generate_vocabulary_game,
                    "speed": generate_speed_challenge,
                    "mistakes": generate_common_mistakes_game,
                    "inference": generate_inference_game
                }

                raw = game_generators[game_type](uid, major)
                question = raw["question"]
                options = raw["options"]
                correct_index = raw["correct_index"]

                if not isinstance(options, list) or len(options) < 2:
                    raise ValueError("Ø¹Ø¯Ø¯ Ø§Ù„Ø®ÙŠØ§Ø±Ø§Øª ØºÙŠØ± ØµØ§Ù„Ø­")

                # Ø­ÙØ¸ Ø®ÙŠØ§Ø±Ø§Øª Ø§Ù„Ø³Ø¤Ø§Ù„ Ø§Ù„Ø¬Ø¯ÙŠØ¯
                game_states[uid]["count"] += 1
                game_states[uid]["options"] = options

                # Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø£Ø²Ø±Ø§Ø±
                keyboard = InlineKeyboardMarkup(row_width=2)
                for i, option in enumerate(options):
                    short_option = (option[:50] + "...") if len(option) > 50 else option
                    callback_data = f"ans_{game_type}_{i}_{correct_index}"
                    keyboard.add(InlineKeyboardButton(short_option, callback_data=callback_data))

                keyboard.row(
                    InlineKeyboardButton("ğŸ”„ Ø³Ø¤Ø§Ù„ Ø¬Ø¯ÙŠØ¯", callback_data=f"new_{game_type}"),
                    InlineKeyboardButton("â¬…ï¸ Ø±Ø¬ÙˆØ¹", callback_data="back_to_games")
                )
                keyboard.add(
                    InlineKeyboardButton(
                        "ğŸ“¤ Ø´Ø§Ø±Ùƒ Ù‡Ø°Ù‡ Ø§Ù„Ù„Ø¹Ø¨Ø©", 
                        switch_inline_query="Ø¬Ø±Ø¨ Ù‡Ø°Ù‡ Ø§Ù„Ù„Ø¹Ø¨Ø© Ø§Ù„Ø±Ø§Ø¦Ø¹Ø© Ù…Ù† @Oiuhelper_bot ğŸ¯")
                )

                # ØªØ¹Ø¯ÙŠÙ„ Ù†ÙØ³ Ø§Ù„Ø±Ø³Ø§Ù„Ø©
                bot.edit_message_text(
                    text=f"ğŸ§  Ø§Ø®ØªØ± Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø©:\n\n{question}",
                    chat_id=c.message.chat.id,
                    message_id=c.message.message_id,
                    reply_markup=keyboard
                )

            except Exception as e:
                logging.error(f"âŒ ÙØ´Ù„ ØªÙˆÙ„ÙŠØ¯ Ø³Ø¤Ø§Ù„ Ø¬Ø¯ÙŠØ¯: {e}")
                bot.answer_callback_query(c.id, "âŒ ÙØ´Ù„ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø³Ø¤Ø§Ù„")

            finally:
                try:
                    bot.delete_message(c.message.chat.id, loading_msg.message_id)
                except:
                    pass

        elif data.startswith("ans_"):
            parts = data.split("_")
            game_type = parts[1]
            selected = int(parts[2])
            correct = int(parts[3])

            options = game_states.get(uid, {}).get("options", [])
            correct_text = options[correct] if correct < len(options) else f"Ø§Ù„Ø®ÙŠØ§Ø± Ø±Ù‚Ù… {correct+1}"

            wrong_responses = [
                "âŒ Ø®Ø·Ø£! Ø¬Ø±Ø¨ Ù…Ø¬Ø¯Ø¯Ù‹Ø§ ğŸ˜‰\nâœ… Ø§Ù„ØµØ­ÙŠØ­: {correct}",
                "ğŸš« Ù„Ù„Ø£Ø³ÙØŒ Ù„ÙŠØ³Øª Ø§Ù„ØµØ­ÙŠØ­Ø©!\nâœ… Ø§Ù„Ø¬ÙˆØ§Ø¨: {correct}",
                "ğŸ˜… Ù„ÙŠØ³Øª Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø©ØŒ Ø§Ù„Ø¬ÙˆØ§Ø¨ Ù‡Ùˆ: {correct}",
                "âŒ Ù„Ø§ØŒ Ø­Ø§ÙˆÙ„ Ù…Ø±Ø© Ø£Ø®Ø±Ù‰!\nâœ”ï¸ Ø§Ù„ØµØ­ÙŠØ­ Ù‡Ùˆ: {correct}"
            ]

            if selected == correct:
                bot.answer_callback_query(c.id, "âœ… Ø¥Ø¬Ø§Ø¨Ø© ØµØ­ÙŠØ­Ø©!", show_alert=False)
            else:
                msg = random.choice(wrong_responses).format(correct=correct_text)
                bot.answer_callback_query(c.id, msg, show_alert=False)


        
        # ÙŠÙ…ÙƒÙ†Ùƒ Ø¥Ø¶Ø§ÙØ© Ø§Ù„Ù…Ø²ÙŠØ¯ Ù…Ù† Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø§Øª Ø§Ù„Ø£Ø®Ø±Ù‰ Ù„Ù€ callback_data Ù‡Ù†Ø§


        elif data.startswith("soon_"):
            feature_name = {
                "soon_review": "ğŸ“š Ù…ÙŠØ²Ø© Ø§Ù„Ù…Ø±Ø§Ø¬Ø¹Ø© Ø§Ù„Ø³Ø±ÙŠØ¹Ø©",
                "soon_summary": "ğŸ“„ Ù…Ù„Ø®ØµØ§Øª PDF",
            }.get(data, "Ù‡Ø°Ù‡ Ø§Ù„Ù…ÙŠØ²Ø©")

            bot.answer_callback_query(c.id)
            bot.send_message(chat_id, f"{feature_name} Ø³ØªÙƒÙˆÙ† Ù…ØªØ§Ø­Ø© Ù‚Ø±ÙŠØ¨Ù‹Ø§... ğŸš§")
        
        elif data.startswith("retry:"):
            quiz_code = data[6:]
            quiz_manager.start_quiz(chat_id, quiz_code, bot)
        
        elif data.startswith("share_quiz:"):
            quiz_code = data.split(":", 1)[1]
            chat_id = c.message.chat.id  # â† ØªØ£ÙƒØ¯ Ù…Ù† ØªØ¹ÙŠÙŠÙ† chat_id Ù‡Ù†Ø§

            share_link = f"https://t.me/Oiuhelper_bot?start=quiz_{quiz_code}"

            msg_text = f"""<b>ğŸ‰ Ø´Ø§Ø±Ùƒ Ù‡Ø°Ø§ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù…Ø¹ Ø²Ù…Ù„Ø§Ø¦Ùƒ!</b>

        Ø§Ù†Ø³Ø® Ø§Ù„Ø±Ø§Ø¨Ø· Ø£Ø¯Ù†Ø§Ù‡ Ø£Ùˆ Ø§Ø¶ØºØ· Ù„ÙØªØ­Ù‡ Ù…Ø¨Ø§Ø´Ø±Ø©:
        ğŸ”— <a href="{share_link}">{share_link}</a>

        ğŸ“ Ø¹Ù†Ø¯ ÙØªØ­ Ø§Ù„Ø±Ø§Ø¨Ø·ØŒ Ø³ÙŠØ¨Ø¯Ø£ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§ Ø¨Ø¥Ø°Ù† Ø§Ù„Ù„Ù‡.
        """

            keyboard = types.InlineKeyboardMarkup()
            keyboard.add(
                types.InlineKeyboardButton("ğŸ”— Ù†Ø³Ø® Ø§Ù„Ø±Ø§Ø¨Ø·", switch_inline_query=share_link),
                types.InlineKeyboardButton("ğŸ  Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©", callback_data="go_back_home")
            )

            bot.send_message(chat_id, msg_text, parse_mode="HTML", reply_markup=keyboard)

    except Exception as e:
        logging.exception("Callback handler error: %s", e)



# -------------------------------------------------------------------
# ------ message & document Handlers ---------------------------------
# ---------------------------------------


@bot.message_handler(func=lambda m: user_states.get(m.from_user.id) == "awaiting_major", content_types=['text'])
def set_custom_major(msg):
    try:
        major = msg.text.strip()
        uid = msg.from_user.id

        # Ø­ÙØ¸ Ø§Ù„ØªØ®ØµØµ ÙÙŠ DB
        save_user_major(uid, major)
    
        # Ø¥Ø®Ø·Ø§Ø± Ø§Ù„Ø£Ø¯Ù…Ù†
        bot.send_message(
            ADMIN_ID,
            f"ğŸ†• ØªØ®ØµØµ Ø¬Ø¯ÙŠØ¯ Ø£ÙØ±Ø³Ù„ Ù…Ù† Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù…:\n"
            f"ğŸ‘¤ @{msg.from_user.username or msg.from_user.id}\n"
            f"ğŸ“š Ø§Ù„ØªØ®ØµØµ: {major}"
        )
        
        # ØªØºÙŠÙŠØ± Ø§Ù„Ø­Ø§Ù„Ø© Ù„Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù„Ø§Ø­Ù‚Ø§Ù‹
        user_states[uid] = "awaiting_simple_test_file"
    
        # Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© ØªØ£ÙƒÙŠØ¯
        bot.send_message(
            uid,
            f"âœ… ØªÙ… ØªØ­Ø¯ÙŠØ¯ ØªØ®ØµØµÙƒ: {major}\n"
            f"Ø§Ù„Ø¢Ù† Ø£Ø±Ø³Ù„ Ù…Ù„Ù (PDF/DOCX/TXT) Ø£Ùˆ Ù†ØµÙ‹Ø§ Ù…Ø¨Ø§Ø´Ø±Ù‹Ø§ Ù„ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±Ùƒ."
        )
    except Exception as e:
        logging.error(f"Error in set_custom_major: {e}")
        bot.send_message(uid, "âŒ Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ù…Ø¹Ø§Ù„Ø¬Ø© Ø·Ù„Ø¨ÙƒØŒ ÙŠØ±Ø¬Ù‰ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù„Ø§Ø­Ù‚Ø§Ù‹")
@bot.message_handler(func=lambda m: user_states.get(m.from_user.id) in [
    "awaiting_major_for_games",
    "awaiting_custom_major"
])
def handle_user_major(msg):
    if msg.chat.type != "private":
        return  # ØªØ¬Ø§Ù‡Ù„ Ø§Ù„Ø±Ø³Ø§Ø¦Ù„ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø§Øª

    uid = msg.from_user.id
    state = user_states.get(uid)
    major = msg.text.strip()
    if uid in user_states and 'message_id' in user_states[uid]:
        stored_message_id = user_states[uid]['message_id']

    if len(major) < 2:
        bot.send_message(uid, "âš ï¸ ÙŠØ±Ø¬Ù‰ Ø¥Ø¯Ø®Ø§Ù„ ØªØ®ØµØµ ØµØ§Ù„Ø­.")
        return

    cursor.execute("INSERT OR REPLACE INTO users(user_id, major) VALUES(?, ?)", (uid, major))
    conn.commit()
    user_states.pop(uid, None)


    if state == "awaiting_major_for_games":
        bot.send_message(uid, f"âœ… ØªÙ… ØªØ³Ø¬ÙŠÙ„ ØªØ®ØµØµÙƒ: {major}\n"
                              "Ø§Ù„Ø¢Ù† ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ø®ØªÙŠØ§Ø± Ù„Ø¹Ø¨Ø© Ù…Ù† Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø£Ù„Ø¹Ø§Ø¨ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠØ©.")
        keyboard = InlineKeyboardMarkup(row_width=1)
        keyboard.add(
            InlineKeyboardButton("ğŸ”’ Ø§Ù„Ø¹Ø¨ ÙÙŠ Ø§Ù„Ø®Ø§Øµ", callback_data="game_private"),
            InlineKeyboardButton("ğŸ‘¥ Ø§Ù„Ø¹Ø¨ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø©", switch_inline_query="game")
        )
        bot.send_message(uid, "ğŸ® Ø§Ø®ØªØ± Ø·Ø±ÙŠÙ‚Ø© Ø§Ù„Ù„Ø¹Ø¨:", reply_markup=keyboard)

    elif state == "awaiting_custom_major":
        sent = bot.send_message(uid, f"âœ… ØªÙ… ØªØ³Ø¬ÙŠÙ„ ØªØ®ØµØµÙƒ: *{major}*", parse_mode="Markdown")
        time.sleep(2)
        try:
            bot.edit_message_text(
                "â¬‡ï¸ Ù‡Ø°Ù‡ Ù‡ÙŠ Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ø§Ù„Ø±Ø¦ÙŠØ³ÙŠØ©:",
                chat_id=sent.chat.id,
                message_id=sent.message_id
            )
            send_main_menu(uid, message_id=sent.message_id)
        except:
            send_main_menu(uid)



# Ù…Ø¬Ø±Ø¯ Ø¥Ù‚Ø±Ø§Ø± Ø³Ø±ÙŠØ¹ ÙÙŠ Ø§Ù„Ù€ handler Ø«Ù… ÙˆØ¶Ø¹ ÙÙŠ Ø§Ù„Ø·Ø§Ø¨ÙˆØ±

@bot.message_handler(content_types=['text', 'document', 'photo'])
def unified_handler(msg):
    if msg.chat.type != "private":
        return
    
    uid = int(msg.from_user.id)

    file_id = getattr(getattr(msg, 'document', None), 'file_id', None)
    if is_request_already_queued(file_id=file_id, user_id=uid, message_id=msg.message_id):
        # Ø³Ø¬Ù„/Ø£Ø±Ø³Ù„ Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù… Ø±Ø³Ø§Ù„Ø© Ù‚ØµÙŠØ±Ø© ØªÙÙŠØ¯ Ø£Ù† Ø§Ù„Ø·Ù„Ø¨ Ù‚ÙŠØ¯ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ø¨Ø§Ù„ÙØ¹Ù„
        bot.reply_to(msg, "â³ Ø·Ù„Ø¨Ùƒ Ù‚ÙŠØ¯ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø¨Ø§Ù„ÙØ¹Ù„.")
        return



    # Ø­ÙØ¸ Ø§Ù„Ø·Ù„Ø¨ Ø¥Ù† Ø£Ø±Ø¯Øª
    try:
        save_request(msg)
    except Exception:
        pass
    
    update_files_and_users(uid, files_count=1)
    update_daily_stats(files=1)


    sent_msg = bot.reply_to(msg, "ğŸ“ Ø¬Ø§Ø±ÙŠ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ØŒ ÙŠØ±Ø¬Ù‰ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù‚Ù„ÙŠÙ„Ø§Ù‹...")
    try:

        request_queue.put_nowait((msg, sent_msg))
    except queue.Full:
        bot.edit_message_text(" Ø§Ù„Ø¹Ø¯ÙŠØ¯ Ù…Ù† Ø§Ù„Ø£Ø´Ø®Ø§Øµ ÙŠÙ‚ÙˆÙ…ÙˆÙ† Ø¨ØªÙˆÙ„ÙŠØ¯ Ù…Ø­ØªÙˆÙ‰ ÙÙŠ Ø§Ù„ÙˆÙ‚Øª Ø§Ù„Ø­Ø§Ù„ÙŠØŒ ÙŠØ±Ø¬Ù‰ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù‚Ù„ÙŠÙ„Ø§Ù‹ ï¸ğŸ•°ï¸. Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø¥Ù†ØªØ¸Ø§Ø± ...", chat_id=sent_msg.chat.id, message_id=sent_msg.message_id)
        time.sleep(random.randint(1, 2))
        while True:
            try:
                request_queue.put((msg, sent_msg), timeout=5)
            
                break
            except Exception as e:
                bot.edit_message_text("Ø­Ø¯Ø« Ø®Ø·Ø£ØŒ ÙŠØ±Ø¬Ù‰ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù‚Ù„ÙŠÙ„Ø§Ù‹...", chat_id=sent_msg.chat.id, message_id=sent_msg.message_id)
                time.sleep(1)


# -------------------------------------------------------------------
# Ø§Ù„Ø¯Ø§Ù„Ø© Ø§Ù„ØªÙŠ ØªÙ†ÙØ° ÙØ¹Ù„ÙŠÙ‹Ø§ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© (ØªØ¹Ù…Ù„ Ø¯Ø§Ø®Ù„ Ø§Ù„Ø¹Ø§Ù…Ù„)
# -------------------------------------------------------------------
def process_message(msg, message_id=None, chat_id=None):
    logging.info("process_message enter: uid=%s type=%s", msg.from_user.id, msg.content_type)

    # ...

    content_type = msg.content_type
    username = msg.from_user.username or "Ø¨Ø¯ÙˆÙ† Ø§Ø³Ù… Ù…Ø³ØªØ®Ø¯Ù…"
    uid = msg.from_user.id
        
    with semaphore:
        if msg.content_type == "document":
            file_id = msg.document.file_id
        else:
            file_id = None  # Ù„Ø§ ØªØ­Ø¯Ù‘Ø« requests Ù‡Ù†Ø§ Ù„Ø£Ù† Ù„Ø§ Ù…Ù„Ù


        try:
            if file_id:
                update_request_status(file_id, 'processing')
                print(f"ğŸš€ Ø¨Ø¯Ø¡ Ù…Ø¹Ø§Ù„Ø¬Ø©: {file_id}")
            time.sleep(random.randint(1, 2))  # Ù…Ø­Ø§ÙƒØ§Ø© Ù…Ø¹Ø§Ù„Ø¬Ø©
            
            if file_id:
                update_request_status(file_id, 'done')
                print(f"âœ… ØªÙ… Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø©: {file_id}")
        except Exception:
            logging.exception("requests status update failed")

    notify_process_info(uid, file_id, username)




    # Ø§Ø¬Ù„Ø¨ Ø§Ù„Ø­Ø§Ù„Ø©
    try:
        with state_lock:
            state = user_states.get(int(uid))
    except NameError:
        # Ø¥Ø°Ø§ Ù…Ø§ Ø¹Ù†Ø¯Ùƒ state_lockØŒ Ø§Ø³ØªØ®Ø¯Ù… Ø¨Ø¯ÙˆÙ† Ù‚ÙÙ„
        state = user_states.get(int(uid))

    logging.info("state for uid %s: %s", uid, state)

    major = fetch_user_major(uid)
    

    content = ""
    path = None

    # log at entry
    logging.info("Message received: uid=%s type=%s text_len=%s", uid, content_type, len(getattr(msg, 'text', '') or ""))


    
        # ÙÙˆØ± Ø§Ù„ØªØ£ÙƒØ¯ Ù…Ù† ÙˆØ¬ÙˆØ¯ Ø­Ø§Ù„Ø©ØŒ Ø§Ø¹Ø±Ø¶ Ù„Ù„Ù‘ÙˆØ¬
    if state is None:
        # ÙŠÙ…ÙƒÙ† Ø§Ù„Ø±Ø¯ Ù„Ù„ØªØ¬Ø±Ø¨Ø© ÙÙ‚Ø· Ø£Ø«Ù†Ø§Ø¡ Ø§Ù„Ø¯ÙŠØ¨Ø§Øº
        # bot.send_message(uid, f"DEBUG: no state set (you are {uid})")
        return


    try:
        if content_type == "text":
            content = msg.text or ""
            coverage = "ÙƒØ§Ù…Ù„Ø© âœ…"

        
        # Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„ØµÙˆØ± (photo)
        elif msg.content_type == "photo":
            if not can_generate(uid):
                return bot.reply_to(msg, "âš ï¸ Ù‡Ø°Ù‡ Ø§Ù„Ù…ÙŠØ²Ø© Ù…ØªØ§Ø­Ø© ÙÙ‚Ø· Ù„Ù„Ù…Ø´ØªØ±ÙƒÙŠÙ†.")
            
            file_id = msg.photo[-1].file_id
            file_info = bot.get_file(file_id)
            file_data = bot.download_file(file_info.file_path)

            os.makedirs("downloads", exist_ok=True)
            path = os.path.join("downloads", f"{uid}_photo.jpg")
            with open(path, "wb") as f:
                f.write(file_data)

            bot.edit_message_text("ğŸ–¼ï¸ Ø¬Ø§Ø±ÙŠ Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ù†Øµ Ù…Ù† Ø§Ù„ØµÙˆØ±Ø©...", chat_id=chat_id, message_id=message_id)


            content, ocr_debug = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng")
            if not content.strip():
                return bot.send_message(uid, f"âŒ ÙØ´Ù„ ÙÙŠ Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ù†Øµ Ù…Ù† Ø§Ù„ØµÙˆØ±Ø©. {ocr_debug}")


        elif msg.content_type == "document":
            # ÙˆØ¶Ø¹ Ø§Ù„Ø·Ù„Ø¨ ÙÙŠ Ø§Ù„Ù‚Ø§Ø¦Ù…Ø© Ù„ÙŠØªÙ… Ù…Ø¹Ø§Ù„Ø¬ØªÙ‡ Ø¨ÙˆØ§Ø³Ø·Ø© Ø§Ù„Ø¹Ù…Ø§Ù„
            file_info = bot.get_file(msg.document.file_id)
            if file_info.file_size > 5 * 1024 * 1024:
                return bot.send_message(uid, "âŒ Ø§Ù„Ù…Ù„Ù ÙƒØ¨ÙŠØ± Ø¬Ø¯Ù‹Ø§ØŒ Ø§Ù„Ø­Ø¯ 5 Ù…ÙŠØºØ§Ø¨Ø§ÙŠØª.")
    
            file_data = bot.download_file(file_info.file_path)
            os.makedirs("downloads", exist_ok=True)
            path = os.path.join("downloads", msg.document.file_name)

            with open(path, "wb") as f:
                f.write(file_data)
            user_files[uid] = path

            ext = path.rsplit(".", 1)[-1].lower()
            # Ø¨Ø¹Ø¯ Ø§Ù„Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ø¹Ø§Ø¯ÙŠ
            if ext == "pdf":
                content_full = extract_text_from_pdf(path)  # Ø§Ù„Ù†Øµ Ø§Ù„ÙƒØ§Ù…Ù„
                full_length = len(content_full)

                # Ø¥Ø°Ø§ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… ØºÙŠØ± Ù…Ø´ØªØ±ÙƒØŒ Ø§Ù‚ØªØ·Ø¹ ÙÙ‚Ø· 3000 Ø­Ø±Ù
                if not can_generate(uid):
                    content = content_full[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% Ù…Ù† Ø§Ù„Ù…Ù„Ù"
                else:
                    content = content_full
                    coverage = "ÙƒØ§Ù…Ù„Ø© âœ…"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(
                            uid,
                            "âš ï¸ Ù„Ø§ ÙŠÙ…ÙƒÙ† Ù‚Ø±Ø§Ø¡Ø© Ù‡Ø°Ø§ Ø§Ù„Ù…Ù„Ù ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§. ØªØªØ·Ù„Ø¨ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…ØªÙ‚Ø¯Ù…Ø© Ø§Ø´ØªØ±Ø§ÙƒÙ‹Ø§ ÙØ¹Ø§Ù„Ù‹Ø§."
                        )
                    bot.reply_to(msg, "â³ ÙŠØªÙ… ØªØ¬Ù‡ÙŠØ² Ø§Ù„Ù…Ù„Ù... Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù„Ø­Ø¸Ø§Øª.")
                    language = detect_language_from_filename(msg.document.file_name)
                    content, ocr_debug = extract_text_from_pdf_with_ocr(path, api_key=OCR_API_KEY, language=language)
                    if not content.strip():
                        bot.send_message(uid, f"âŒ ÙØ´Ù„ ÙÙŠ Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ù†Øµ Ù…Ù† Ø§Ù„Ù…Ù„Ù. {ocr_debug}")
                        return
                    preview = content[:1500]
                    bot.send_message(uid, f"ğŸ“„ ØªÙ… Ø§Ø³ØªØ®Ø±Ø§Ø¬ Ø§Ù„Ù†Øµ Ø¨Ù†Ø¬Ø§Ø­ (Ø¬Ø²Ø¡ Ù…Ù†Ù‡):\n\n{preview}")
            elif ext == "docx":
                content_full = extract_text_from_pdf(path)  # Ø§Ù„Ù†Øµ Ø§Ù„ÙƒØ§Ù…Ù„
                full_length = len(content_full)
                # Ø¥Ø°Ø§ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… ØºÙŠØ± Ù…Ø´ØªØ±ÙƒØŒ Ø§Ù‚ØªØ·Ø¹ ÙÙ‚Ø· 3000 Ø­Ø±Ù
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% Ù…Ù† Ø§Ù„Ù…Ù„Ù"
                else:
                    content = content_full
                    coverage = "ÙƒØ§Ù…Ù„Ø© âœ…"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "âš ï¸ Ù„Ø§ ÙŠÙ…ÙƒÙ† Ù‚Ø±Ø§Ø¡Ø© Ù‡Ø°Ø§ Ø§Ù„Ù…Ù„Ù ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§. ØªØªØ·Ù„Ø¨ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…ØªÙ‚Ø¯Ù…Ø© Ø§Ø´ØªØ±Ø§ÙƒÙ‹Ø§ ÙØ¹Ø§Ù„Ù‹Ø§.")
                    bot.edit_message_text("â³ ÙŠØªÙ… ØªØ¬Ù‡ÙŠØ² Ø§Ù„Ù…Ù„Ù... Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù„Ø­Ø¸Ø§Øª.", chat_id=chat_id, message_id=message_id)
                    language = detect_language_from_filename(msg.document.file_name)
                    content = extract_text_from_pdf_with_ocr(path, api_key=OCR_API_KEY, language=language)
            elif ext == "txt":
                content_full = extract_text_from_pdf(path)  # Ø§Ù„Ù†Øµ Ø§Ù„ÙƒØ§Ù…Ù„
                full_length = len(content_full)
                # Ø¥Ø°Ø§ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… ØºÙŠØ± Ù…Ø´ØªØ±ÙƒØŒ Ø§Ù‚ØªØ·Ø¹ ÙÙ‚Ø· 3000 Ø­Ø±Ù
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% Ù…Ù† Ø§Ù„Ù…Ù„Ù"
                else:
                    content = content_full
                    coverage = "ÙƒØ§Ù…Ù„Ø© âœ…"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "âš ï¸ Ù„Ø§ ÙŠÙ…ÙƒÙ† Ù‚Ø±Ø§Ø¡Ø© Ù‡Ø°Ø§ Ø§Ù„Ù…Ù„Ù ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§. ØªØªØ·Ù„Ø¨ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…ØªÙ‚Ø¯Ù…Ø© Ø§Ø´ØªØ±Ø§ÙƒÙ‹Ø§ ÙØ¹Ø§Ù„Ù‹Ø§.")
                    bot.edit_message_text("â³ ÙŠØªÙ… ØªØ¬Ù‡ÙŠØ² Ø§Ù„Ù…Ù„Ù... Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù„Ø­Ø¸Ø§Øª.", chat_id=chat_id, message_id=message_id)
                    content = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng+ara")
                
            elif ext == "pptx":
                content_full = extract_text_from_pdf(path)  # Ø§Ù„Ù†Øµ Ø§Ù„ÙƒØ§Ù…Ù„
                full_length = len(content_full)
                
                # Ø¥Ø°Ø§ Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… ØºÙŠØ± Ù…Ø´ØªØ±ÙƒØŒ Ø§Ù‚ØªØ·Ø¹ ÙÙ‚Ø· 3000 Ø­Ø±Ù
                if not can_generate(uid):
                    content = content[:3000]
                    coverage_ratio = (len(content) / full_length) * 100 if full_length > 0 else 0
                    coverage = f"{coverage_ratio:.1f}% Ù…Ù† Ø§Ù„Ù…Ù„Ù"
                else:
                    content = content_full
                    coverage = "ÙƒØ§Ù…Ù„Ø© âœ…"

                if is_text_empty(content):
                    if not can_generate(uid):
                        return bot.send_message(uid, "âš ï¸ Ù„Ø§ ÙŠÙ…ÙƒÙ† Ù‚Ø±Ø§Ø¡Ø© Ù‡Ø°Ø§ Ø§Ù„Ù…Ù„Ù ØªÙ„Ù‚Ø§Ø¦ÙŠÙ‹Ø§. ØªØªØ·Ù„Ø¨ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…ØªÙ‚Ø¯Ù…Ø© Ø§Ø´ØªØ±Ø§ÙƒÙ‹Ø§ ÙØ¹Ø§Ù„Ù‹Ø§.")
                    bot.edit_message_text("â³ ÙŠØªÙ… ØªØ¬Ù‡ÙŠØ² Ø§Ù„Ù…Ù„Ù... Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù„Ø­Ø¸Ø§Øª.", chat_id=chat_id, message_id=message_id)
                    language = detect_language_from_filename(msg.document.file_name)
                    content = extract_text_from_pptx_with_ocr(path, api_key=OCR_API_KEY, language=language)

            elif ext in ("jpg", "png"):
                if not can_generate(uid):
                    return bot.send_message(uid, "âš ï¸ Ù‡Ø°Ù‡ Ø§Ù„Ù…ÙŠØ²Ø© Ù…ØªØ§Ø­Ø© ÙÙ‚Ø· Ù„Ù„Ù…Ø´ØªØ±ÙƒÙŠÙ†.")
                bot.edit_message_text("â³ Ø¬Ø§Ø±ÙŠ ØªØ­Ù„ÙŠÙ„ Ø§Ù„ØµÙˆØ±Ø©...", chat_id=chat_id, message_id=message_id)
                content, ocr_debug = extract_text_with_ocr_space(path, api_key=OCR_API_KEY, language="eng")
            
               
            else:
                return bot.send_message(uid, "âš ï¸ Ù†ÙˆØ¹ Ø§Ù„Ù…Ù„Ù ØºÙŠØ± Ù…Ø¯Ø¹ÙˆÙ…. Ø£Ø±Ø³Ù„ PDF Ø£Ùˆ Word Ø£Ùˆ TXT.")

        else:
            try:
                os.remove(path)
            except Exception as e:
                print(f"[WARNING] Ù„Ù… ÙŠØªÙ… Ø­Ø°Ù Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…Ø¤Ù‚Øª: {e}")
        if not state:
        # Ù„Ø§ Ø­Ø§Ù„Ø©: Ù„Ø§ ØªÙØ¹Ù„ Ø´ÙŠØ¡
            return

        if not content or not content.strip():
            return bot.send_message(uid, "âš ï¸ Ù„Ù… Ø£ØªÙ…ÙƒÙ† Ù…Ù† Ù‚Ø±Ø§Ø¡Ø© Ù…Ø­ØªÙˆÙ‰ Ø§Ù„Ù…Ù„Ù Ø£Ùˆ Ø§Ù„Ù†Øµ.")
        print(f">>> Content preview: {content[:300]}")

        waiting_messages_anki = [
            "ğŸ§  ÙŠØªÙ… ØªØ­Ù„ÙŠÙ„ Ø§Ù„Ù†Øµ Ù„Ø§Ø³ØªØ®Ù„Ø§Øµ Ø§Ù„Ù†Ù‚Ø§Ø· Ø§Ù„Ø£Ø³Ø§Ø³ÙŠØ©...",
            "âœ¨ Ø¬Ø§Ø±ÙŠ ØªØ­ÙˆÙŠÙ„ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø¥Ù„Ù‰ Ø¨Ø·Ø§Ù‚Ø§Øª Ø°ÙƒÙŠØ© Ø³Ù‡Ù„Ø© Ø§Ù„Ù…Ø±Ø§Ø¬Ø¹Ø©...",
            "ğŸ“š ØªÙ†Ø¸ÙŠÙ… Ø§Ù„Ù…Ø¹Ù„ÙˆÙ…Ø§Øª Ù„ØªØ¹Ø²ÙŠØ² Ù‚Ø¯Ø±ØªÙƒ Ø¹Ù„Ù‰ Ø§Ù„ØªØ°ÙƒØ±...",
            "ğŸ¨ ØªØµÙ…ÙŠÙ… Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª Ø¨Ø£Ø³Ù„ÙˆØ¨ ÙŠØ³Ø§Ø¹Ø¯ Ø¹Ù„Ù‰ Ø§Ù„ØªØ¹Ù„Ù… Ø§Ù„Ø³Ø±ÙŠØ¹...",
            "âš™ï¸ Ø¥Ø¹Ø¯Ø§Ø¯ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª Ù„ØªÙƒÙˆÙ† Ø¬Ø§Ù‡Ø²Ø© Ù„Ù„Ù…Ø±Ø§Ø¬Ø¹Ø© Ø§Ù„ÙØ¹Ø§Ù„Ø©...",
            "ğŸ” Ø§Ù„ØªØ±ÙƒÙŠØ² Ø¹Ù„Ù‰ Ø§Ù„Ù…ÙØ§Ù‡ÙŠÙ… Ø§Ù„Ø¬ÙˆÙ‡Ø±ÙŠØ© Ù„Ø¨Ù†Ø§Ø¡ Ø£Ø³Ø§Ø³ Ù‚ÙˆÙŠ...",
            "ğŸš€ ØªØ­Ø³ÙŠÙ† Ù‡ÙŠÙƒÙ„ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª Ù„ØªØ¬Ø±Ø¨Ø© Ø¯Ø±Ø§Ø³ÙŠØ© Ø³Ù„Ø³Ø©...",
            "ğŸ’¡ Ø¬Ø§Ø±ÙŠ ØµÙŠØ§ØºØ© Ø§Ù„Ø£Ø³Ø¦Ù„Ø© ÙˆØ§Ù„Ø£Ø¬ÙˆØ¨Ø© Ø¨ÙˆØ¶ÙˆØ­ ÙˆØ¯Ù‚Ø©..."
        ]


        waiting_messages_quiz = [
            "ğŸ“ Ø¬Ø§Ø±ÙŠ ØµÙŠØ§ØºØ© Ø£Ø³Ø¦Ù„Ø© Ø¯Ù‚ÙŠÙ‚Ø© Ù…Ù† Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø§Ù„Ù…Ù‚Ø¯Ù…...",
            "ğŸ¯ ØªØ­Ø¯ÙŠØ¯ Ø£Ù‡Ù… Ø§Ù„Ø£ÙÙƒØ§Ø± Ù„ÙˆØ¶Ø¹ Ø£Ø³Ø¦Ù„Ø© ØªÙ‚ÙŠØ³ Ø§Ù„ÙÙ‡Ù… Ø§Ù„Ø­Ù‚ÙŠÙ‚ÙŠ...",
            "âš–ï¸ Ù…ÙˆØ§Ø²Ù†Ø© Ø£Ù†ÙˆØ§Ø¹ Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ù„ØªØºØ·ÙŠØ© Ø´Ø§Ù…Ù„Ø© Ù„Ù„Ù…ÙˆØ¶ÙˆØ¹...",
            "ğŸ§  ØªØµÙ…ÙŠÙ… Ø§Ø®ØªØ¨Ø§Ø± Ø°ÙƒÙŠ ÙŠØªØ­Ø¯Ù‰ Ù…Ø¹Ø±ÙØªÙƒ Ø¨Ø´ÙƒÙ„ Ø¨Ù†Ù‘Ø§Ø¡...",
            "ğŸ“Š ØªÙ†Ø¸ÙŠÙ… Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ù„ØªÙ‚Ø¯ÙŠÙ… ØªÙ‚ÙŠÙŠÙ… Ù…ØªÙƒØ§Ù…Ù„ Ù„Ù…Ø³ØªÙˆØ§Ùƒ...",
            "ğŸ§© Ø¨Ù†Ø§Ø¡ Ø§Ø®ØªØ¨Ø§Ø± Ù…ØªÙ…Ø§Ø³Ùƒ ÙŠÙ‚ÙŠØ³ Ù…Ø®ØªÙ„Ù Ø¬ÙˆØ§Ù†Ø¨ Ø§Ù„Ù…Ø¹Ø±ÙØ©...",
            "ğŸ” ÙØ­Øµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø¨Ø¹Ù†Ø§ÙŠØ© Ù„ØµÙŠØ§ØºØ© Ø£Ø³Ø¦Ù„Ø© Ù…Ø­ÙƒÙ…Ø©...",
            "âœ… Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø¬ÙˆØ¯Ø© Ø§Ù„Ø£Ø³Ø¦Ù„Ø© Ù„Ø¶Ù…Ø§Ù† Ø§Ø®ØªØ¨Ø§Ø± Ø¹Ø§Ø¯Ù„..."
        ]


        progress_messages = [
            "â³ Ù„Ø­Ø¸Ø§Øª ÙˆÙ†Ø¨Ø¯Ø£... 25% Ù…Ù† Ø§Ù„Ø·Ø±ÙŠÙ‚",
            "âš¡ï¸ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø© ÙÙŠ Ù…Ù†ØªØµÙÙ‡Ø§... 50% Ù…Ù† Ø§Ù„Ø·Ø±ÙŠÙ‚",
            "ğŸš€ Ø§Ù‚ØªØ±Ø¨Ù†Ø§ Ù…Ù† Ø§Ù„Ø¥Ù†Ø¬Ø§Ø²... 75% Ù…Ù† Ø§Ù„Ø·Ø±ÙŠÙ‚",
            "ğŸ‰ Ø§Ù„Ù„Ù…Ø³Ø§Øª Ø§Ù„Ø£Ø®ÙŠØ±Ø©... 90% Ù…Ù† Ø§Ù„Ø·Ø±ÙŠÙ‚"
        ]




         # ============================
        # Awaiting AI Anki
        # ============================
        if state == "awaiting_anki_file_ai":
            logging.info("Handling awaiting_anki_file_ai for uid=%s", uid)
            if not can_generate(uid):
                return bot.send_message(uid, "âš ï¸ Ù„Ù‚Ø¯ Ø§Ø³ØªÙ†ÙØ¯Øª 3 Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ù…Ø¬Ø§Ù†ÙŠØ© Ù‡Ø°Ø§ Ø§Ù„Ø´Ù‡Ø±.")


    
            # Ø¥Ø¹Ø¯Ø§Ø¯ Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ø£ÙˆÙ„ÙŠØ©
            loading_msg = safe_edit_or_send("ğŸ”„ Ø¬Ø§Ø±ÙŠ Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…Ù„Ù...", chat_id, message_id)
    
            try:
                # Ø¥Ø°Ø§ ÙƒØ§Ù† Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙƒØ¨ÙŠØ±Ø§Ù‹
                if len(content) > 10000:
                    try:
                        # ØªØ­Ø¯ÙŠØ« Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ù„Ø¹Ù…Ù„ÙŠØ© Ø§Ù„ØªÙ„Ø®ÙŠØµ
                        bot.edit_message_text(
                            chat_id=chat_id,
                            message_id=message_id,
                            text="ğŸ“š Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙƒØ¨ÙŠØ± Ø¬Ø¯Ø§Ù‹\nğŸ” Ø¬Ø§Ø±ÙŠ ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰..."
                        )
                
                        content = summarize_long_text(content)
                
                        # ØªØ£ÙƒÙŠØ¯ Ù†Ø¬Ø§Ø­ Ø§Ù„ØªÙ„Ø®ÙŠØµ
                        bot.edit_message_text(
                            chat_id=uid,
                            message_id=loading_msg.message_id,
                            text="âœ… ØªÙ… ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø¨Ù†Ø¬Ø§Ø­\nâ³ Ø¬Ø§Ø±ÙŠ Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª..."
                        )
                        time.sleep(1)
                
                    except Exception as e:
                        print("[ERROR] ÙØ´Ù„ ÙÙŠ ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰:", e)
                        return bot.edit_message_text(
                            chat_id=uid,
                            message_id=loading_msg.message_id,
                            text="âŒ ÙØ´Ù„ ÙÙŠ ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰.\n\nÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ù Ø£ØµØºØ± Ø£Ùˆ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù„Ø§Ø­Ù‚Ø§Ù‹."
                        )
        
                # Ù…Ø¤Ø´Ø± ØªÙ‚Ø¯Ù… Ù…ØªØ­Ø±Ùƒ
                progress_phrases = [
                    "ğŸ“– Ø¬Ø§Ø±ÙŠ ØªØ­Ù„ÙŠÙ„ Ø§Ù„Ù…Ø­ØªÙˆÙ‰...",
                    "ğŸ§  Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…Ø¹Ù„ÙˆÙ…Ø§Øª...",
                    "ğŸ› ï¸ Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª...",
                    "âœ¨ Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªÙ†Ø³ÙŠÙ‚ Ø§Ù„Ù†Ù‡Ø§Ø¦ÙŠ..."
                ]
        
                for i, phrase in enumerate(progress_phrases):
                    # Ø¥Ø¶Ø§ÙØ© Ø´Ø±ÙŠØ· ØªÙ‚Ø¯Ù… Ø¨ØµØ±ÙŠ
                    progress_bar = "[" + "=" * (i+1) + " " * (len(progress_phrases)-i-1) + "]"
            
                    bot.edit_message_text(
                        chat_id=uid,
                        message_id=loading_msg.message_id,
                        text=f"{progress_bar}\n\n{phrase}\n\nâ³ ÙŠØ±Ø¬Ù‰ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø±..."
                    )
                    time.sleep(1.5)
        
                # Ø¥Ø¶Ø§ÙØ© Ø±Ø³Ø§Ù„Ø© Ø§Ù†ØªØ¸Ø§Ø± Ø¬Ø°Ø§Ø¨Ø©
                bot.edit_message_text(
                    chat_id=uid,
                    message_id=loading_msg.message_id,
                    text=f"ğŸ¯ {random.choice(waiting_messages_anki)}\n\nâš¡ Ø¬Ø§Ø±ÙŠ Ø§Ù„Ø§Ù†ØªÙ‡Ø§Ø¡ Ù…Ù† Ø§Ù„ØªØ­Ø¶ÙŠØ±..."
                )
                time.sleep(random.randint(2, 5))
        
                # Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª
                if not can_generate(uid):
                    cards, title = generate_anki_cards_from_text(content, major=major, user_id=uid)
                    
                else:
                    cards, title = generate_special_anki_cards_from_text(content, major=major, user_id=uid)

                if not cards:
                    return bot.edit_message_text(
                        chat_id=uid,
                        message_id=loading_msg.message_id,
                        text="âŒ Ù„Ù… Ø£ØªÙ…ÙƒÙ† Ù…Ù† Ø¥Ù†Ø´Ø§Ø¡ Ø£ÙŠ Ø¨Ø·Ø§Ù‚Ø§Øª.\n\nÙ‚Ø¯ ÙŠÙƒÙˆÙ† Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ØºÙŠØ± Ù…Ù†Ø§Ø³Ø¨ Ø£Ùˆ Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø©."
                    )
                    

               # Ù‚Ø¯ ØªØ­ØªÙˆÙŠ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª Ø¹Ù„Ù‰ image_hint ÙÙ‚Ø·Ø› Ø­ÙˆÙ„Ù‡Ø§ Ù„ØµÙŠØº URL Ù‚Ø¨Ù„ Ø§Ù„Ø­ÙØ¸ Ø¥Ù† Ø£Ø±Ø¯Øª:
                for c in cards:
                    hint = c.get("image_hint", "").strip()
                    if hint and not c.get("image_url"):
                        c["image_url"] = search_image(hint)  # ÙŠÙ…ÙƒÙ† Ø±Ø¬ÙˆØ¹ "" Ø¥Ù† Ù„Ù… ØªÙˆØ¬Ø¯ ØµÙˆØ±Ø©

        
                # ØªÙ†Ø¸ÙŠÙ Ø§Ù„Ø¹Ù†ÙˆØ§Ù† Ù„ÙŠÙƒÙˆÙ† Ø§Ø³Ù… Ù…Ù„Ù ØµØ§Ù„Ø­
                safe_title = re.sub(r'[^a-zA-Z0-9_\u0600-\u06FF]', '_', title)[:40]
                filename = f"{safe_title}_{uid}.apkg"
        
                # Ø­ÙØ¸ Ø§Ù„Ù…Ù„Ù ÙˆØ¥Ø±Ø³Ø§Ù„Ù‡ Ù…Ø¹ ØªØ­Ø¯ÙŠØ« Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø³Ø§Ø¨Ù‚Ø©
                filepath = save_cards_to_apkg(cards, filename=filename, deck_name=title)
        
                 # ØªØ­Ø±ÙŠØ± Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø£Ø®ÙŠØ±Ø© Ù„Ø¥Ø¸Ù‡Ø§Ø± Ù†Ø¬Ø§Ø­ Ø§Ù„Ø¹Ù…Ù„ÙŠØ©
                bot.edit_message_text(
                    chat_id=uid,
                    message_id=loading_msg.message_id,
                    text=f"âœ… ØªÙ… Ø¥Ù†Ø´Ø§Ø¡ {len(cards)} Ø¨Ø·Ø§Ù‚Ø© Ø¨Ù†Ø¬Ø§Ø­!\n\nğŸ“š Ø§Ù„Ø¹Ù†ÙˆØ§Ù†: {title}\n\nâš¡ Ø¬Ø§Ø±ÙŠ Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù…Ù„Ù..."
                )
                increment_count(uid)
                notify_admin("ØªÙˆÙ„ÙŠØ¯ Ø£Ù†ÙƒÙŠ Ø¢Ù„ÙŠ", username, uid)

                # Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù…Ù„Ù Ù…Ø¹ caption
                with open(filepath, 'rb') as file:
                    bot.send_document(
                        chat_id=uid,
                        document=file,
                        caption=f"ğŸ“‚ {title}\n\nğŸ´ Ø¹Ø¯Ø¯ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª: {len(cards)}\n\nØ§Ø³ØªÙ…ØªØ¹ Ø¨Ø§Ù„Ø¯Ø±Ø§Ø³Ø©!",
                        reply_to_message_id=loading_msg.message_id
                    )
                    with state_lock:
                        user_states.pop(uid, None)
                    logging.info("Finished ai_anki for uid=%s", uid)



            except Exception:
                logging.exception("Error while processing ai anki for uid=%s", uid)
                with state_lock:
                    user_states.pop(uid, None)
                    bot.send_message(uid, "Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª. Ø­Ø§ÙˆÙ„ Ù„Ø§Ø­Ù‚Ù‹Ø§.")
            return


        # ============================
        # Awaiting manual anki
        # ============================
        elif state == "awaiting_anki_file_manual":
            if msg.content_type == "text":
            # Ù‚Ø§Ø¦Ù…Ø© Ø±Ø³Ø§Ø¦Ù„ Ø§Ù†ØªØ¸Ø§Ø± Ù…ØªØ­Ø±ÙƒØ©
                waiting_messages = [
                    "ğŸŒ± Ø¬Ø§Ø±Ù ØªØ¬Ù‡ÙŠØ² Ù…Ø¯Ø®Ù„Ø§ØªÙƒ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠØ©...",
                    "ğŸ§  Ø§Ù„Ù†Ø¸Ø§Ù… Ø§Ù„Ù…Ø¹Ø±ÙÙŠ ÙŠØ¹Ù…Ù„ Ø¨ÙƒØ§Ù…Ù„ Ø·Ø§Ù‚ØªÙ‡...",
                    "ğŸ”® ÙŠØªÙ… ØªÙ†Ø¸ÙŠÙ… Ø§Ù„Ù…Ø­ØªÙˆÙ‰ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠ...",
                    "ğŸš€ ÙŠØ¬Ø±ÙŠ Ø§Ù„Ø¥Ø¹Ø¯Ø§Ø¯ Ù„Ù…Ø±Ø­Ù„Ø© Ø§Ù„Ø¥Ø·Ù„Ø§Ù‚...",
                    "ğŸ© ØªØªÙ… Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ø¨ØµÙˆØ±Ø© Ø¯Ù‚ÙŠÙ‚Ø©..."
                ]
        
                # Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± Ù…Ø¹ Ù…Ø¤Ø«Ø±Ø§Øª Ø¨ØµØ±ÙŠØ©
                waiting_msg = bot.send_message(chat_id, "â³ **Ø¬Ø§Ø±Ù Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø©**\n`0% Ø§ÙƒØªÙ…Ø§Ù„`", 
                                      parse_mode="Markdown")
        
                # Ø®Ø·ÙˆØ§Øª Ø§Ù„ØªÙ‚Ø¯Ù… Ù…Ø¹ Ø±Ù…ÙˆØ² Ø¥Ø¨Ø¯Ø§Ø¹ÙŠØ©
                processing_steps = [
                    {"icon": "ğŸ”", "text": "ØªØ­Ù„ÙŠÙ„ Ø§Ù„Ù†ØµÙˆØµ Ø§Ù„Ù…Ø¯Ø®Ù„Ø©", "delay": 0.8},
                    {"icon": "ğŸ§©", "text": "Ø¨Ù†Ø§Ø¡ Ø§Ù„Ø¨Ù†ÙŠØ© Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠØ©", "delay": 1.2},
                    {"icon": "ğŸ¨", "text": "Ø¥Ø¹Ø¯Ø§Ø¯ Ø¨Ø·Ø§Ù‚Ø§Øª Ø§Ù„Ù…Ø±Ø§Ø¬Ø¹Ø©", "delay": 1.0},
                    {"icon": "âš¡", "text": "ØªØ¬Ù‡ÙŠØ² Ø§Ù„Ø­Ø²Ù…Ø© Ø§Ù„Ù†Ù‡Ø§Ø¦ÙŠØ©", "delay": 0.7},
                    {"icon": "ğŸš€", "text": "Ø¥Ø·Ù„Ø§Ù‚ Ø§Ù„Ù…Ù„Ù Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠ", "delay": 1.5}
                ]
        
                # Ù…Ø­Ø§ÙƒØ§Ø© Ø§Ù„ØªÙ‚Ø¯Ù… Ø§Ù„ØªØ¯Ø±ÙŠØ¬ÙŠ
                progress = 0
                step_size = 100 // len(processing_steps)
        
                for idx, step in enumerate(processing_steps):
                    # Ø­Ø³Ø§Ø¨ Ø§Ù„Ù†Ø³Ø¨Ø© Ø§Ù„Ù…Ø¦ÙˆÙŠØ©
                    progress = min(100, (idx + 1) * step_size)
                    progress_bar = "ğŸŸ©" * (progress // 10) + "â¬œ" * (10 - progress // 10)
            
                    # Ø¨Ù†Ø§Ø¡ Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ù…Ø¹ ØªØ£Ø«ÙŠØ± Ø§Ù„ØªØ±Ø§ÙƒÙ…
                    message_text = (
                        f"â³ **Ø¬Ø§Ø±Ù Ø§Ù„Ù…Ø¹Ø§Ù„Ø¬Ø©**\n"
                        f"`{progress}% Ø§ÙƒØªÙ…Ø§Ù„`\n"
                        f"{progress_bar}\n\n"
                        f"{step['icon']} **Ø§Ù„Ù…Ø±Ø­Ù„Ø© {idx+1}:** {step['text']}"
                    )
            
                    # ØªØ­Ø¯ÙŠØ« Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ù…Ø¹ ØªØ£Ø«ÙŠØ± Ø§Ù„ØªØ¯Ø±Ø¬
                    try:
                        bot.edit_message_text(
                            message_text,
                            chat_id=waiting_msg.chat.id,
                            message_id=waiting_msg.message_id,
                            parse_mode="Markdown"
                )
                    except:
                        pass
            
                    # ØªØ£Ø®ÙŠØ± Ø¯ÙŠÙ†Ø§Ù…ÙŠÙƒÙŠ Ø¨ÙŠÙ† Ø§Ù„Ø®Ø·ÙˆØ§Øª
                    time.sleep(step['delay'])
        
                # Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…Ù„Ù Ø§Ù„ÙØ¹Ù„ÙŠØ©
                cards = parse_manual_anki_input(msg.text)
                if cards:
                    # Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ù…Ù„Ù
                    output_file = f"{uid}_manual_anki.apkg"
                    save_cards_to_apkg(cards, filename=output_file, deck_name="Ù…ÙƒØªØ¨ØªÙƒ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠØ©")
            
                    # Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù…Ù„Ù Ù…Ø¹ Ø±Ø³Ø§Ù„Ø© Ø±Ø³Ù…ÙŠØ©
                    with open(output_file, 'rb') as file:
                        bot.send_document(
                            chat_id=uid,
                            document=file,
                            caption=(
                                f"ğŸŒ¿ *ØªÙ… Ø¥Ù†Ø´Ø§Ø¡ Ù…Ù„ÙÙƒ Ø§Ù„ØªØ¹Ù„ÙŠÙ…ÙŠ Ø¨Ù†Ø¬Ø§Ø­.*\n"
                                f"Ø¹Ø¯Ø¯ Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª: {len(cards)} Ø¨Ø·Ø§Ù‚Ø©\n"
                                f"Ù…Ø¯Ø© Ø§Ù„ØªÙ†ÙÙŠØ°: {random.randint(3,7)} Ø«ÙˆØ§Ù†Ù\n\n"
                                f"ğŸ“š Ù…Ù„Ù Ø§Ù„Ù…Ø±Ø§Ø¬Ø¹Ø© Ø¬Ø§Ù‡Ø² Ù„Ù„Ø§Ø³ØªØ®Ø¯Ø§Ù…."
                            ),
                            reply_to_message_id=message_id,
                            parse_mode="Markdown"
                        )
                        notify_admin("ØªÙˆÙ„ÙŠØ¯ Ø£Ù†ÙƒÙŠ ÙŠØ¯ÙˆÙŠ", username, uid)

            
                    # Ø­Ø°Ù Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªÙ‚Ø¯Ù… Ø¨Ø¹Ø¯ Ø§Ù„Ø¥Ø±Ø³Ø§Ù„
                    try:
                        bot.delete_message(chat_id, waiting_msg.message_id)
                    except:
                        pass
            
                    # Ø±Ø³Ø§Ù„Ø© Ø®ØªØ§Ù…ÙŠØ© Ø±Ø³Ù…ÙŠØ©
                    if random.random() < 0.1:  # 30% Ø§Ø­ØªÙ…Ø§Ù„
                        bot.send_message(
                            chat_id,
                            "âœ¨ *Ø£Ø­Ø³Ù†Øª! ØªÙ… Ø¥Ù†Ø´Ø§Ø¡ Ø¨Ø·Ø§Ù‚Ø§ØªÙƒ Ø¨Ù†Ø¬Ø§Ø­.*\n"
                            "ğŸš€ ÙˆØ§ØµÙ„ Ø§Ù„Ù…Ø±Ø§Ø¬Ø¹Ø© Ø¨Ø§Ù†ØªØ¸Ø§Ù…ØŒ ÙˆØ³ØªØªÙØ§Ø¬Ø£ Ø¨Ø³Ø±Ø¹Ø© ØªÙ‚Ø¯Ù…Ùƒ.\n"
                            "â”€â”€â”€ â‹†â‹…â˜†â‹…â‹† â”€â”€\n"
                            "ğŸ’¡ *ØªØ°ÙƒÙŠØ±:* Ø±Ø§Ø¬Ø¹ Ø¨Ø·Ø§Ù‚Ø§ØªÙƒ ØºØ¯Ù‹Ø§ØŒ ÙØ§Ù„ØªÙƒØ±Ø§Ø± Ù‡Ùˆ Ø³Ø± ØªØ«Ø¨ÙŠØª Ø§Ù„Ù…Ø¹Ø±ÙØ©!",
                            reply_to_message_id=message_id,
                            parse_mode="Markdown"
                    )
                    with state_lock:
                        user_states.pop(uid, None)

                else:
                    # Ø­Ø°Ù Ø±Ø³Ø§Ø¦Ù„ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø± ÙˆØ§Ù„Ø®Ø·ÙˆØ§Øª ÙˆØ¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø®Ø·Ø£
                    bot.delete_message(uid, waiting_msg.message_id)
                    for step_msg in step_messages:
                        bot.delete_message(uid, step_msg.message_id)
                    bot.send_message(uid, "âŒ Ù„Ù… ÙŠØªÙ… Ø§Ù„Ø¹Ø«ÙˆØ± Ø¹Ù„Ù‰ Ø£ÙŠ Ø¨Ø·Ø§Ù‚Ø§Øª ØµØ§Ù„Ø­Ø©")
            else:
                bot.send_message(uid, "âŒ ÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù†Øµ ÙÙ‚Ø· Ù„Ø¥Ù†Ø´Ø§Ø¡ Ø¨Ø·Ø§Ù‚Ø§Øª Anki ÙŠØ¯ÙˆÙŠÙ‹Ø§.")

        # ============================
        # Awaiting advanced test file
        # ============================
        # Ø§Ø³ØªØ®Ø¯Ù… Ù…ÙƒØªØ¨Ø© traceback Ù„ØªØ´Ø®ÙŠØµ Ø¯Ù‚ÙŠÙ‚
        elif state == "awaiting_advanced_test_file":
            import traceback
    
            # Ø§Ø­ØªÙØ¸ Ø¨Ø§Ù„Ù€ ID Ø§Ù„Ø£ØµÙ„ÙŠ Ù„Ù„Ø±Ø³Ø§Ù„Ø© ÙˆØ§Ù„Ø¯Ø±Ø¯Ø´Ø© ÙÙŠ Ù…ØªØºÙŠØ±Ø§Øª Ø¢Ù…Ù†Ø©
            original_chat_id = msg.chat.id
            original_message_id = message_id

            try:
                # 1. Ø§Ø¨Ø¯Ø£ Ø¨ØªØ¹Ø¯ÙŠÙ„ Ø§Ù„Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø£ÙˆÙ„Ù‰
                bot.edit_message_text("ğŸ¤– Ø¬Ø§Ø±ÙŠ Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ù…Ù„Ù ÙˆØ¥Ù†Ø´Ø§Ø¡ Ø§Ø®ØªØ¨Ø§Ø± Ø°ÙƒÙŠ...", chat_id=original_chat_id, message_id=original_message_id)
        
                if not can_generate(uid):
                    return bot.send_message(uid, "âš ï¸ Ù„Ù‚Ø¯ Ø§Ø³ØªÙ†ÙØ¯Øª 3 Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ù…Ø¬Ø§Ù†ÙŠØ© Ù‡Ø°Ø§ Ø§Ù„Ø´Ù‡Ø±.")

                if len(content) > 10000:
                    bot.edit_message_text("ğŸ” Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙƒØ¨ÙŠØ±ØŒ Ø¬Ø§Ø±ÙŠ ØªÙ„Ø®ÙŠØµÙ‡...", chat_id=original_chat_id, message_id=original_message_id)
                    content = summarize_long_text(content)
        
                # 2. Ø¹Ø±Ø¶ Ø±Ø³Ø§Ø¦Ù„ Ø§Ù„ØªÙ‚Ø¯Ù… (Ù„Ø§ Ù†ØºÙŠØ± Ù‚ÙŠÙ…Ø© Ø§Ù„Ù…ØªØºÙŠØ±Ø§Øª Ø§Ù„Ø£ØµÙ„ÙŠØ©)
                bot.edit_message_text("ğŸ§  Ø¬Ø§Ø±ÙŠ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±ØŒ Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø±...", chat_id=original_chat_id, message_id=original_message_id)
                for progress_msg in progress_messages:
                    bot.edit_message_text(progress_msg, chat_id=original_chat_id, message_id=original_message_id)
                    time.sleep(1.5)

                bot.edit_message_text(random.choice(waiting_messages_quiz), chat_id=original_chat_id, message_id=original_message_id)
                time.sleep(2)

                # 3. ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
                print("[ADVANCED_QUIZ] Ø¨Ø¯Ø¡ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø§Ù„Ø·Ø¨ÙŠ Ø§Ù„Ù…ØªÙ‚Ø¯Ù…...")
                quiz_data = generate_Medical_quizzes(content=content, major="General Medicine", user_id=uid)
        
                # Ø·Ø¨Ø§Ø¹Ø© Ù„Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø§Ù„Ù‚ÙŠÙ… Ù‚Ø¨Ù„ Ø§Ù„Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù†Ù‡Ø§Ø¦ÙŠ
                print(f"[DEBUG] chat_id: {original_chat_id}, message_id: {original_message_id}, quiz_data is not None: {quiz_data is not None}")
    
                if quiz_data:
                    # 4. Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù†ØªÙŠØ¬Ø© Ø§Ù„Ù†Ù‡Ø§Ø¦ÙŠØ©
                    send_quiz_to_user(original_chat_id, quiz_data)
                    
                    notify_admin("ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø± Ø°ÙƒÙŠ", username, uid)
                    update_top_user(uid, tests=1)
                else:
                    bot.edit_message_text("âŒ ÙØ´Ù„ ÙÙŠ Ø¥Ù†Ø´Ø§Ø¡ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±. Ù‚Ø¯ ÙŠÙƒÙˆÙ† Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ØºÙŠØ± Ù…Ù†Ø§Ø³Ø¨. ÙŠØ±Ø¬Ù‰ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù„Ø§Ø­Ù‚Ø§Ù‹.", chat_id=original_chat_id, message_id=original_message_id)

            except Exception as e:
                # Ø·Ø¨Ø§Ø¹Ø© Ø§Ù„Ø®Ø·Ø£ Ø§Ù„ÙƒØ§Ù…Ù„ ÙÙŠ Ø§Ù„ÙƒÙˆÙ†Ø³ÙˆÙ„ Ù„ØªØ´Ø®ÙŠØµÙ‡
                print("!!!!!!!!!!!!!!!!!! Ø®Ø·Ø£ ÙØ§Ø¯Ø­ ÙÙŠ Ù…Ø¹Ø§Ù„Ø¬Ø© Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø§Ù„Ù…ØªÙ‚Ø¯Ù… !!!!!!!!!!!!!!!!!!")
                traceback.print_exc()
                print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        
                error_message = "âš ï¸ Ø­Ø¯Ø« Ø®Ø·Ø£ ØªÙ‚Ù†ÙŠ. ØªÙ… Ø¥Ø¨Ù„Ø§Øº Ø§Ù„Ù…Ø·ÙˆØ±ÙŠÙ†."
                try:
                    bot.edit_message_text(error_message, chat_id=original_chat_id, message_id=original_message_id)
                except:
                    bot.send_message(original_chat_id, error_message)
            finally:
                # ØªØ£ÙƒØ¯ Ù…Ù† Ø¥Ø²Ø§Ù„Ø© Ø­Ø§Ù„Ø© Ø§Ù„Ù…Ø³ØªØ®Ø¯Ù… Ù„ØªØ¬Ù†Ø¨ Ø¨Ù‚Ø§Ø¦Ù‡ Ø¹Ø§Ù„Ù‚Ø§Ù‹
                with state_lock:
                    user_states.pop(uid, None)

                
        # ============================
        # Awaiting simple test file
        # ============================
        elif state == "awaiting_simple_test_file":
            print("[QUIZ] Ø¨Ø¯Ø¡ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ù„Ù„Ù…Ø³ØªØ®Ø¯Ù…:", uid)
            if not can_generate(uid):
                return bot.send_message(uid, "âš ï¸ Ù„Ù‚Ø¯ Ø§Ø³ØªÙ†ÙØ¯Øª 3 Ø§Ø®ØªØ¨Ø§Ø±Ø§Øª Ù…Ø¬Ø§Ù†ÙŠØ© Ù‡Ø°Ø§ Ø§Ù„Ø´Ù‡Ø±.")


            if len(content) > 10000:
                loading_msg = bot.edit_message_text("ğŸ” Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙƒØ¨ÙŠØ±ØŒ Ø¬Ø§Ø±ÙŠ ØªÙ„Ø®ÙŠØµÙ‡...", chat_id=chat_id, message_id=message_id)
                try:
                    print("[QUIZ] Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙƒØ¨ÙŠØ±ØŒ Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªÙ„Ø®ÙŠØµ...")
                    content = summarize_long_text(content)
                except Exception as e:
                    print("[ERROR] ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰ ÙØ´Ù„:", e)
                    return bot.send_message(uid, "âŒ ÙØ´Ù„ ÙÙŠ ØªÙ„Ø®ÙŠØµ Ø§Ù„Ù…Ø­ØªÙˆÙ‰. Ø£Ø±Ø³Ù„ Ù…Ù„ÙÙ‹Ø§ Ø£ØµØºØ± Ø£Ùˆ Ø­Ø§ÙˆÙ„ Ù„Ø§Ø­Ù‚Ù‹Ø§.")

            else:
                loading_msg = bot.edit_message_text("ğŸ§  Ø¬Ø§Ø±ÙŠ ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±ØŒ Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø±...", chat_id=chat_id, message_id=message_id)

                # Ø¹Ø±Ø¶ Ø±Ø³Ø§Ø¦Ù„ Ø§Ù„ØªØ­Ù…ÙŠÙ„
            for progress_msg in progress_messages:
                try:
                    bot.edit_message_text(chat_id=uid, message_id=loading_msg.message_id, text=progress_msg)
                except Exception as e:
                    logging.exception("[QUIZ] ÙØ´Ù„ Ø£Ø«Ù†Ø§Ø¡ ØªØ­Ø¯ÙŠØ« Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªÙ‚Ø¯Ù‘Ù…")

                    time.sleep(1.5)
            try:
                bot.edit_message_text(chat_id=uid, message_id=loading_msg.message_id,
                              text=random.choice(waiting_messages_quiz))
            except Exception:
                pass

            time.sleep(2)

            print("[QUIZ] Ø§Ø³ØªØ¯Ø¹Ø§Ø¡ generate_quizzes_from_text...")
            quizzes = generate_quizzes_from_text(content, major=major, user_id=uid, num_quizzes=10)
            print("[QUIZ] Ø±Ø¬Ø¹:", type(quizzes), "Ø¨Ø·ÙˆÙ„:", len(quizzes) if quizzes else "None")

        
            if isinstance(quizzes, list) and len(quizzes) > 0:
                
                try:
                    print(f"ØªÙ… ØªÙˆÙ„ÙŠØ¯ {len(quizzes)} Ø³Ø¤Ø§Ù„Ø§")
                    # ØªØ®Ø²ÙŠÙ† Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø± Ø£ÙˆÙ„Ø§Ù‹
                    quiz_code = store_quiz(uid, quizzes, bot)
                    print("[QUIZ] ÙƒÙˆØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±:", quiz_code)

                    if not quiz_code:
                        raise Exception("Failed to store quiz")
                        
                    waiting_quiz = loading_msg.message_id
                    major = fetch_user_major(uid)
                    file_path = user_files[uid]
                    level = "Ù…ØªÙˆØ³Ø·"

                    # Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© "Ø¥Ø®ØªØ¨Ø§Ø±Ùƒ Ø¬Ø§Ù‡Ø²" Ù…Ø¹ Ø±Ø§Ø¨Ø· Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
                    quiz_link = f"https://t.me/QuizzyAI_bot?start=quiz_{quiz_code}"
                    estimated_time = len(quizzes) * 30

                    # Ø¥Ø±Ø³Ø§Ù„ Ø±Ø³Ø§Ù„Ø© "Ø¥Ø®ØªØ¨Ø§Ø±Ùƒ Ø¬Ø§Ù‡Ø²" Ù…Ø¹ Ø±Ø§Ø¨Ø· Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±
                    markup = InlineKeyboardMarkup()
                    btn = InlineKeyboardButton("ÙØªØ­ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±", url=quiz_link)
                    markup.add(btn)

                    quiz_msg = (
                    "âœ¨âœ”ï¸ <b>Ø¥Ø®ØªØ¨Ø§Ø±Ùƒ Ø¬Ø§Ù‡Ø²!</b>\n"
                    "â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€\n"
                    f"ğŸ“‚ <b>Ø§Ù„Ø¹Ù†ÙˆØ§Ù†:</b> {msg.document.file_name}\n\n"
                    f"ğŸ“‹ <b>Ø¹Ø¯Ø¯ Ø§Ù„Ø£Ø³Ø¦Ù„Ø©:</b> {len(quizzes)}\n"
                    f"â±ï¸ <b>Ø§Ù„Ø²Ù…Ù† Ø§Ù„ÙƒÙ„ÙŠ:</b> {estimated_time // 60} Ø¯Ù‚ÙŠÙ‚Ø© Ùˆ {estimated_time % 60} Ø«Ø§Ù†ÙŠØ©\n"
                    f"ğŸ“ <b>Ø§Ù„ØªØ®ØµØµ:</b> {major} \n"
                    "ğŸ“¦ <b>Ù†ÙˆØ¹ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±:</b> Ø®Ø§Øµ\n\n"
                    f"ğŸ“‰ <b>Ø§Ù„ØªØºØ·ÙŠØ©:</b> {coverage}\n"
                    "ğŸ’¡ <b>Ù…ÙŠØ²Ø© Ø§Ù„Ø´Ø±Ø­:</b> ØºÙŠØ± Ù…ØªÙˆÙØ±Ø©\n"
                    f"ğŸ“Š <b>Ø§Ù„Ù…Ø³ØªÙˆÙ‰:</b> {level}\n\n"
                    "â“Ù‡Ù„ Ø£Ù†Øª Ø¬Ø§Ù‡Ø² Ù„Ù„Ø¥Ø®ØªØ¨Ø§Ø±\n"
                    f"ğŸ‘ˆ <a href=\"{quiz_link}\">Ø§Ø¶ØºØ· Ù‡Ù†Ø§ Ù„Ù„Ø¨Ø¯Ø¡</a>"
                    )
                    try:
                        bot.delete_message(chat_id=chat_id, message_id=loading_msg.message_id)
                    except Exception as del_err:
                        print(f"Ù„Ù… ÙŠØªÙ…ÙƒÙ† Ù…Ù† Ø­Ø°Ù Ø±Ø³Ø§Ù„Ø© Ø§Ù„ØªØ­Ù…ÙŠÙ„: {del_err}")
                
                    bot.send_message(chat_id, quiz_msg, reply_markup=markup, parse_mode="HTML", disable_web_page_preview=True)
                    

                    update_top_user(uid, tests=1)
                    notify_admin("ØªÙˆÙ„ÙŠØ¯ Ø§Ø®ØªØ¨Ø§Ø±", username, uid)
                    
                    

                    with state_lock:
                        user_states.pop(uid, None)

                except Exception as e:
                    print(f"Error in quiz generation: {e}")
                    bot.send_message(uid, "Ø­Ø¯Ø« Ø®Ø·Ø£ ØºÙŠØ± Ù…ØªÙˆÙ‚Ø¹   .")
                    
            else:
                print("[QUIZ] Ø§Ù„ØªÙˆÙ„ÙŠØ¯ ÙØ´Ù„ Ø£Ùˆ Ø±Ø¬Ø¹ None")
                bot.send_message(uid, "âŒ Ù„Ù… ÙŠØªÙ…ÙƒÙ† Ø§Ù„Ø¨ÙˆØª Ù…Ù† ØªÙˆÙ„ÙŠØ¯ Ø§Ù„Ø§Ø®ØªØ¨Ø§Ø±.")


        else:
            bot.reply_to(msg, "Ø¹Ø°Ø±Ù‹Ø§ØŒ Ù„Ø§ Ø£Ø³ØªØ·ÙŠØ¹ ÙÙ‡Ù… Ø·Ù„Ø¨Ùƒ ÙÙŠ Ø§Ù„ÙˆÙ‚Øª Ø§Ù„Ø­Ø§Ù„ÙŠ. ÙŠØ±Ø¬Ù‰ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© Ù„Ø§Ø­Ù‚Ù‹Ø§ Ø£Ùˆ Ø§Ù„Ø§ØªØµØ§Ù„ Ø¨Ø§Ù„Ø¯Ø¹Ù….")
            # Ø£Ùˆ ÙŠÙ…ÙƒÙ†Ùƒ ØªØ³Ø¬ÙŠÙ„ Ø°Ù„Ùƒ ÙÙŠ Ù‚Ø§Ø¹Ø¯Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ø£Ùˆ Ù…Ù„Ù Ø§Ù„Ø³Ø¬Ù„Ø§Øª
            print("No state for user", uid, "state:", state)


    except Exception as e:
        import traceback
        logging.exception("process_message error: %s", e)
        print("!!!!!!!!!!!!!!!!! Ø­Ø¯Ø« Ø®Ø·Ø£ !!!!!!!!!!!!!!!!!!")
        traceback.print_exc() # Ù‡Ø°Ø§ Ø§Ù„Ø³Ø·Ø± Ø³ÙŠØ·Ø¨Ø¹ Ø§Ù„Ø®Ø·Ø£ Ø§Ù„ÙƒØ§Ù…Ù„ ÙˆÙ…ÙƒØ§Ù†Ù‡ Ø¨Ø§Ù„Ø¶Ø¨Ø·
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        bot.send_message(uid, "Ø­Ø¯Ø« Ø®Ø·Ø£ ØºÙŠØ± Ù…ØªÙˆÙ‚Ø¹   .")
    finally:
        # Ø­Ø°Ù Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…Ø¤Ù‚Øª Ø¥Ù† ÙˆÙØ¬Ø¯
        if path and os.path.exists(path):
            try:
                os.remove(path)
            except Exception as e:
                print(f"[WARNING] Ù„Ù… ÙŠØªÙ… Ø­Ø°Ù Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…Ø¤Ù‚Øª: {e}")



known_channels = set()


@bot.channel_post_handler(func=lambda msg: True)
def handle_channel_post(msg):
    channel_id = msg.chat.id

    if channel_id in known_channels:
        return  # ØªÙ… Ù…Ø¹Ø§Ù„Ø¬ØªÙ‡ Ù…Ù† Ù‚Ø¨Ù„

    known_channels.add(channel_id)

    try:
        bot.send_message(
            ADMIN_ID,
            f"ğŸ“¢ ØªÙ… Ø§ÙƒØªØ´Ø§Ù Ù‚Ù†Ø§Ø© Ø¬Ø¯ÙŠØ¯Ø©:\n\n"
            f"*Ø§Ù„Ø§Ø³Ù…:* {msg.chat.title}\n"
            f"*ID:* `{channel_id}`",
            parse_mode="Markdown"
        )
    except Exception as e:
        print(f"[ERROR] Ø¥Ø±Ø³Ø§Ù„ Ø§Ù„Ù…Ø¹Ø±Ù ÙØ´Ù„: {e}")


# -------------------------------------------------------------------
#                   inference handler
# -------------------------------------------------------------------


@bot.message_handler(commands=['submit_inference'])
def handle_submit_inference(msg):
    if msg.chat.type != "private":
        return  # ØªØ¬Ø§Ù‡Ù„ Ø§Ù„Ø±Ø³Ø§Ø¦Ù„ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø§Øª
    uid = msg.from_user.id
    user_states[uid] = {"state": "awaiting_inference_question", "temp": {}}
    bot.send_message(uid, "ğŸ§  Ø£Ø±Ø³Ù„ Ø§Ù„Ø¢Ù† Ø³ÙŠÙ†Ø§Ø±ÙŠÙˆ Ø£Ùˆ Ø³Ø¤Ø§Ù„Ù‹Ø§ Ù„Ù„Ø§Ø¹Ø¨ÙŠÙ† (Ù…Ø«Ø§Ù„: ÙƒÙŠÙ ØªØªØµØ±Ù ÙÙŠ Ù‡Ø°Ø§ Ø§Ù„Ù…ÙˆÙ‚ÙØŸ)")

@bot.message_handler(func=lambda m: user_states.get(m.from_user.id, {}).get("state") in [
    "awaiting_inference_question", "awaiting_inference_options", "awaiting_inference_correct"])
def handle_inference_submission(msg):
    if msg.hat.type != "private":
        return  # ØªØ¬Ø§Ù‡Ù„ Ø§Ù„Ø±Ø³Ø§Ø¦Ù„ ÙÙŠ Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø§Øª
    uid = msg.from_user.id
    state = user_states.get(uid, {})
    temp = state.get("temp", {})

    if state["state"] == "awaiting_inference_question":
        temp["question"] = msg.text.strip()
        user_states[uid] = {"state": "awaiting_inference_options", "temp": temp}
        bot.send_message(uid, "âœï¸ Ø£Ø±Ø³Ù„ Ø§Ù„Ø¢Ù† 4 Ø®ÙŠØ§Ø±Ø§ØªØŒ ÙƒÙ„ Ø®ÙŠØ§Ø± ÙÙŠ Ø³Ø·Ø± Ù…Ù†ÙØµÙ„.")

    elif state["state"] == "awaiting_inference_options":
        options = msg.text.strip().split("\n")
        if len(options) != 4:
            return bot.send_message(uid, "âš ï¸ ÙŠØ¬Ø¨ Ø£Ù† ØªØ±Ø³Ù„ 4 Ø®ÙŠØ§Ø±Ø§Øª ÙÙ‚Ø·ØŒ ÙƒÙ„ Ø®ÙŠØ§Ø± ÙÙŠ Ø³Ø·Ø±.")
        temp["options"] = options
        user_states[uid] = {"state": "awaiting_inference_correct", "temp": temp}
        bot.send_message(uid, "âœ… Ù…Ø§ Ù‡Ùˆ Ø±Ù‚Ù… Ø§Ù„Ø¥Ø¬Ø§Ø¨Ø© Ø§Ù„ØµØ­ÙŠØ­Ø©ØŸ (Ù…Ù† 0 Ø¥Ù„Ù‰ 3)")

    elif state["state"] == "awaiting_inference_correct":
        try:
            correct = int(msg.text.strip())
            if correct not in [0, 1, 2, 3]:
                raise ValueError()
        except:
            return bot.send_message(uid, "âš ï¸ Ø£Ø±Ø³Ù„ Ø±Ù‚Ù…Ù‹Ø§ ØµØ­ÙŠØ­Ù‹Ø§ Ù…Ù† 0 Ø¥Ù„Ù‰ 3 ÙÙ‚Ø·.")
        
        # Ø­ÙØ¸ ÙÙŠ Ù‚Ø§Ø¹Ø¯Ø© Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª
        q = temp["question"]
        options = temp["options"]
        options_str = json.dumps(options)
        cursor.execute("""
        INSERT INTO inference_questions (question, options, correct_index, submitted_by)
        VALUES (?, ?, ?, ?)
        """, (q, options_str, correct, uid))
        conn.commit()

        user_states.pop(uid, None)
        bot.send_message(uid, "ğŸ‰ ØªÙ… Ø§Ø³ØªÙ„Ø§Ù… Ø§Ù‚ØªØ±Ø§Ø­Ùƒ Ø¨Ù†Ø¬Ø§Ø­! Ø³ÙŠØªÙ… Ù…Ø±Ø§Ø¬Ø¹ØªÙ‡ Ù‚Ø±ÙŠØ¨Ù‹Ø§. Ø´ÙƒØ±Ø§Ù‹ Ù„Ù…Ø³Ø§Ù‡Ù…ØªÙƒ ğŸ™")
# -------------------------------------------------------------------
#                           Run Bot
# -------------------------------------------------------------------

@app.route('/anki_preview')
def anki_preview():
    user_cards = generate_anki_cards_from_text(text)[:5]  # â† Ù†Ø­ØµÙ„ Ø¹Ù„Ù‰ Ø£ÙˆÙ„ 5 Ø¨Ø·Ø§Ù‚Ø§Øª
    session['cards'] = user_cards
    session['index'] = 0
    session['show_back'] = False
    return redirect('/anki')
    
app.secret_key = 'anki_secret'  # Ø³Ø± Ø§Ù„Ø¬Ù„Ø³Ø© Ù„ØªØ®Ø²ÙŠÙ† Ø§Ù„Ø¨ÙŠØ§Ù†Ø§Øª Ù…Ø¤Ù‚ØªÙ‹Ø§


@app.route('/anki', methods=['GET', 'POST'])
def anki_cards():
    content = session.get('anki_content')
    major = session.get('anki_major', 'General')
    if 'cards' not in session:
        session['cards'] = example_cards[:5]
        session['index'] = 0
        session['show_back'] = False

    if request.method == 'POST':
        action = request.form.get('action')
        if action == 'show':
            session['show_back'] = True
        elif action == 'next':
            session['index'] += 1
            session['show_back'] = False

    index = session['index']
    cards = session['cards']

    if index >= len(cards):
        session.clear()
        return "<h2>ğŸ‰ Ø§Ù†ØªÙ‡ÙŠØª Ù…Ù† Ø§Ù„Ø¨Ø·Ø§Ù‚Ø§Øª! Ø£Ø­Ø³Ù†Øª.</h2><a href='/anki'>ğŸ” Ø§Ø¨Ø¯Ø£ Ù…Ù† Ø¬Ø¯ÙŠØ¯</a>"

    return render_template('anki_viewer.html',
                           card=cards[index],
                           index=index,
                           total=len(cards),
                           show_back=session['show_back'])
# Ø¨Ø¯Ø¡ Ø§Ù„Ø¨ÙˆØª
start_workers()



import json
from datetime import datetime

def insert_sample_quiz_if_not_exists(db_path='quiz_users.db'):
    conn = sqlite3.connect(db_path, check_same_thread=False)
    cursor = conn.cursor()
        
    """
    Inserts a sample quiz into the database if it doesn't already exist.

    Args:
        cursor: The database cursor object.
        conn: The database connection object.
    """
    cursor.execute("SELECT quiz_code FROM sample_quizzes WHERE quiz_code = ?", ("sample",))
    if cursor.fetchone() is None:
        # The list of questions should be a single list of dictionaries, not a list within a list.
        sample_quiz_data = [
            {
                "question": "Ù…Ø§ Ù‡Ùˆ Ø£Ø·ÙˆÙ„ Ø¨Ø±Ø¬ ÙÙŠ Ø§Ù„Ø¹Ø§Ù„Ù…ØŸ",
                "options": ["Ø¨Ø±Ø¬ Ø®Ù„ÙŠÙØ©", "Ø¨Ø±Ø¬ Ø¥ÙŠÙÙ„", "Ø¨Ø±Ø¬ Ø¨ÙŠØ²Ø§", "Ø¨Ø±Ø¬ Ø´Ù†ØºÙ‡Ø§ÙŠ"],
                "correct_index": 0,
                "explanation": "Ø¨Ø±Ø¬ Ø®Ù„ÙŠÙØ© ÙÙŠ Ø¯Ø¨ÙŠ Ù‡Ùˆ Ø£Ø·ÙˆÙ„ Ø¨Ø±Ø¬ ÙÙŠ Ø§Ù„Ø¹Ø§Ù„Ù… Ù…Ù†Ø° Ø§ÙƒØªÙ…Ø§Ù„Ù‡ Ø¹Ø§Ù… 2010."
            },
            {
                "question": "Ù…Ø§ Ù‡Ùˆ Ù…Ø¬Ù…ÙˆØ¹ 7 + 5ØŸ",
                "options": ["10", "12", "13", "14"],
                "correct_index": 1,
                "explanation": "7 + 5 = 12."
            },
            {
                "question": "Ù…Ø§ Ù‡ÙŠ Ø¹Ø§ØµÙ…Ø© ÙØ±Ù†Ø³Ø§ØŸ",
                "options": ["Ø¨Ø§Ø±ÙŠØ³", "Ø±ÙˆÙ…Ø§", "Ø¨Ø±Ù„ÙŠÙ†", "Ù…Ø¯Ø±ÙŠØ¯"],
                "correct_index": 0,
                "explanation": "Ø¨Ø§Ø±ÙŠØ³ Ù‡ÙŠ Ø¹Ø§ØµÙ…Ø© ÙØ±Ù†Ø³Ø§ ÙˆØ£Ø´Ù‡Ø± Ù…Ø¯Ù†Ù‡Ø§."
            },
            {
                "question": "Ù…Ø§ Ù‡Ùˆ Ø£ÙƒØ¨Ø± Ù…Ø­ÙŠØ· ÙÙŠ Ø§Ù„Ø¹Ø§Ù„Ù…ØŸ",
                "options": ["Ø§Ù„Ù…Ø­ÙŠØ· Ø§Ù„Ø£Ø·Ù„Ø³ÙŠ", "Ø§Ù„Ù…Ø­ÙŠØ· Ø§Ù„Ù‡Ù†Ø¯ÙŠ", "Ø§Ù„Ù…Ø­ÙŠØ· Ø§Ù„Ù‡Ø§Ø¯Ø¦", "Ø§Ù„Ù…Ø­ÙŠØ· Ø§Ù„Ù…ØªØ¬Ù…Ø¯ Ø§Ù„Ø´Ù…Ø§Ù„ÙŠ"],
                "correct_index": 2,
                "explanation": "Ø§Ù„Ù…Ø­ÙŠØ· Ø§Ù„Ù‡Ø§Ø¯Ø¦ Ù‡Ùˆ Ø£ÙƒØ¨Ø± Ù…Ø­ÙŠØ· Ø¹Ù„Ù‰ Ø§Ù„Ø£Ø±Ø¶."
            },
            {
                "question": "ÙƒÙ… Ø¹Ø¯Ø¯ Ø£ÙŠØ§Ù… Ø§Ù„Ø£Ø³Ø¨ÙˆØ¹ØŸ",
                "options": ["5", "6", "7", "8"],
                "correct_index": 2,
                "explanation": "Ø§Ù„Ø£Ø³Ø¨ÙˆØ¹ ÙŠØ­ØªÙˆÙŠ Ø¹Ù„Ù‰ 7 Ø£ÙŠØ§Ù…."
            },
            {
                "question": "Ù…Ø§ Ù‡Ùˆ Ø§Ù„ÙƒÙˆÙƒØ¨ Ø§Ù„Ø£Ø­Ù…Ø±ØŸ",
                "options": ["Ø§Ù„Ù…Ø´ØªØ±ÙŠ", "Ø§Ù„Ù…Ø±ÙŠØ®", "Ø§Ù„Ø²Ù‡Ø±Ø©", "Ø¹Ø·Ø§Ø±Ø¯"],
                "correct_index": 1,
                "explanation": "Ø§Ù„Ù…Ø±ÙŠØ® ÙŠØ³Ù…Ù‰ Ø§Ù„ÙƒÙˆÙƒØ¨ Ø§Ù„Ø£Ø­Ù…Ø± Ø¨Ø³Ø¨Ø¨ Ù„ÙˆÙ†Ù‡."
            },
            {
                "question": "Ù…Ø§ Ù‡Ùˆ Ø§Ù„ØºØ§Ø² Ø§Ù„Ø°ÙŠ Ù†ØªÙ†ÙØ³Ù‡ØŸ",
                "options": ["Ø§Ù„Ø£ÙƒØ³Ø¬ÙŠÙ†", "Ø«Ø§Ù†ÙŠ Ø£ÙƒØ³ÙŠØ¯ Ø§Ù„ÙƒØ±Ø¨ÙˆÙ†", "Ø§Ù„Ù‡ÙŠØ¯Ø±ÙˆØ¬ÙŠÙ†", "Ø§Ù„Ù†ÙŠØªØ±ÙˆØ¬ÙŠÙ†"],
                "correct_index": 0,
                "explanation": "Ø§Ù„Ø£ÙƒØ³Ø¬ÙŠÙ† Ù‡Ùˆ Ø§Ù„ØºØ§Ø² Ø§Ù„Ø£Ø³Ø§Ø³ÙŠ Ø§Ù„Ø°ÙŠ Ù†ØªÙ†ÙØ³Ù‡."
            },
            {
                "question": "ÙƒÙ… Ø¹Ø¯Ø¯ Ø§Ù„Ø­ÙˆØ§Ø³ Ø¹Ù†Ø¯ Ø§Ù„Ø¥Ù†Ø³Ø§Ù†ØŸ",
                "options": ["4", "5", "6", "7"],
                "correct_index": 1,
                "explanation": "Ø§Ù„Ø¥Ù†Ø³Ø§Ù† Ù„Ø¯ÙŠÙ‡ Ø®Ù…Ø³ Ø­ÙˆØ§Ø³ Ø±Ø¦ÙŠØ³ÙŠØ©."
            }
        ]
        
        # Convert the Python list of dictionaries to a JSON string.
        # `ensure_ascii=False` is important for correctly handling Arabic characters.
        sample_quiz_json = json.dumps(sample_quiz_data, ensure_ascii=False)
        
        cursor.execute(
            "INSERT INTO sample_quizzes (quiz_code, quiz_data, created_at) VALUES (?, ?, ?)",
            ("sample", sample_quiz_json, datetime.utcnow().isoformat())
        )
        conn.commit()
    conn.close()




# Ù†Ù‚Ø·Ø© Ù†Ù‡Ø§ÙŠØ© Ø§Ù„ÙˆÙŠØ¨ Ù‡ÙˆÙƒ
@app.route('/' + os.getenv('BOT_TOKEN'), methods=['POST'])
def webhook_bot():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot.process_new_updates([update])
        return 'ok', 200
    return 'Method Not Allowed', 405



@app.route('/' + os.getenv('BOT_TOKEN_2'), methods=['POST'])
def webhook_bot2():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot2.process_new_updates([update])
        return 'ok', 200

    return 'Method Not Allowed', 405

@app.route('/' + os.getenv('BOT_TOKEN_3'), methods=['POST'])
def webhook_bot3():
    if request.method == "POST":
        update = telebot.types.Update.de_json(request.stream.read().decode('utf-8'))
        bot3.process_new_updates([update])
        return 'ok', 200

    return 'Method Not Allowed', 405

def set_webhook():
    bot.remove_webhook()
    bot.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN)
    logging.info(f"ğŸŒ ØªÙ… ØªØ¹ÙŠÙŠÙ† Ø§Ù„ÙˆÙŠØ¨ Ù‡ÙˆÙƒ Ø¹Ù„Ù‰: {WEBHOOK_URL}/{BOT_TOKEN}")


    bot2.remove_webhook()
    bot2.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN_2)
    logging.info(f"ğŸŒ ØªÙ… ØªØ¹ÙŠÙŠÙ† Ø§Ù„ÙˆÙŠØ¨ Ù‡ÙˆÙƒ Ø¹Ù„Ù‰: {WEBHOOK_URL}/{BOT_TOKEN_2}")
    
    bot3.remove_webhook()
    bot3.set_webhook(url=WEBHOOK_URL + '/' + BOT_TOKEN_3)
    logging.info(f"ğŸŒ ØªÙ… ØªØ¹ÙŠÙŠÙ† Ø§Ù„ÙˆÙŠØ¨ Ù‡ÙˆÙƒ Ø¹Ù„Ù‰: {WEBHOOK_URL}/{BOT_TOKEN_3}")



import schedule
import time
import threading
from datetime import datetime

# Ø§Ø³ØªØ¯Ø¹Ø§Ø¡ Ø§Ù„Ø¯ÙˆØ§Ù„ Ù…Ø¨Ø§Ø´Ø±Ø© Ø¥Ø°Ø§ ÙƒØ§Ù†Øª ÙÙŠ Ù†ÙØ³ Ø§Ù„Ù…Ù„Ù
# Ø¥Ø°Ø§ ÙƒØ§Ù†Øª ÙÙŠ Ù…Ù„Ù Ø«Ø§Ù†ÙŠ Ø§Ø³ØªØ¨Ø¯Ù„ Ø¨Ù€ from stats import ...

def run_reports():
    """ØªØ´ØºÙŠÙ„ ØªÙ‚Ø§Ø±ÙŠØ± Ø§Ù„ÙŠÙˆÙ…"""
    send_daily_report()
    send_top_users_report(5)  # Ù†Ø±Ø³Ù„ Ø£ÙØ¶Ù„ 5

# Ø¬Ø¯ÙˆÙ„Ø© Ø§Ù„ØªÙ‚Ø§Ø±ÙŠØ± Ù„ØªØ¹Ù…Ù„ ÙŠÙˆÙ…ÙŠØ§Ù‹ Ø§Ù„Ø³Ø§Ø¹Ø© 23:59
schedule.every().day.at("23:59").do(run_reports)

def run_scheduler():
    """Ø§Ù„ØªØ´ØºÙŠÙ„ Ø§Ù„Ù…Ø³ØªÙ…Ø± Ù„Ù„Ø¬Ø¯ÙˆÙ„Ø©"""
    while True:
        schedule.run_pending()
        time.sleep(1)

# ØªØ´ØºÙŠÙ„ Ø§Ù„Ù€ scheduler ÙÙŠ thread Ù…Ù†ÙØµÙ„
def start_scheduler():
    t = threading.Thread(target=run_scheduler, daemon=True)
    t.start()
    print(f"ğŸ“… Scheduler started at {datetime.now()}")


start_scheduler()

if __name__ == "__main__":
    init_all_dbs()
    insert_sample_quiz_if_not_exists()
    set_webhook()
    port = int(os.environ.get('PORT', 10000))  # Render ÙŠØ³ØªØ®Ø¯Ù… 10000
    app.run(host='0.0.0.0', port=port)
